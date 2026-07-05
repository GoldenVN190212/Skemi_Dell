import json
import logging
import httpx
import asyncio
import importlib.util
import inspect
from datetime import datetime, timedelta
from urllib.parse import urlparse
try:
    from trusted_sources import get_quality_boost
except ImportError:
    def get_quality_boost(domain): return 1.0
from session_context import (
    create_session,
    append_message as session_append,
    get_context as session_get_context,
    touch_session,
    delete_session,
    cleanup_expired_sessions,
)
from fastapi import FastAPI, Request, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
import contextlib
from contextlib import asynccontextmanager
from typing import List, Dict, Optional, Any
from urllib.parse import urlparse
import time
import os
import sqlite3
import sys
import re
import io
import csv
import base64
import hashlib
from pathlib import Path
import subprocess
import shutil
import math
import tempfile
import difflib
from collections import Counter
import statistics
import unicodedata
try:
    from sentence_transformers import SentenceTransformer, CrossEncoder
    ST_AVAILABLE = True
except Exception:
    SentenceTransformer = None
    CrossEncoder = None
    ST_AVAILABLE = False

# Optional advanced modules
MULTI_TIER_OCR_AVAILABLE = False
try:
    import multi_tier_ocr
    MULTI_TIER_OCR_AVAILABLE = True
except Exception:
    pass

ADVANCED_CACHE_AVAILABLE = False
try:
    import ai_news_ecosystem
    ADVANCED_CACHE_AVAILABLE = True
except Exception:
    pass

LANGDETECT_AVAILABLE = False
try:
    from langdetect import detect as _langdetect_detect
    from langdetect import DetectorFactory as _LangDetectFactory
    from langdetect.lang_detect_exception import LangDetectException

    _LangDetectFactory.seed = 0
    LANGDETECT_AVAILABLE = True
except Exception:
    _langdetect_detect = None
    _LangDetectFactory = None

    class LangDetectException(Exception):
        pass

# Force UTF-8 console on Windows to avoid mojibake in terminal logs.
if os.name == "nt":
    try:
        subprocess.run("chcp 65001 > nul", shell=True, capture_output=True, text=True)
    except Exception:
        pass
    try:
        if hasattr(sys.stdout, "reconfigure"):
            sys.stdout.reconfigure(encoding="utf-8", errors="replace")
        if hasattr(sys.stderr, "reconfigure"):
            sys.stderr.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

MOJIBAKE_TOKENS = ("Ã", "ðŸ", "ï¸", "âœ", "âš", "áº", "á»")

def _repair_mojibake_text(text: str) -> str:
    if not isinstance(text, str):
        return text
    if not any(tok in text for tok in MOJIBAKE_TOKENS):
        return text
    try:
        fixed = text.encode("latin1", errors="ignore").decode("utf-8", errors="ignore")
        bad_before = sum(text.count(tok) for tok in MOJIBAKE_TOKENS)
        bad_after = sum(fixed.count(tok) for tok in MOJIBAKE_TOKENS)
        if fixed and bad_after < bad_before:
            return fixed
    except Exception:
        pass
    return text

class MojibakeSafeFormatter(logging.Formatter):
    def format(self, record: logging.LogRecord) -> str:
        message = super().format(record)
        return _repair_mojibake_text(message)

# --- Local SQLite Cache System ---
CACHE_DB = os.path.join("cache", "cache.db")
os.makedirs("cache", exist_ok=True)

class SQLiteCache:
    @staticmethod
    def _get_conn():
        conn = sqlite3.connect(CACHE_DB)
        conn.execute("CREATE TABLE IF NOT EXISTS message_cache (key TEXT PRIMARY KEY, value TEXT, expires REAL)")
        conn.execute("CREATE TABLE IF NOT EXISTS file_cache (hash TEXT PRIMARY KEY, type TEXT, value TEXT, expires REAL)")
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS semantic_response_cache (
                cache_key TEXT PRIMARY KEY,
                question_norm TEXT NOT NULL,
                context_fingerprint TEXT NOT NULL,
                response TEXT NOT NULL,
                response_gist TEXT,
                style_hint TEXT,
                model_name TEXT,
                created_at REAL NOT NULL,
                expires REAL NOT NULL
            )
            """
        )
        conn.execute("CREATE INDEX IF NOT EXISTS idx_semantic_question ON semantic_response_cache(question_norm)")
        conn.execute("CREATE INDEX IF NOT EXISTS idx_semantic_context ON semantic_response_cache(context_fingerprint)")
        return conn

    @classmethod
    async def get(cls, table, key, key_col="key", extra_filters: Optional[Dict[str, Any]] = None):
        try:
            with cls._get_conn() as conn:
                cur = conn.cursor()
                clauses = [f"{key_col} = ?", "expires > ?"]
                params: List[Any] = [key, time.time()]
                for filter_key, filter_value in (extra_filters or {}).items():
                    clauses.append(f"{filter_key} = ?")
                    params.append(filter_value)
                cur.execute(f"SELECT value FROM {table} WHERE {' AND '.join(clauses)}", tuple(params))
                row = cur.fetchone()
                return json.loads(row[0]) if row else None
        except Exception: return None

    @classmethod
    async def set(cls, table, key, value, ttl, key_col="key"):
        try:
            with cls._get_conn() as conn:
                conn.execute(f"INSERT OR REPLACE INTO {table} ({key_col}, value, expires) VALUES (?, ?, ?)",
                             (key, json.dumps(value), time.time() + ttl))
        except Exception: pass

class MessageCache:
    @staticmethod
    async def get_message(q): return await SQLiteCache.get("message_cache", q)
    @staticmethod
    async def set_message(q, v, ttl): await SQLiteCache.set("message_cache", q, v, ttl)

class FileCache:
    @staticmethod
    async def get_file(h, t): return await SQLiteCache.get("file_cache", h, key_col="hash", extra_filters={"type": t})
    @staticmethod
    async def set_file(h, t, v, ttl): await SQLiteCache.set("file_cache", h, v, ttl, "hash")

class CacheType:
    TEXT_MESSAGE = "messages"
    IMAGE_FILE = "images"
    DOCUMENT_FILE = "docs"

class advanced_cache:
    @staticmethod
    async def get_stats():
        class Stats:
            def __init__(self):
                self.total_entries = 0
                self.total_size_bytes = 0
                self.hit_rate = 0
                self.eviction_count = 0
                self.cache_type_stats = {}

        stats = Stats()
        try:
            with SQLiteCache._get_conn() as conn:
                now_ts = time.time()
                message_count = int(conn.execute("SELECT COUNT(*) FROM message_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                file_count = int(conn.execute("SELECT COUNT(*) FROM file_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                semantic_count = int(conn.execute("SELECT COUNT(*) FROM semantic_response_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                message_size = int(conn.execute("SELECT COALESCE(SUM(LENGTH(value)), 0) FROM message_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                file_size = int(conn.execute("SELECT COALESCE(SUM(LENGTH(value)), 0) FROM file_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                semantic_size = int(conn.execute("SELECT COALESCE(SUM(LENGTH(response) + LENGTH(COALESCE(response_gist, ''))), 0) FROM semantic_response_cache WHERE expires > ?", (now_ts,)).fetchone()[0] or 0)
                stats.total_entries = message_count + file_count + semantic_count
                stats.total_size_bytes = message_size + file_size + semantic_size
                stats.cache_type_stats = {
                    "messages": message_count,
                    "files": file_count,
                    "semantic": semantic_count,
                }
        except Exception:
            pass
        return stats
    
    @staticmethod
    async def clear_by_type(t):
        try:
            with SQLiteCache._get_conn() as conn:
                if t == CacheType.TEXT_MESSAGE:
                    conn.execute("DELETE FROM message_cache")
                    conn.execute("DELETE FROM semantic_response_cache")
                else:
                    conn.execute("DELETE FROM file_cache")
            return 0
        except Exception: return 0

    @staticmethod
    async def clear_all():
        try:
            with SQLiteCache._get_conn() as conn:
                conn.execute("DELETE FROM message_cache")
                conn.execute("DELETE FROM file_cache")
                conn.execute("DELETE FROM semantic_response_cache")
            return 0
        except Exception: return 0
    
    @staticmethod
    def start_background_cleanup(): pass
    
    @staticmethod
    def stop_background_cleanup(): pass

CACHE_AVAILABLE = True
MAX_SESSION_HISTORY_MESSAGES = int(os.getenv("SKEMI_MAX_SESSION_HISTORY_MESSAGES", "160"))
MAX_HISTORY_MESSAGES_TO_MODEL = int(os.getenv("SKEMI_MAX_HISTORY_MESSAGES_TO_MODEL", "18"))
MAX_HISTORY_MESSAGE_CHARS = int(os.getenv("SKEMI_MAX_HISTORY_MESSAGE_CHARS", "320"))
MAX_HISTORY_SUMMARY_CHARS = int(os.getenv("SKEMI_MAX_HISTORY_SUMMARY_CHARS", "2200"))
SEMANTIC_CACHE_TTL_SECONDS = int(os.getenv("SKEMI_SEMANTIC_CACHE_TTL_SECONDS", "21600"))
SEMANTIC_CACHE_DIRECT_HIT_THRESHOLD = float(os.getenv("SKEMI_SEMANTIC_CACHE_DIRECT_HIT_THRESHOLD", "0.96"))
SEMANTIC_CACHE_REUSE_THRESHOLD = float(os.getenv("SKEMI_SEMANTIC_CACHE_REUSE_THRESHOLD", "0.84"))
_MODEL_CONTEXT_WINDOW_CACHE: Dict[str, Dict[str, Any]] = {}

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import camelot
except Exception:
    camelot = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

# Multi-tier OCR fallback
run_multi_tier_ocr = None
multi_tier_ocr_system = None
MULTI_TIER_OCR_AVAILABLE = False
OCR_STARTUP_DIAGNOSTICS = {"ready": True, "critical_missing": [], "required_engine_errors": [], "tiers": []}

# --- Rate Limiting ---
request_queue = asyncio.Queue()
user_requests: Dict[str, Dict] = {}
MAX_REQUESTS_PER_DAY = 10
PROMPT_TOKEN_LIMIT = 30000  # 30k tokens for input
MAX_OUTPUT_TOKENS = 50000  # 50k tokens for output
MAX_FILE_BYTES = 50 * 1024 * 1024  # 50MB
MAX_FILE_TOKENS = 30000  # 30k tokens for file content
OCR_MIN_CHARS = 20
OCR_FAIL_FAST_STARTUP = os.getenv("SKEMI_OCR_FAIL_FAST", "0").lower() in ("1", "true", "yes")
FORCE_ENGLISH = os.getenv("SKEMI_FORCE_ENGLISH", "0").lower() in ("1", "true", "yes")

OLLAMA_URL = "http://127.0.0.1:11434/api/chat"
GENERATE_URL = "http://127.0.0.1:11434/api/generate"
MODEL_MAIN = os.getenv("SKEMI_MODEL_MAIN", "devstral-2:123b-cloud")
MODEL_MAIN_CANARY = os.getenv("SKEMI_MODEL_MAIN_CANARY", "").strip()
MODEL_ROUTER = os.getenv("SKEMI_MODEL_ROUTER", "devstral-2:123b-cloud")
# Cloud-readiness, resolved once at startup by probing Ollama-Cloud auth. The
# advertised MODEL_MAIN / MODEL_ROUTER keep their cloud (":cloud") names so the
# product ships configured for the cloud model and activates the moment `ollama signin`
# is done. Until then, generation transparently uses LOCAL_PRIMARY_FALLBACK so
# nothing breaks. See check_ollama_status() and _select_generation_model().
MODEL_MAIN_CLOUD_READY = False
MODEL_ROUTER_CLOUD_READY = False
LOCAL_PRIMARY_FALLBACK = ""
MODEL_STATUS = os.getenv("SKEMI_MODEL_STATUS", "devstral-2:123b-cloud").strip() or "devstral-2:123b-cloud"
STATUS_MODEL_ACTIVE = False
# Use Ollama vision model instead of premium cloud model
MODEL_VISION = os.getenv("SKEMI_MODEL_VISION", "moondream:latest")
CHAT_TIMEOUT_SECONDS = float(os.getenv("SKEMI_CHAT_TIMEOUT", "120"))
IMAGE_FORCE_VISION = os.getenv("SKEMI_IMAGE_FORCE_VISION", "1").lower() in ("1", "true", "yes")
IMAGE_MIN_WORDS = int(os.getenv("SKEMI_IMAGE_MIN_WORDS", os.getenv("SKEMI_IMAGE_MIN_TOKENS", "500")))
IMAGE_MAX_WORDS = int(os.getenv("SKEMI_IMAGE_MAX_WORDS", os.getenv("SKEMI_IMAGE_MAX_TOKENS", "5000")))
EMBED_MODEL_NAME = os.getenv("SKEMI_EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
RERANK_MODEL_NAME = os.getenv("SKEMI_RERANK_MODEL", "cross-encoder/ms-marco-MiniLM-L-6-v2")
AUTO_LEARNING_INTERVAL_SECONDS = int(os.getenv("SKEMI_AUTO_LEARNING_INTERVAL_SECONDS", "180"))
AUTO_LEARNING_ENABLED = os.getenv("SKEMI_AUTO_LEARNING_ENABLED", "1").lower() in ("1", "true", "yes")
CANARY_ENABLED = os.getenv("SKEMI_CANARY_ENABLED", "0").lower() in ("1", "true", "yes")
CANARY_PERCENT = max(0.0, min(100.0, float(os.getenv("SKEMI_CANARY_PERCENT", "5"))))
_embed_model = None
_rerank_model = None
_status_text_cache: Dict[str, str] = {}

DB_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data")
os.makedirs(DB_DIR, exist_ok=True)
DB_FILE = os.path.join(DB_DIR, "chat_memory.db")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s', force=True)
logger = logging.getLogger("Skemi")
for noisy_logger in (
    "uvicorn.access",
    "uvicorn.protocols.websockets.websockets_impl",
    "uvicorn.error",
):
    logging.getLogger(noisy_logger).setLevel(logging.WARNING)
for handler in logging.getLogger().handlers:
    handler.setFormatter(MojibakeSafeFormatter('%(asctime)s - %(levelname)s - %(message)s'))

OCR_STARTUP_DIAGNOSTICS: Dict[str, Any] = {
    "ready": False,
    "critical_missing": ["OCR startup has not run yet"],
    "required_engine_errors": [],
}
ANALYSIS_MODES = {"assistant", "structured", "concise"}

# Cache helper functions
def _trim_to_max_words(text: str, max_words: int) -> str:
    """Trim text to maximum word count"""
    if not text:
        return text
    
    words = text.split()
    if len(words) <= max_words:
        return text
    
    return ' '.join(words[:max_words]) + "..."

def generate_content_hash(content: str) -> str:
    """Generate hash for content"""
    return hashlib.md5(content.encode('utf-8')).hexdigest()

def generate_file_hash(file_data: bytes) -> str:
    """Generate hash for file data"""
    return hashlib.md5(file_data).hexdigest()

def _normalize_cache_text(text: str) -> str:
    raw = unicodedata.normalize("NFKD", str(text or ""))
    raw = "".join(ch for ch in raw if not unicodedata.combining(ch))
    raw = raw.replace("\u0111", "d").replace("\u0110", "d")
    raw = re.sub(r"[^\w\s]", " ", raw.lower())
    return re.sub(r"\s+", " ", raw).strip()

def _tokenize_cache_text(text: str) -> set:
    return {token for token in _normalize_cache_text(text).split() if len(token) > 1}

def _cache_similarity_score(a: str, b: str) -> float:
    left = _normalize_cache_text(a)
    right = _normalize_cache_text(b)
    if not left or not right:
        return 0.0
    seq_score = difflib.SequenceMatcher(None, left, right).ratio()
    left_tokens = _tokenize_cache_text(left)
    right_tokens = _tokenize_cache_text(right)
    union = left_tokens | right_tokens
    token_score = (len(left_tokens & right_tokens) / max(1, len(union))) if union else 0.0
    return seq_score * 0.65 + token_score * 0.35

def _compact_memory_text(text: str, max_chars: int = MAX_HISTORY_MESSAGE_CHARS) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(cleaned) <= max_chars:
        return cleaned
    return cleaned[:max_chars].rstrip() + "..."

def _build_response_gist(text: str, max_chars: int = 700) -> str:
    if not text:
        return ""
    lines = [line.strip() for line in str(text).splitlines() if line.strip()]
    if not lines:
        return _compact_memory_text(text, max_chars=max_chars)
    picked: List[str] = []
    total = 0
    for line in lines:
        compact = _compact_memory_text(line, max_chars=min(240, max_chars))
        if not compact:
            continue
        if total + len(compact) > max_chars and picked:
            break
        picked.append(compact)
        total += len(compact) + 1
        if len(picked) >= 4:
            break
    return "\n".join(picked)[:max_chars]

def _normalize_language_code(code: str) -> str:
    raw = str(code or "").strip().lower().replace("_", "-")
    if not raw:
        return ""
    if raw.startswith("zh"):
        return "zh"
    if raw.startswith("ja"):
        return "ja"
    if raw.startswith("ko"):
        return "ko"
    if raw.startswith("vi"):
        return "vi"
    if raw.startswith("en"):
        return "en"
    if raw.startswith("pt"):
        return "pt"
    if raw.startswith("es"):
        return "es"
    if raw.startswith("fr"):
        return "fr"
    if raw.startswith("de"):
        return "de"
    if raw.startswith("it"):
        return "it"
    if raw.startswith("ru"):
        return "ru"
    if raw.startswith("uk"):
        return "uk"
    if raw.startswith("ar"):
        return "ar"
    if raw.startswith("hi"):
        return "hi"
    if raw.startswith("th"):
        return "th"
    if raw.startswith("id"):
        return "id"
    if raw.startswith("ms"):
        return "ms"
    if raw.startswith("tr"):
        return "tr"
    return raw.split("-", 1)[0]
def _detect_language_hint_from_text(text: str) -> str:
    sample = str(text or "").strip()
    if not sample:
        return ""
    if FORCE_ENGLISH:
        return "en"

    if re.search(r"[\u3040-\u30ff]", sample):
        return "ja"
    if re.search(r"[\uac00-\ud7af]", sample):
        return "ko"
    if re.search(r"[\u4e00-\u9fff]", sample):
        return "zh"
    if re.search(r"[\u0600-\u06ff]", sample):
        return "ar"
    if re.search(r"[\u0400-\u04ff]", sample):
        return "ru"
    if re.search(r"[\u0e00-\u0e7f]", sample):
        return "th"
    if re.search(r"[\u0900-\u097f]", sample):
        return "hi"

    lowered = sample.lower()
    if any(ch in lowered for ch in "ăâđêôơưáàảãạắằẳẵặấầẩẫậéèẻẽẹếềểễệíìỉĩịóòỏõọốồổỗộớờởỡợúùủũụứừửữựýỳỷỹỵ"):
        return "vi"

    normalized = _normalize_cache_text(sample)
    token_count = len(normalized.split())
    if LANGDETECT_AVAILABLE and (len(normalized) >= 18 or token_count >= 4):
        try:
            return _normalize_language_code(_langdetect_detect(sample))
        except LangDetectException:
            return ""
        except Exception:
            return ""

    return ""

def _infer_response_style(text: str) -> str:
    sample = str(text or "")
    if re.search(r"^\|.+\|$", sample, re.MULTILINE):
        return "table"
    if re.search(r"^\s*[-*•]\s+", sample, re.MULTILINE):
        return "bullet"
    if re.search(r"^\s*#{2,4}\s+", sample, re.MULTILINE):
        return "sectioned"
    return "plain"

def _build_context_fingerprint(
    question: str,
    history_summary: str = "",
    recent_messages: Optional[List[Dict[str, str]]] = None,
    force_search: bool = False,
    deep_research: bool = False,
    has_image: bool = False,
    language_hint: str = "",
) -> str:
    recent_lines = []
    for item in (recent_messages or [])[-4:]:
        role = str(item.get("role", "user")).strip().lower()
        content = _compact_memory_text(item.get("content", ""), max_chars=180)
        if content:
            recent_lines.append(f"{role}:{content}")
    payload = "|".join(
        [
            _normalize_cache_text(question),
            _normalize_cache_text(history_summary)[:700],
            " || ".join(recent_lines),
            f"force={int(bool(force_search))}",
            f"deep={int(bool(deep_research))}",
            f"image={int(bool(has_image))}",
            f"lang={_normalize_language_code(language_hint)}",
        ]
    )
    return hashlib.md5(payload.encode("utf-8")).hexdigest()

def _make_exact_cache_key(question_norm: str, context_fingerprint: str) -> str:
    return hashlib.md5(f"{context_fingerprint}|{question_norm}".encode("utf-8")).hexdigest()

async def get_cached_response(question: str, context_fingerprint: str = "") -> Optional[Dict[str, Any]]:
    """Get cached response or reusable gist for a semantically similar question."""
    if not CACHE_AVAILABLE:
        return None
    question_norm = _normalize_cache_text(question)
    if not question_norm:
        return None

    try:
        exact_key = _make_exact_cache_key(question_norm, context_fingerprint or "global")
        cached_response = await MessageCache.get_message(exact_key)
        if cached_response and isinstance(cached_response, dict) and cached_response.get("response"):
            logger.info(f" Exact cache HIT for question: {question[:50]}...")
            return {"mode": "exact", **cached_response}
        if cached_response and isinstance(cached_response, str):
            logger.info(f" Legacy cache HIT for question: {question[:50]}...")
            return {"mode": "exact", "response": cached_response, "gist": _build_response_gist(cached_response), "style": _infer_response_style(cached_response)}

        with SQLiteCache._get_conn() as conn:
            rows = conn.execute(
                """
                SELECT question_norm, response, response_gist, style_hint
                FROM semantic_response_cache
                WHERE context_fingerprint = ? AND expires > ?
                ORDER BY created_at DESC
                LIMIT 40
                """,
                (context_fingerprint or "global", time.time()),
            ).fetchall()
        if not rows:
            return None

        best_row = None
        best_score = 0.0
        for cached_question, cached_answer, cached_gist, style_hint in rows:
            if cached_question == question_norm:
                logger.info(f" Semantic cache exact-normalized HIT for question: {question[:50]}...")
                return {
                    "mode": "exact",
                    "response": cached_answer,
                    "gist": cached_gist or _build_response_gist(cached_answer),
                    "style": style_hint or _infer_response_style(cached_answer),
                }
            score = _cache_similarity_score(question_norm, cached_question)
            if score > best_score:
                best_score = score
                best_row = (cached_answer, cached_gist, style_hint)

        if best_row and best_score >= SEMANTIC_CACHE_DIRECT_HIT_THRESHOLD:
            cached_answer, cached_gist, style_hint = best_row
            logger.info(f" Semantic direct HIT ({best_score:.2f}) for question: {question[:50]}...")
            return {
                "mode": "direct",
                "response": cached_answer,
                "gist": cached_gist or _build_response_gist(cached_answer),
                "style": style_hint or _infer_response_style(cached_answer),
            }
        if best_row and best_score >= SEMANTIC_CACHE_REUSE_THRESHOLD:
            cached_answer, cached_gist, style_hint = best_row
            logger.info(f" Semantic gist reuse HIT ({best_score:.2f}) for question: {question[:50]}...")
            return {
                "mode": "gist",
                "response": None,
                "gist": cached_gist or _build_response_gist(cached_answer),
                "style": style_hint or _infer_response_style(cached_answer),
            }
        return None
    except Exception as e:
        logger.error(f"Cache error: {e}")
        return None

async def cache_response(
    question: str,
    response: str,
    context_fingerprint: str = "",
    model_name: str = "",
    ttl_seconds: int = SEMANTIC_CACHE_TTL_SECONDS,
):
    """Cache response, gist, and style so similar follow-ups can be served cheaper."""
    if not CACHE_AVAILABLE:
        return
    question_norm = _normalize_cache_text(question)
    if not question_norm or not response:
        return

    try:
        gist = _build_response_gist(response)
        style_hint = _infer_response_style(response)
        exact_key = _make_exact_cache_key(question_norm, context_fingerprint or "global")
        await MessageCache.set_message(
            exact_key,
            {
                "response": response,
                "gist": gist,
                "style": style_hint,
                "question_norm": question_norm,
            },
            ttl_seconds,
        )
        with SQLiteCache._get_conn() as conn:
            conn.execute(
                """
                INSERT OR REPLACE INTO semantic_response_cache
                (cache_key, question_norm, context_fingerprint, response, response_gist, style_hint, model_name, created_at, expires)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    exact_key,
                    question_norm,
                    context_fingerprint or "global",
                    response,
                    gist,
                    style_hint,
                    model_name,
                    time.time(),
                    time.time() + ttl_seconds,
                ),
            )
        logger.info(f" Cached response for question: {question[:50]}...")
    except Exception as e:
        logger.error(f"Cache save error: {e}")

async def get_cached_file_analysis(file_hash: str, file_type: str) -> Optional[Dict]:
    """Get cached file analysis"""
    if not CACHE_AVAILABLE:
        return None
    
    try:
        cached_analysis = await FileCache.get_file(file_hash, file_type)
        if cached_analysis:
            logger.info(f" File Cache HIT for {file_type}: {file_hash[:16]}...")
            return cached_analysis
        return None
    except Exception as e:
        logger.error(f"File cache error: {e}")
        return None

async def cache_file_analysis(file_hash: str, file_type: str, analysis: Dict, ttl_seconds: int = 7200):
    """Cache file analysis"""
    if not CACHE_AVAILABLE:
        return
    
    try:
        await FileCache.set_file(file_hash, file_type, analysis, ttl_seconds)
        logger.info(f" Cached {file_type} analysis: {file_hash[:16]}...")
    except Exception as e:
        logger.error(f"File cache save error: {e}")

async def get_cache_stats():
    """Get cache statistics"""
    if not CACHE_AVAILABLE:
        return {"cache_available": False}
    
    try:
        stats = await advanced_cache.get_stats()
        return {
            "cache_available": True,
            "total_entries": stats.total_entries,
            "total_size_mb": round(stats.total_size_bytes / (1024*1024), 2),
            "hit_rate": round(stats.hit_rate * 100, 2),
            "evictions": stats.eviction_count,
            "cache_type_stats": stats.cache_type_stats
        }
    except Exception as e:
        logger.error(f"Cache stats error: {e}")
        return {"cache_available": False, "error": str(e)}

async def check_ollama_status():
    global MODEL_ROUTER, MODEL_MAIN, MODEL_STATUS, MODEL_VISION, STATUS_MODEL_ACTIVE
    global MODEL_MAIN_CLOUD_READY, MODEL_ROUTER_CLOUD_READY, LOCAL_PRIMARY_FALLBACK

    async def _probe() -> tuple[bool, list[str], list[dict[str, Any]]]:
        async with httpx.AsyncClient() as client:
            resp = await client.get("http://127.0.0.1:11434/api/tags", timeout=10.0)
            if resp.status_code != 200:
                return False, [], []
            model_entries = resp.json().get("models", [])
            models = [m.get("name", "") for m in model_entries if m.get("name")]
            return True, models, model_entries

    def _find_ollama_binary() -> str:
        candidates = [
            os.getenv("SKEMI_OLLAMA_BIN", "").strip(),
            shutil.which("ollama") or "",
            r"D:\Ollama\ollama.exe",
            r"C:\Program Files\Ollama\ollama.exe",
            r"C:\Users\DELL\AppData\Local\Programs\Ollama\ollama.exe",
        ]
        for candidate in candidates:
            if candidate and os.path.exists(candidate):
                return candidate
        return ""

    async def _try_start_ollama() -> bool:
        ollama_bin = _find_ollama_binary()
        if not ollama_bin:
            return False
        try:
            creation_flags = 0
            if os.name == "nt":
                creation_flags = getattr(subprocess, "DETACHED_PROCESS", 0) | getattr(subprocess, "CREATE_NEW_PROCESS_GROUP", 0)
            subprocess.Popen(
                [ollama_bin, "serve"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                stdin=subprocess.DEVNULL,
                creationflags=creation_flags,
                close_fds=False if os.name == "nt" else True,
            )
        except Exception as exc:
            logger.warning(f"Could not auto-start Ollama: {exc}")
            return False
        for _ in range(12):
            await asyncio.sleep(1.0)
            try:
                ok, _, _ = await _probe()
            except Exception:
                ok = False
            if ok:
                logger.info(" Auto-started Ollama service successfully.")
                return True
        return False

    try:
        ok, models, model_entries = await _probe()
    except Exception as e:
        logger.error(f"Ollama check error: {e}")
        ok = False
        models = []
        model_entries = []
    if not ok:
        logger.warning(" Ollama is offline at startup. Attempting auto-start...")
        with contextlib.suppress(Exception):
            await _try_start_ollama()
        try:
            ok, models, model_entries = await _probe()
        except Exception as e:
            logger.error(f"Ollama recheck error: {e}")
            ok = False
            models = []
            model_entries = []
    if ok:
        logger.info(f" Ollama is ONLINE. Available models: {models}")

        if model_entries:
            def _parse_params(raw: str) -> int:
                try:
                    digits = re.sub(r"[^0-9]", "", str(raw or ""))
                    return int(digits) if digits else 0
                except Exception:
                    return 0

            sorted_by_size = sorted(
                model_entries,
                key=lambda m: _parse_params((m.get("details") or {}).get("parameter_size", "")),
            )
            router_candidate = sorted_by_size[0].get("name", models[0] if models else "llama3.2:latest")
            main_candidate = sorted_by_size[-1].get("name", models[0] if models else "llama3.2:latest")

            vision_candidate = ""
            for m in sorted_by_size:
                fam = str((m.get("details") or {}).get("family", "")).lower()
                if "vl" in fam:
                    vision_candidate = m.get("name", "")
                    break

            # Cloud models (name ends with ":cloud") require Ollama-Cloud auth that may
            # be absent on this machine (calls return HTTP 401, empty output). Prefer a
            # LOCAL model for MODEL_MAIN / MODEL_ROUTER whenever one is installed so core
            # generation (chat, prompt-agent, studio, synthesis) keeps working offline.
            def _is_cloud(name: str) -> bool:
                return str(name or "").strip().lower().endswith(":cloud")

            # Source of truth = the plain model-name list returned by Ollama.
            local_model_names = [n for n in models if n and not _is_cloud(n)]

            def _pref_local(preferred: str) -> str:
                if preferred in local_model_names:
                    return preferred
                return local_model_names[0] if local_model_names else ""

            local_main = _pref_local("qwen2.5:3b")
            local_router = _pref_local("qwen2.5:1.5b")

            # A ":cloud" model only works when Ollama-Cloud auth is present (otherwise
            # calls return HTTP 401). Probe it once: if it answers, keep the stronger
            # cloud model; if it fails, fall back to a local model so generation still
            # works offline. After `ollama signin`, a restart auto-activates cloud.
            async def _cloud_usable(name: str) -> bool:
                if not _is_cloud(name):
                    return True
                try:
                    return bool(await _raw_generate_once(name, "ping", 20.0, 1))
                except Exception:
                    return False

            # Local fallback used by generation whenever a configured cloud model is
            # not yet authenticated (keeps streaming chat / studio / prompt-agent alive).
            LOCAL_PRIMARY_FALLBACK = local_main or local_router or ""

            # MODEL_MAIN keeps its configured name (per product decision: ship the
            # gemini-3-flash-preview:cloud tag and run cloud after `ollama signin`).
            # We only swap it away when the name is not a cloud model AND is missing
            # from the catalog (a genuinely invalid local name). For cloud models we
            # keep the name and record readiness; _select_generation_model() routes
            # actual generation to LOCAL_PRIMARY_FALLBACK until cloud auth is present.
            main_ok = await _cloud_usable(MODEL_MAIN)
            if _is_cloud(MODEL_MAIN):
                MODEL_MAIN_CLOUD_READY = bool(main_ok)
                if main_ok:
                    logger.info(f"Cloud main model '{MODEL_MAIN}' authenticated and active.")
                else:
                    logger.warning(
                        f"Cloud main model '{MODEL_MAIN}' kept as configured but NOT yet authenticated "
                        f"(HTTP 401). Generation uses local fallback '{LOCAL_PRIMARY_FALLBACK or 'none'}' "
                        f"until `ollama signin`."
                    )
            elif MODEL_MAIN not in models:
                MODEL_MAIN = local_main or main_candidate
                MODEL_MAIN_CLOUD_READY = False
                logger.warning(f"Main model unusable; using '{MODEL_MAIN}'")
            else:
                MODEL_MAIN_CLOUD_READY = True  # local model, always "ready"

            router_ok = await _cloud_usable(MODEL_ROUTER)
            if _is_cloud(MODEL_ROUTER):
                MODEL_ROUTER_CLOUD_READY = bool(router_ok)
                if router_ok:
                    logger.info(f"Cloud router model '{MODEL_ROUTER}' authenticated and active.")
                else:
                    logger.warning(
                        f"Cloud router model '{MODEL_ROUTER}' kept as configured but NOT yet authenticated; "
                        f"router uses local fallback '{local_router or local_main or 'none'}' until `ollama signin`."
                    )
            elif MODEL_ROUTER not in models:
                MODEL_ROUTER = local_router or local_main or router_candidate
                MODEL_ROUTER_CLOUD_READY = False
                logger.warning(f"Router model unusable; using '{MODEL_ROUTER}'")
            else:
                MODEL_ROUTER_CLOUD_READY = True

            # Keep the last-resort LOCAL fallback list in sync with what is
            # ACTUALLY installed. Hardcoding names (qwen2.5:3b/1.5b) caused
            # repeated HTTP 404s once those models were no longer pulled on this
            # machine. An empty list means "no local model available" so generation
            # fails fast with a clear message instead of spamming 404s at every tier.
            global LOCAL_FALLBACK_MODELS
            _preferred_local = [m for m in ("qwen2.5:3b", "qwen2.5:1.5b") if m in local_model_names]
            _other_local = [m for m in local_model_names if m not in _preferred_local]
            LOCAL_FALLBACK_MODELS = _preferred_local + _other_local
            if LOCAL_FALLBACK_MODELS:
                logger.info(f"Local fallback models available: {LOCAL_FALLBACK_MODELS}")
            else:
                logger.warning(
                    "No LOCAL model is installed and cloud models need `ollama signin`. "
                    "AI features (search synthesis, chat, prompt-agent, quiz) will return a "
                    "'model unavailable' message until you run `ollama signin` OR "
                    "`ollama pull qwen2.5:3b`."
                )

            STATUS_MODEL_ACTIVE = bool(MODEL_STATUS and MODEL_STATUS in models)
            if MODEL_STATUS not in models:
                logger.warning(f"Status model {MODEL_STATUS} not found!")
                logger.info("Status generation will use instant built-in status text instead of a model.")
            else:
                logger.info(f"Status model ready: {MODEL_STATUS}")

            if MODEL_VISION and MODEL_VISION not in models:
                logger.warning(f"Vision model {MODEL_VISION} not found!")
                MODEL_VISION = vision_candidate
                if MODEL_VISION:
                    logger.info(f"Auto-selected vision model: {MODEL_VISION}")
                else:
                    logger.warning("No multimodal vision model detected in catalog; image analysis falls back to OCR/text mode.")

        return True, models
    logger.error(" Ollama is OFFLINE!")
    return False, []

def _default_context_window_for_model(model_name: str) -> int:
    normalized = str(model_name or "").strip().lower()
    if not normalized:
        return 16384
    if any(token in normalized for token in ("128k", "131072")):
        return 131072
    if any(token in normalized for token in ("64k", "65536")):
        return 65536
    if any(token in normalized for token in ("32k", "32768")):
        return 32768
    if any(token in normalized for token in ("16k", "16384")):
        return 16384
    if any(token in normalized for token in ("qwen2.5", "qwq", "gpt-oss:120b", "devstral-2:123b", "minimax-m2")):
        return 32768
    if any(token in normalized for token in ("llama3.1", "llama3.2", "mistral", "gemma3")):
        return 16384
    return 12288

def _extract_context_candidates(payload: Any) -> List[int]:
    candidates: List[int] = []
    seen_paths: set = set()

    def _collect(obj: Any, path: str = ""):
        if isinstance(obj, dict):
            for key, value in obj.items():
                next_path = f"{path}.{key}" if path else str(key)
                key_norm = str(key).lower()
                if any(marker in key_norm for marker in ("context", "num_ctx", "n_ctx", "window", "sequence", "position")):
                    digits = re.sub(r"[^\d]", "", str(value or ""))
                    if digits:
                        try:
                            number = int(digits)
                            if 2048 <= number <= 1048576 and next_path not in seen_paths:
                                candidates.append(number)
                                seen_paths.add(next_path)
                        except Exception:
                            pass
                _collect(value, next_path)
        elif isinstance(obj, list):
            for idx, item in enumerate(obj[:24]):
                _collect(item, f"{path}[{idx}]")

    _collect(payload)
    return candidates

async def get_model_context_window(model_name: str) -> int:
    normalized = str(model_name or "").strip()
    if not normalized:
        return _default_context_window_for_model(normalized)

    cached = _MODEL_CONTEXT_WINDOW_CACHE.get(normalized)
    now_ts = time.time()
    if cached and (now_ts - float(cached.get("checked_at", 0.0) or 0.0)) < 900:
        return int(cached.get("context_window", _default_context_window_for_model(normalized)))

    detected = _default_context_window_for_model(normalized)
    try:
        async with httpx.AsyncClient(timeout=8.0) as client:
            response = await client.post(
                "http://127.0.0.1:11434/api/show",
                json={"model": normalized},
            )
        if response.status_code == 200:
            candidates = _extract_context_candidates(response.json())
            if candidates:
                detected = max(candidates)
    except Exception as e:
        logger.debug(f"Context window lookup failed for {normalized}: {e}")

    detected = max(4096, min(int(detected), 1048576))
    _MODEL_CONTEXT_WINDOW_CACHE[normalized] = {
        "context_window": detected,
        "checked_at": now_ts,
    }
    return detected

def _estimate_message_token_budget(messages: List[Dict[str, str]]) -> int:
    total = 0
    for message in messages or []:
        total += estimate_tokens(str(message.get("content", "") or "")) + 12
    return total

async def fit_context_window_for_model(
    model_name: str,
    desired_ctx: int,
    messages: List[Dict[str, str]],
    max_output_tokens: int,
) -> int:
    available_ctx = await get_model_context_window(model_name)
    prompt_budget = _estimate_message_token_budget(messages)
    target_ctx = max(4096, prompt_budget + int(max_output_tokens) + 1024)
    return int(min(available_ctx, max(4096, min(int(desired_ctx), target_ctx))))

def _sanitize_history_messages(
    messages: Optional[List[Dict[str, Any]]],
    limit: int = MAX_SESSION_HISTORY_MESSAGES,
) -> List[Dict[str, str]]:
    sanitized: List[Dict[str, str]] = []
    for raw in messages or []:
        role = str((raw or {}).get("role", "user")).strip().lower()
        if role not in {"user", "assistant", "system"}:
            role = "user"
        content = _compact_memory_text((raw or {}).get("content", ""), max_chars=MAX_HISTORY_MESSAGE_CHARS)
        if not content:
            continue
        sanitized.append({"role": role, "content": content})
    return sanitized[-max(1, int(limit)):]

def _merge_history_messages(*message_groups: Optional[List[Dict[str, Any]]]) -> List[Dict[str, str]]:
    merged: List[Dict[str, str]] = []
    seen: set = set()
    for group in message_groups:
        for item in _sanitize_history_messages(group, limit=MAX_SESSION_HISTORY_MESSAGES):
            key = (item["role"], _normalize_cache_text(item["content"]))
            if not key[1] or key in seen:
                continue
            seen.add(key)
            merged.append(item)
    return merged[-MAX_SESSION_HISTORY_MESSAGES:]

def _build_conversation_summary_from_messages(
    messages: Optional[List[Dict[str, Any]]],
    max_chars: int = MAX_HISTORY_SUMMARY_CHARS,
) -> str:
    lines: List[str] = []
    seen: set = set()
    total = 0
    for item in _sanitize_history_messages(messages, limit=32):
        role = item.get("role", "user")
        label = "User" if role == "user" else "Assistant" if role == "assistant" else role.title()
        snippet = _compact_memory_text(item.get("content", ""), max_chars=220)
        normalized = _normalize_cache_text(snippet)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        line = f"- {label}: {snippet}"
        if total + len(line) + 1 > max_chars and lines:
            break
        lines.append(line[: max_chars - total])
        total += len(lines[-1]) + 1
    return "\n".join(lines).strip()

def _merge_context_summaries(*summaries: str, max_chars: int = MAX_HISTORY_SUMMARY_CHARS) -> str:
    lines: List[str] = []
    seen: set = set()
    total = 0
    for summary in summaries:
        for raw_line in re.split(r"[\r\n]+", str(summary or "")):
            line = _compact_memory_text(raw_line.lstrip("- ").strip(), max_chars=260)
            normalized = _normalize_cache_text(line)
            if not normalized or normalized in seen:
                continue
            seen.add(normalized)
            formatted = f"- {line}"
            if total + len(formatted) + 1 > max_chars and lines:
                return "\n".join(lines)
            lines.append(formatted)
            total += len(formatted) + 1
    return "\n".join(lines)

def _build_router_conversation_context(history_summary: str, recent_messages: List[Dict[str, str]]) -> str:
    parts: List[str] = []
    if history_summary:
        parts.append(f"Conversation summary:\n{history_summary}")
    if recent_messages:
        recent_lines = [
            f"- {msg.get('role', 'user')}: {_compact_memory_text(msg.get('content', ''), max_chars=180)}"
            for msg in recent_messages[-6:]
            if msg.get("content")
        ]
        if recent_lines:
            parts.append("Recent turns:\n" + "\n".join(recent_lines))
    return "\n\n".join(parts).strip()

def _should_bypass_response_cache(
    question: str,
    force_search: bool = False,
    deep_research: bool = False,
    has_image: bool = False,
) -> bool:
    return bool(
        force_search
        or deep_research
        or has_image
        or _query_requires_live_search(question)
    )


def _detect_text_language_for_cache(text: str) -> str:
    sample = str(text or "").strip()
    if not sample:
        return ""
    hinted = _detect_language_hint_from_text(sample)
    if hinted:
        return hinted

    normalized = _normalize_cache_text(sample)
    token_count = len(normalized.split())
    if LANGDETECT_AVAILABLE and (len(normalized) >= 18 or token_count >= 4):
        try:
            return _normalize_language_code(_langdetect_detect(sample))
        except LangDetectException:
            return ""
        except Exception:
            return ""
    return ""


def _script_family_hint(text: str) -> str:
    sample = str(text or "").strip()
    if not sample:
        return ""
    if re.search(r"[\u4e00-\u9fff]", sample):
        return "zh"
    if re.search(r"[\u3040-\u30ff]", sample):
        return "ja"
    if re.search(r"[\uac00-\ud7af]", sample):
        return "ko"
    if re.search(r"[\u0600-\u06ff]", sample):
        return "ar"
    if re.search(r"[\u0400-\u04ff]", sample):
        return "ru"
    if re.search(r"[\u0e00-\u0e7f]", sample):
        return "th"
    if re.search(r"[\u0900-\u097f]", sample):
        return "hi"
    if re.search(r"[A-Za-zÀ-ỹ]", sample):
        return "latin"
    return ""


def _is_cached_response_language_compatible(text: str, expected_language: str) -> bool:
    expected = _normalize_language_code(expected_language)
    detected = _detect_text_language_for_cache(text)
    if expected:
        if not detected:
            return True
        return detected == expected

    response_script = _script_family_hint(text)
    if not response_script:
        return True
    return response_script == "latin"

if sys.platform == 'win32':
    from asyncio.proactor_events import _ProactorBasePipeTransport
    def _silence_connection_lost(self, exc):
        try: self._call_connection_lost_original(exc)
        except (ConnectionResetError, ConnectionAbortedError, BrokenPipeError): pass
    if not hasattr(_ProactorBasePipeTransport, '_call_connection_lost_original'):
        _ProactorBasePipeTransport._call_connection_lost_original = _ProactorBasePipeTransport._call_connection_lost
        _ProactorBasePipeTransport._call_connection_lost = _silence_connection_lost

async def detect_language_with_model(text: str) -> str:
    """Detect the user's message language for answer-language selection."""
    if FORCE_ENGLISH:
        return "en"

    hinted = _detect_language_hint_from_text(text)
    if hinted:
        return hinted

    prompt = f"""Detect the language of this user text and return JSON only.
Use ISO 639-1 when possible.
Text: "{text}"

JSON format:
{{
  "language": "en",
  "confidence": 0.0-1.0
}}
"""
    try:
        async with httpx.AsyncClient(timeout=8.0) as client:
            response = await client.post(
                OLLAMA_URL,
                json={
                    "model": MODEL_ROUTER,
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False,
                    "options": {"temperature": 0.0, "num_predict": 50},
                    "format": "json",
                },
            )
        if response.status_code == 200:
            content = response.json().get("message", {}).get("content", "").strip()
            match = re.search(r"\{.*\}", content, re.DOTALL)
            if match:
                parsed = json.loads(match.group(0))
                language = _normalize_language_code(parsed.get("language", ""))
                if re.fullmatch(r"[a-z]{2,3}", language or ""):
                    return language
    except Exception as e:
        logger.warning(f"Language auto-detect failed: {e}")

    if LANGDETECT_AVAILABLE:
        try:
            return _normalize_language_code(_langdetect_detect(str(text or "").strip()))
        except LangDetectException:
            pass
        except Exception:
            pass

    return hinted or "en"

async def detect_query_context_with_model(text: str) -> Dict[str, Any]:
    """Detect interaction context with model-only classifier (no keyword routing rules)."""
    prompt = f"""Classify user input context and return JSON only.
Text: "{text}"

JSON format:
{{
  "is_image_request": true/false,
  "has_inline_analysis": true/false,
  "has_embedded_file_context": true/false,
  "confidence": 0.0-1.0
}}
"""
    try:
        async with httpx.AsyncClient(timeout=8.0) as client:
            response = await client.post(
                OLLAMA_URL,
                json={
                    "model": MODEL_ROUTER,
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False,
                    "options": {"temperature": 0.0, "num_predict": 80},
                },
            )
        if response.status_code == 200:
            content = response.json().get("message", {}).get("content", "").strip()
            match = re.search(r"\{.*\}", content, re.DOTALL)
            if match:
                parsed = json.loads(match.group(0))
                return {
                    "is_image_request": bool(parsed.get("is_image_request", False)),
                    "has_inline_analysis": bool(parsed.get("has_inline_analysis", False)),
                    "has_embedded_file_context": bool(parsed.get("has_embedded_file_context", False)),
                    "confidence": float(parsed.get("confidence", 0.5) or 0.5),
                }
    except Exception as e:
        logger.warning(f"Context auto-detect failed: {e}")
    return {
        "is_image_request": False,
        "has_inline_analysis": False,
        "has_embedded_file_context": False,
        "confidence": 0.0,
    }
def get_current_datetime() -> dict:
    from datetime import timezone, timedelta
    tz = timezone(timedelta(hours=7))
    now = datetime.now(tz=tz)
    weekdays = {
        0: "Thứ Hai",
        1: "Thứ Ba",
        2: "Thứ Tư",
        3: "Thứ Năm",
        4: "Thứ Sáu",
        5: "Thứ Bảy",
        6: "Chủ Nhật",
    }
    return {"date": now.strftime("%d/%m/%Y"), "time": now.strftime("%H:%M:%S"), "full": f"{weekdays[now.weekday()]}, {now.strftime('%d/%m/%Y %H:%M:%S')}"}

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    question: str
    history: List[Message] = Field(default_factory=list)
    history_summary: str = ""
    user_id: str = "default_user"
    session_id: Optional[str] = None  # for fast session context (temporary)
    language: Optional[str] = None
    force_search: bool = False
    deep_research: bool = False  # Deep research mode for comprehensive search
    confirmed_plan: Optional[List[str]] = None

class ClickFeedback(BaseModel):
    query: str
    doc_url: str
    source: str = "hybrid"
    user_id: str = "default_user"
    dwell_ms: int = 0
    clicked: bool = True

class DummySearchEngine:
    async def smart_search(self, *args, **kwargs): return ""
    async def smart_search_multi(self, *args, **kwargs): return {"urls": [], "context": ""}
    def get_engine_info(self): return {}
    async def close(self): pass
try:
    from web_mcp_server import WebMCPServer
    from database_mcp_server import DatabaseMCPServer
    MCP_BACKENDS_AVAILABLE = True
except Exception as e:
    MCP_BACKENDS_AVAILABLE = False
    WebMCPServer = None
    DatabaseMCPServer = None
    logger.info(f"MCP backend optional modules unavailable ({e}). Running without MCP backends.")

def _load_shared_search_engine():
    candidate_paths = [
        Path(__file__).resolve().parent / "search_engine.py",
        Path(__file__).resolve().parent.parent / "Skemma-main (1)" / "Skemma-main" / "Skemma-main (1) (1)" / "search_engine.py",
    ]
    last_error = None
    for engine_path in candidate_paths:
        try:
            if not engine_path.exists():
                continue
            module_name = f"skemi_shared_search_engine_{abs(hash(str(engine_path)))}"
            spec = importlib.util.spec_from_file_location(module_name, engine_path)
            if spec is None or spec.loader is None:
                continue
            module = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(module)
            engine_cls = getattr(module, "SmartSearchEngine", None)
            if engine_cls is None:
                continue
            engine = engine_cls()
            if engine_path.parent == Path(__file__).resolve().parent:
                logger.info("Loaded canonical SmartSearchEngine from local Skemi repo.")
            else:
                logger.info(f"Loaded shared SmartSearchEngine from: {engine_path}")
            return engine
        except Exception as exc:
            last_error = exc
    raise RuntimeError(last_error or "SmartSearchEngine module not found")

try:
    search_engine = _load_shared_search_engine()
except Exception as e:
    search_engine = DummySearchEngine()
    logger.warning(f"SmartSearchEngine unavailable, using DummySearchEngine fallback: {e}")

web_mcp_backend = WebMCPServer() if MCP_BACKENDS_AVAILABLE else None
db_mcp_backend = DatabaseMCPServer() if MCP_BACKENDS_AVAILABLE else None

def _normalize_analysis_mode(mode: Optional[str]) -> str:
    normalized = str(mode or "assistant").strip().lower()
    if normalized not in ANALYSIS_MODES:
        return "assistant"
    return normalized

def _analysis_style_instructions(mode: str) -> str:
    if mode == "structured":
        return (
            "Use a strict structure with headings: Summary, Findings, Citations, Uncertainty, Next Step. "
            "Keep each section short and factual."
        )
    if mode == "concise":
        return (
            "Be concise: 3-5 bullets total. Include at least one citation and one uncertainty note."
        )
    return (
        "Use natural assistant style, concise and practical. "
        "No rigid UI inventory unless explicitly requested."
    )

def _layout_summary_text(layout: Dict[str, Any]) -> str:
    summary = layout.get("summary", {}) if isinstance(layout, dict) else {}
    if not summary:
        return "No layout summary available."
    hint = summary.get("layout_hint", "unknown")
    region_count = int(summary.get("region_count", 0) or 0)
    text_regions = int(summary.get("text_region_count", 0) or 0)
    non_text_regions = int(summary.get("non_text_region_count", 0) or 0)
    return (
        f"layout_hint={hint}; region_count={region_count}; "
        f"text_regions={text_regions}; non_text_regions={non_text_regions}"
    )

def _build_ocr_evidence(spans: List[Dict[str, Any]], max_spans: int = 10) -> str:
    if not spans:
        return "(no OCR spans)"
    ranked = sorted(
        spans,
        key=lambda s: float(s.get("confidence", 0.0) or 0.0) * max(1, len(str(s.get("text", "")))),
        reverse=True,
    )
    lines: List[str] = []
    for idx, span in enumerate(ranked[:max_spans], start=1):
        text = str(span.get("text", "")).strip()
        if not text:
            continue
        box = span.get("bbox", {}) or {}
        lines.append(
            f"[S{idx}] text='{text[:140]}' conf={float(span.get('confidence', 0.0) or 0.0):.2f} "
            f"bbox=({int(box.get('x', 0))},{int(box.get('y', 0))},{int(box.get('w', 0))},{int(box.get('h', 0))})"
        )
    return "\n".join(lines) if lines else "(no OCR spans)"

def _ocr_uncertainty_label(confidence: float, text: str) -> str:
    chars = len((text or "").strip())
    if confidence >= 0.8 and chars >= 80:
        return "low"
    if confidence >= 0.55 and chars >= 30:
        return "medium"
    return "high"

@asynccontextmanager
async def lifespan(app: FastAPI):
    global OCR_STARTUP_DIAGNOSTICS
    logger.info("Server starting up...")
    
    if CACHE_AVAILABLE:
        try:
            advanced_cache.start_background_cleanup()
            logger.info(" Cache system started")
        except Exception as e:
            logger.info(f" Cache system start note: {e}")
    
    worker_task = asyncio.create_task(ai_worker(), name="ai_worker")
    learning_task = asyncio.create_task(auto_learning_loop(), name="auto_learning_loop")
    await check_ollama_status()
    extra_startup = getattr(getattr(app, "state", None), "skemi_extra_startup", None)
    if callable(extra_startup):
        try:
            result = extra_startup()
            if inspect.isawaitable(result):
                await result
        except Exception as exc:
            logger.error(f"Extra startup hook failed: {exc}")
    
    yield
    
    try:
        extra_shutdown = getattr(getattr(app, "state", None), "skemi_extra_shutdown", None)
        if callable(extra_shutdown):
            try:
                result = extra_shutdown()
                if inspect.isawaitable(result):
                    await result
            except Exception as exc:
                logger.error(f"Extra shutdown hook failed: {exc}")
        if CACHE_AVAILABLE:
            try:
                advanced_cache.stop_background_cleanup()
                logger.info(" Cache system stopped")
            except Exception as e:
                logger.error(f" Cache system stop error: {e}")
        
        await search_engine.close()
        logger.info(" Search engine closed")
        worker_task.cancel()
        learning_task.cancel()
        try:
            await worker_task
        except asyncio.CancelledError:
            logger.info(" AI worker task cancelled.")
        try:
            await learning_task
        except asyncio.CancelledError:
            logger.info(" Auto-learning task cancelled.")
    except Exception as e:
        logger.error(f"Error during server shutdown: {e}")

app = FastAPI(title="Skemi AI Core DEBUG", version="4.5", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

_HEALTH_CHECK_LOGGED = False
_QUIET_REQUEST_PATHS = {
    "/health",
    "/api/health",
    "/api/computer/status",
    "/api/computer/history",
    "/api/computer/webrtc/offer",
    "/api/global/status",
    "/api/local-computer/status",
    "/api/local-computer/events",
    "/favicon.ico",
    "/end_session",
    "/Computer.html",
}
_QUIET_REQUEST_PREFIXES = (
    "/api/computer/live",
    "/api/computer/surface",
    "/api/computer/history",
    "/api/computer/webrtc",
    "/api/computer/stream",
    "/api/local-computer/live",
    "/api/local-computer/mjpeg",
    "/api/local-computer/stream",
    "/api/local-computer/voice",
    "/ws/computer/surface",
    "/Css/",
    "/Js/",
)

@app.middleware("http")
async def log_requests(request: Request, call_next):
    global _HEALTH_CHECK_LOGGED
    path = request.url.path
    debug_enabled = str(os.getenv("SKEMI_DEBUG_REQUESTS", "") or "").strip().lower() in {"1", "true", "yes", "on"}
    if not debug_enabled:
        return await call_next(request)
    if path in _QUIET_REQUEST_PATHS or any(path.startswith(prefix) for prefix in _QUIET_REQUEST_PREFIXES):
        if path == "/health" and not _HEALTH_CHECK_LOGGED:
            # print(f" \n>>> [PYTHON DEBUG] Incoming: {request.method} {path} (further health logs silenced) <<<")
            # logger.info(f"🔍 [PYTHON DEBUG] Request to {path} (further health logs silenced)")
            _HEALTH_CHECK_LOGGED = True
        return await call_next(request)
    if path == "/health":
        if not _HEALTH_CHECK_LOGGED:
            # print(f" \n>>> [PYTHON DEBUG] Incoming: {request.method} {path} (further health logs silenced) <<<")
            # logger.info(f"🔍 [PYTHON DEBUG] Request to {path} (further health logs silenced)")
            _HEALTH_CHECK_LOGGED = True
    else:
        # Silence spammy endpoints
        spam_endpoints = ["/api/local-computer/stream", "/api/local-computer/status", "/api/computer/history/manifest", "/api/computer/mjpeg"]
        if any(endpoint in path for endpoint in spam_endpoints):
            return await call_next(request)
        else:
            # print(f" \n>>> [PYTHON DEBUG] Incoming: {request.method} {path} <<<")
            # logger.info(f"🔍 [PYTHON DEBUG] Request to {path}")
            pass
    return await call_next(request)

@app.get("/health")
async def health():
    return {
        "status": "ok", 
        "version": "4.5", 
        "name": "Skemi",
        "ocr_ready": bool(OCR_STARTUP_DIAGNOSTICS.get("ready", False)),
    }

@app.get("/ocr/diagnostics")
async def ocr_diagnostics():
    return {
        "ocr_available": MULTI_TIER_OCR_AVAILABLE,
        "startup_diagnostics": OCR_STARTUP_DIAGNOSTICS,
    }

async def analyze_image_with_4tier_ocr_bundle(image_data: bytes, language: str = "eng") -> Dict[str, Any]:
    """Run multi-tier OCR and return structured output for grounding/citations."""
    if not MULTI_TIER_OCR_AVAILABLE:
        return {
            "success": False,
            "tier": 0,
            "engine": "none",
            "text": "",
            "confidence": 0.0,
            "spans": [],
            "layout": {},
            "all_results": [],
            "error": "multi-tier OCR not available",
        }
    try:
        result = await run_multi_tier_ocr(image_data, language)
        if not isinstance(result, dict):
            logger.error(f"4-tier OCR returned unexpected type: {type(result).__name__}")
            return {
                "success": False,
                "tier": 0,
                "engine": "none",
                "text": "",
                "confidence": 0.0,
                "spans": [],
                "layout": {},
                "all_results": [],
                "error": "invalid response format",
            }
        return result
    except Exception as e:
        logger.error(f"4-Tier OCR error: {e}")
        return {
            "success": False,
            "tier": 0,
            "engine": "none",
            "text": "",
            "confidence": 0.0,
            "spans": [],
            "layout": {},
            "all_results": [],
            "error": str(e),
        }

async def analyze_image_with_4tier_ocr(image_data: bytes, language: str = "eng") -> str:
    """Run multi-tier OCR and return extracted text only."""
    result = await analyze_image_with_4tier_ocr_bundle(image_data, language)
    if result.get("success"):
        extracted_text = (result.get("text") or "").strip()
        logger.info(
            f"OCR successful - Tier {result.get('tier')}, "
            f"Engine: {result.get('engine')}, chars={len(extracted_text)}, "
            f"confidence={float(result.get('confidence') or 0):.2f}"
        )
        return extracted_text
    logger.error(f"All OCR tiers failed: {result.get('error', 'no details')}")
    return ""

# Last-resort LOCAL fallback list, used when the primary model fails (e.g. a
# ":cloud" model returns HTTP 401 because Ollama-Cloud auth is missing).
# DELIBERATELY EMPTY by default: check_ollama_status() is the single source of
# truth and repopulates this at startup with the models ACTUALLY installed
# (see the `global LOCAL_FALLBACK_MODELS` block). Hardcoding names here
# (previously qwen2.5:3b/1.5b) caused endless HTTP 404 spam on machines where
# those models were never pulled. Empty means "no local fallback" so generation
# fails fast with a clear 'model unavailable' message instead of probing ghosts.
LOCAL_FALLBACK_MODELS: List[str] = []


async def _raw_generate_once(model: str, prompt: str, timeout: float, num_predict: int) -> str:
    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            response = await client.post(
                GENERATE_URL,
                json={
                    "model": model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {"temperature": 0.25, "num_predict": num_predict},
                },
            )
        if response.status_code != 200:
            logger.error(f"Single-shot generation failed ({model}): {response.status_code}")
            return ""
        full_response = ""
        prompt_tokens = 0
        completion_tokens = 0
        for line in response.content.decode("utf-8", errors="replace").splitlines():
            if not line.strip():
                continue
            try:
                json_part = json.loads(line)
                if "response" in json_part:
                    full_response += json_part["response"]
                # ollama reports token counts on the final ("done") chunk.
                if "prompt_eval_count" in json_part:
                    prompt_tokens = int(json_part.get("prompt_eval_count") or 0)
                if "eval_count" in json_part:
                    completion_tokens = int(json_part.get("eval_count") or 0)
                if json_part.get("done"):
                    break
            except json.JSONDecodeError:
                continue
        # Attribute token spend to the active account (no-op if unavailable).
        try:
            import entitlements as _ent
            _ent.note_model_usage(prompt_tokens, completion_tokens)
        except Exception:
            pass
        return full_response.strip()
    except Exception as e:
        logger.error(f"Single-shot generation error ({model}): {e}")
        return ""


async def _generate_text_once(model: str, prompt: str, timeout: float = 60.0, num_predict: int = 900) -> str:
    # Build ordered candidate list: requested model first, then local fallbacks.
    candidates: List[str] = []
    if model:
        candidates.append(model)
    for fb in LOCAL_FALLBACK_MODELS:
        if fb not in candidates:
            candidates.append(fb)
    for idx, candidate in enumerate(candidates):
        text = await _raw_generate_once(candidate, prompt, timeout, num_predict)
        if text:
            if idx > 0:
                logger.warning(f"Primary model '{candidates[0]}' failed/empty; used local fallback '{candidate}'")
            return text
    return ""

async def analyze_image_with_vision_model(
    image_data: bytes,
    output_language: str = "en",
    analysis_mode: str = "assistant",
) -> str:
    """Analyze image using a 4-tier OCR pipeline, then vision fallback."""
    if not image_data:
        return "No image data provided."
    mode = _normalize_analysis_mode(analysis_mode)

    # Run OCR pipeline first.
    ocr_result = ""
    ocr_bundle: Dict[str, Any] = {}
    ocr_spans: List[Dict[str, Any]] = []
    ocr_layout: Dict[str, Any] = {}
    ocr_confidence = 0.0
    if MULTI_TIER_OCR_AVAILABLE:
        logger.info("Using 4-tier OCR pipeline for image analysis")
        try:
            ocr_bundle = await analyze_image_with_4tier_ocr_bundle(image_data)
            ocr_result = str(ocr_bundle.get("text") or "").strip()
            ocr_spans = list(ocr_bundle.get("spans") or [])
            ocr_layout = dict(ocr_bundle.get("layout") or {})
            ocr_confidence = float(ocr_bundle.get("confidence") or 0.0)
            if not (ocr_result and len(ocr_result.strip()) >= OCR_MIN_CHARS):
                logger.warning("4-tier OCR did not return usable text, trying vision fallback")
        except Exception as e:
            logger.error(f"4-tier OCR error: {e}")
            ocr_result = ""
    
    # Fallback to Ollama vision model
    if not MODEL_VISION:
        return "No vision model is configured. Install a multimodal model and try again."
    
    logger.info(f"Fallback to vision model: {MODEL_VISION}")
    encoded_image = base64.b64encode(image_data).decode('utf-8')

    language_instruction = "Trả lời bằng tiếng Việt." if output_language == "vi" else "Answer in English."
    newline = "\n"
    prompt_text = f"""Bạn là AI phân tích hình ảnh. Hãy suy luận theo từng bước (step-by-step reasoning).
{language_instruction}

Nhiệm vụ: Phân tích chi tiết hình ảnh này.

Hướng dẫn:
1. MÔ TẢ TỔNG QUAN: Ảnh chứa gì? (giao diện, văn bản, biểu đồ, người, sản phẩm, mã code, v.v.)
2. NỘI DUNG CHI TIẾT: Đọc và trích xuất MỌI văn bản, số liệu, nhãn hiệu, màu sắc, bố cục.
3. PHÂN TÍCH SÂU: Ý nghĩa, mục đích, ngữ cảnh của ảnh.
4. NHẬN XÉT: Đưa ra nhận xét, đề xuất nếu phù hợp.

Chế độ phân tích: {mode}

{('Văn bản OCR tham khảo (có thể không chính xác):' + newline + ocr_result[:5000]) if ocr_result else '(Không có OCR)'}

Hãy trả lời đầy đủ, chi tiết và suy luận kỹ càng. KHÔNG bỏ qua bất kỳ thông tin nào trong ảnh.
"""

    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(GENERATE_URL, json={
                "model": MODEL_VISION,
                "prompt": prompt_text,
                "images": [encoded_image],
                "stream": False,
                "options": {
                    "temperature": 0.3,
                    "num_predict": 3000
                }
            })
            
            if response.status_code == 200:
                full_response = ""
                for line in response.content.decode("utf-8", errors="replace").splitlines():
                    if line.strip():
                        try:
                            json_part = json.loads(line)
                            if "response" in json_part:
                                full_response += json_part["response"] + "\n"
                            if json_part.get("done"):
                                break
                        except json.JSONDecodeError:
                            continue
                
                result = full_response.strip()
                if result and result != "NO_TEXT":
                    logger.info("Vision model analysis successful")
                    return result
                if ocr_result:
                    fallback_prompt = f"""You are given OCR text and OCR evidence spans.
Return only the final answer with short rationale.
Do not reveal chain-of-thought.
Mode: {mode}. {_analysis_style_instructions(mode)}

OCR text:
{ocr_result[:5000]}

Evidence:
{evidence_block}

Required:
- Include citations with [Sx] from evidence list when possible.
- Include an uncertainty label: low/medium/high."""
                    router_summary = await _generate_text_once(MODEL_ROUTER, fallback_prompt, timeout=40.0, num_predict=700)
                    if router_summary:
                        return router_summary
                return "No readable text detected in the image."
            else:
                if ocr_result:
                    fallback_prompt = f"""Interpret OCR text for user.
Return final answer only with short rationale.
Mode: {mode}. {_analysis_style_instructions(mode)}
OCR text:
{ocr_result[:5000]}
Evidence:
{evidence_block}
Uncertainty baseline: {uncertainty_level}"""
                    router_summary = await _generate_text_once(MODEL_ROUTER, fallback_prompt, timeout=40.0, num_predict=700)
                    if router_summary:
                        return router_summary
                return f"Vision model error ({response.status_code}): {response.text[:200]}"
                
    except Exception as e:
        logger.error(f"Vision model error: {e}")
        if ocr_result:
            fallback_prompt = f"""Interpret OCR text for user with grounded citations.
Keep chain-of-thought hidden.
Mode: {mode}. {_analysis_style_instructions(mode)}
OCR text:
{ocr_result[:5000]}
Evidence:
{evidence_block}
Uncertainty baseline: {uncertainty_level}"""
            router_summary = await _generate_text_once(MODEL_ROUTER, fallback_prompt, timeout=40.0, num_predict=700)
            if router_summary:
                return router_summary
        return f"Image analysis failed: {str(e)}"

async def analyze_file_content_with_router(
    file_name: str,
    extracted_text: str,
    file_type: str,
    analysis_mode: str = "assistant",
) -> str:
    """Analyze textual file content with the router model."""
    if not extracted_text or len(extracted_text.strip()) < 10:
        return f"File {file_name} does not contain enough extracted text for analysis."
    mode = _normalize_analysis_mode(analysis_mode)

    prompt_text = f"""Phân tích file: {file_name} (loại: {file_type})

Hãy suy luận từng bước (step-by-step reasoning) để phân tích nội dung này.
    
NỘI DUNG TRÍCH XUẤT:
{extracted_text[:8000]}

YÊU CẦU PHÂN TÍCH:
1. Hiểu ý định/mục đích của nội dung
2. Tóm tắt ý nghĩa bằng ngôn ngữ dễ hiểu
3. Trích xuất các dữ kiện, thông tin quan trọng
4. Đề xuất hành động tiếp theo

Phong cách: {mode}. {_analysis_style_instructions(mode)}
Trả lời chi tiết, đầy đủ, dùng markdown nếu cần."""

    payload = {
        "model": MODEL_ROUTER,
        "prompt": prompt_text,
        "stream": False,
        "options": {
            "temperature": 0.3,
            "num_predict": 2000
        }
    }

    try:
        async with httpx.AsyncClient(timeout=45.0) as client:
            logger.info(f"Analyzing file {file_name} with {MODEL_ROUTER}")
            response = await client.post(GENERATE_URL, json=payload)

            if response.status_code == 200:
                full_response = ""
                for line in response.content.decode("utf-8", errors="replace").splitlines():
                    if line.strip():
                        try:
                            json_part = json.loads(line)
                            full_response += json_part.get("response", "")
                            if json_part.get("done"):
                                break
                        except json.JSONDecodeError:
                            logger.warning(f"Couldn't decode JSON line from model: {line}")
                            continue
                
                logger.info(f" File analysis successful for {file_name}")
                return full_response.strip()
            else:
                error_msg = f"File analysis error ({response.status_code}): {response.text}"
                logger.error(error_msg)
                return error_msg
    except httpx.TimeoutException:
        error_msg = f"File analysis timeout with {MODEL_ROUTER}"
        logger.error(error_msg)
        return error_msg
    except Exception as e:
        logger.error(f"File analysis failed: {e}", exc_info=True)
        return f"File analysis failed: {e}"

async def detect_content_language(extracted_text: str, fallback_hint: str = "") -> str:
    sample = (extracted_text or "").strip()
    if len(sample) > 400:
        sample = sample[:400]
    if not sample:
        sample = fallback_hint or "text"
    return await detect_language_with_model(sample)

def _sse_pack(payload: dict) -> str:
    return f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"


def _event_ts() -> str:
    return datetime.utcnow().isoformat(timespec="seconds") + "Z"


def _safe_domain_from_url(url: str) -> str:
    try:
        return urlparse(str(url or "")).netloc.lower().replace("www.", "")
    except Exception:
        return ""


def _favicon_url_for_domain(domain: str) -> str:
    clean = str(domain or "").strip()
    if not clean:
        return ""
    return f"https://icons.duckduckgo.com/ip3/{clean}.ico"


def _build_source_event_payload(item: Dict[str, Any]) -> Dict[str, Any]:
    url = str(item.get("url") or "").strip()
    domain = _safe_domain_from_url(url)
    return {
        "title": str(item.get("title") or domain or "Source").strip(),
        "url": url,
        "domain": domain,
        "favicon_url": _favicon_url_for_domain(domain),
        "snippet": str(item.get("snippet") or "").strip(),
    }


def _fallback_stream_status_text(
    *,
    event_name: str,
    phase: str,
    label: str,
    detail: str,
    query: str,
    provider: str,
    result_count: Optional[int],
    source: Optional[Dict[str, Any]],
    language: str = "",
) -> str:
    lang = str(language or "").strip().lower()
    query_short = str(query or "").strip()
    if len(query_short) > 44:
        query_short = query_short[:41].rstrip() + "..."
    source_title = str((source or {}).get("title") or (source or {}).get("domain") or "").strip()
    if len(source_title) > 34:
        source_title = source_title[:31].rstrip() + "..."

    if lang.startswith("vi"):
        phase_map = {
            "analyzing": "Đang hiểu yêu cầu",
            "routing": "Đang chọn hướng xử lý",
            "planning_search": "Đang lên lộ trình tìm kiếm",
            "searching": f"Đang quét web cho {query_short}" if query_short else "Đang quét web",
            "fetching": "Đang đọc và lọc nguồn",
            "ranking": "Đang chọn nguồn tốt nhất",
            "summarizing": "Đang chắt lọc thông tin",
            "streaming_reply": "Đang viết câu trả lời",
        }
        if event_name == "provider_result" and isinstance(result_count, int):
            return f"Đã gom {result_count} kết quả đầu"
        if event_name == "query_started" and query_short:
            return f"Đang tra cứu: {query_short}"
        if event_name in {"source_candidate", "source_accepted"} and source_title:
            return f"Đang soi {source_title}"
        return phase_map.get(phase) or label or detail or "Đang xử lý"

    phase_map = {
        "analyzing": "Reading the request",
        "routing": "Choosing the best path",
        "planning_search": "Sketching the search path",
        "searching": f"Scanning the web for {query_short}" if query_short else "Scanning the web",
        "fetching": "Reading and filtering sources",
        "ranking": "Picking the strongest sources",
        "summarizing": "Distilling the evidence",
        "streaming_reply": "Writing the answer",
    }
    if event_name == "provider_result" and isinstance(result_count, int):
        return f"Gathered {result_count} live hits"
    if event_name == "query_started" and query_short:
        return f"Looking up {query_short}"
    if event_name in {"source_candidate", "source_accepted"} and source_title:
        return f"Checking {source_title}"
    return phase_map.get(phase) or label or detail or "Working"


def _clean_dynamic_status_text(value: Any, fallback: str) -> str:
    text = str(value or "").strip()
    if not text:
        return fallback
    text = re.sub(r"^```(?:text)?|```$", "", text, flags=re.IGNORECASE | re.MULTILINE).strip()
    text = re.sub(r"^[\"'`•\-\s]+", "", text)
    text = re.sub(r"[\"'`]+$", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return fallback
    if len(text) > 72:
        text = text[:69].rstrip(" ,.;:") + "..."
    return text


def _build_backend_failure_message(kind: str, language: str = "en", status_code: Optional[int] = None) -> str:
    lang = _normalize_language_code(language or "en")
    if lang == "vi":
        if kind == "unavailable":
            if status_code:
                return f"Mình chưa thể hoàn tất phản hồi vì dịch vụ AI đang tạm thời gián đoạn (HTTP {status_code}). Bạn hãy thử lại sau ít phút hoặc gửi lại câu hỏi ngắn hơn."
            return "Mình chưa thể hoàn tất phản hồi vì dịch vụ AI đang tạm thời gián đoạn. Bạn hãy thử lại sau ít phút hoặc gửi lại câu hỏi ngắn hơn."
        if kind == "timeout":
            return "Mình đang mất nhiều thời gian hơn bình thường để tạo phản hồi. Bạn hãy thử lại sau ít phút hoặc gửi lại câu hỏi ngắn gọn hơn."
        if kind == "empty":
            return "Mình chưa nhận được nội dung phản hồi từ mô hình cho câu hỏi này. Bạn hãy thử lại hoặc diễn đạt lại câu hỏi ngắn hơn."
        return "Mình gặp sự cố khi tạo phản hồi cho yêu cầu này. Bạn hãy thử lại sau ít phút."

    if kind == "unavailable":
        if status_code:
            return f"I could not finish the reply because the AI service is temporarily unavailable (HTTP {status_code}). Please try again in a moment or resend a shorter question."
        return "I could not finish the reply because the AI service is temporarily unavailable. Please try again in a moment or resend a shorter question."
    if kind == "timeout":
        return "I am taking longer than usual to complete this reply. Please try again in a moment or resend a shorter question."
    if kind == "empty":
        return "I did not receive any usable reply from the model for this question. Please try again or resend a shorter version."
    return "I ran into a temporary problem while generating this reply. Please try again in a moment."


async def _generate_dynamic_stream_status(
    *,
    event_name: str,
    phase: str,
    label: str,
    detail: str,
    query: str,
    provider: str,
    result_count: Optional[int],
    source: Optional[Dict[str, Any]],
    language: str = "",
) -> str:
    fallback = _fallback_stream_status_text(
        event_name=event_name,
        phase=phase,
        label=label,
        detail=detail,
        query=query,
        provider=provider,
        result_count=result_count,
        source=source,
        language=language,
    )

    if event_name not in {"phase_updated", "query_started", "provider_result", "ranking_updated", "summary_started"}:
        return fallback
    if not STATUS_MODEL_ACTIVE:
        return fallback

    cache_key = "|".join(
        [
            str(language or "en").strip().lower() or "en",
            event_name,
            str(phase or "").strip().lower(),
            str(label or "").strip().lower(),
            str(query or "").strip().lower()[:120],
            str(provider or "").strip().lower(),
            str(result_count if result_count is not None else ""),
        ]
    )
    if cache_key in _status_text_cache:
        return _status_text_cache[cache_key]

    lang = str(language or "").strip().lower()
    if lang.startswith("vi"):
        language_instruction = "Write the UI status in Vietnamese."
    elif lang.startswith("en"):
        language_instruction = "Write the UI status in English."
    else:
        language_instruction = f"Write the UI status in the user's language ({lang}), or English if needed."

    prompt = f"""You generate ultra-short live status text for an AI chat interface.
{language_instruction}

Style:
- 2 to 7 words
- concise, elegant, alive
- present tense or progressive when natural
- no emojis
- no quotes
- no markdown
- do not mention models, tools, or internal systems

Context:
- Event: {event_name}
- Phase: {phase or 'unknown'}
- Label: {label or 'none'}
- Detail: {detail or 'none'}
- Query: {query or 'none'}
- Provider: {provider or 'none'}
- Result count: {result_count if result_count is not None else 'none'}
- Source: {((source or {}).get('title') or (source or {}).get('domain') or 'none')}

Return only the short status text."""

    try:
        raw = await _generate_text_once(MODEL_STATUS, prompt, timeout=3.0, num_predict=24)
        cleaned = _clean_dynamic_status_text(raw, fallback)
        if cleaned:
            _status_text_cache[cache_key] = cleaned
            if len(_status_text_cache) > 256:
                oldest_key = next(iter(_status_text_cache))
                _status_text_cache.pop(oldest_key, None)
            return cleaned
    except Exception:
        pass

    _status_text_cache[cache_key] = fallback
    return fallback


async def _emit_stream_event(
    result_queue: asyncio.Queue,
    event_name: str,
    *,
    phase: str = "",
    label: str = "",
    detail: str = "",
    progress: Optional[float] = None,
    query: str = "",
    provider: str = "",
    result_count: Optional[int] = None,
    source: Optional[Dict[str, Any]] = None,
    language: str = "",
) -> None:
    status_text = await _generate_dynamic_stream_status(
        event_name=event_name,
        phase=phase,
        label=label,
        detail=detail,
        query=query,
        provider=provider,
        result_count=result_count,
        source=source,
        language=language,
    )
    payload: Dict[str, Any] = {
        "event": event_name,
        "phase": phase,
        "label": label,
        "detail": detail,
        "status": status_text,
        "ts": _event_ts(),
    }
    if progress is not None:
        payload["progress"] = progress
    if query:
        payload["query"] = query
    if provider:
        payload["provider"] = provider
    if result_count is not None:
        payload["result_count"] = int(result_count)
    if source:
        payload["source"] = source
    try:
        summary_bits = [event_name]
        if phase:
            summary_bits.append(f"phase={phase}")
        if label:
            summary_bits.append(f"label={label}")
        if status_text:
            summary_bits.append(f"status={status_text[:120]}")
        if query:
            summary_bits.append(f"query={query[:120]}")
        if provider:
            summary_bits.append(f"provider={provider}")
        if result_count is not None:
            summary_bits.append(f"results={int(result_count)}")
        if source and source.get("url"):
            summary_bits.append(f"url={str(source.get('url'))[:140]}")
        logger.info("[STREAM] " + " | ".join(summary_bits))
    except Exception:
        pass
    await result_queue.put(payload)

def _sse_single_response(payload: dict) -> StreamingResponse:
    async def _single_event():
        yield _sse_pack(payload)
    return StreamingResponse(_single_event(), media_type="text/event-stream")

@app.post("/parse_file_stream")
async def parse_file_stream(
    file: UploadFile = File(...),
    analysis_mode: str = Form("assistant"),
):
    if not file:
        raise HTTPException(status_code=400, detail="No file uploaded")
    normalized_mode = _normalize_analysis_mode(analysis_mode)

    async def event_stream():
        try:
            yield _sse_pack({"progress": 5, "status": "Receiving file..."})
            data = await file.read()
            if len(data) > MAX_FILE_BYTES:
                yield _sse_pack({"error": "File too large (max 50MB)", "complete": True})
                return

            yield _sse_pack({"progress": 15, "status": "Extracting content..."})
            extracted, detected_type = extract_text_by_extension(file.filename or "", data)
            extracted = extracted or ""
            
            # Check token limit for file content
            if estimate_tokens(extracted) > MAX_FILE_TOKENS:
                extracted = extracted[:MAX_FILE_TOKENS * 4] + "\n\n[Nội dung quá dài, đã rút gọn]"
                yield _sse_pack({"warning": "Nội dung file quá dài, đã rút gọn."})
            
            extracted, truncated = _truncate_text(extracted, 12000)

            yield _sse_pack({"progress": 35, "status": "Preparing analysis..."})

            analysis_task = None
            file_analysis = ""

            if detected_type == "image":
                target_language = await detect_content_language(extracted, file.filename or "")
                if _has_usable_ocr(extracted):
                    yield _sse_pack({"progress": 60, "status": "Basic OCR text extracted"})
                logger.info("Image file detected, starting 4-tier image analysis (stream)...")
                analysis_task = asyncio.create_task(
                    analyze_image_with_vision_model(
                        data,
                        output_language=target_language,
                        analysis_mode=normalized_mode,
                    )
                )
            else:
                logger.info(f"Document file detected ({detected_type}), starting AI analysis (stream)...")
                analysis_task = asyncio.create_task(
                    analyze_file_content_with_router(
                        file.filename or "",
                        extracted,
                        detected_type,
                        analysis_mode=normalized_mode,
                    )
                )

            if analysis_task:
                progress = 40
                yield _sse_pack({"progress": progress, "status": "Analyzing..."})

                while not analysis_task.done():
                    await asyncio.sleep(0.6)
                    if progress < 90:
                        progress += 1
                        yield _sse_pack({"progress": progress, "status": "Analyzing..."})

                file_analysis = await analysis_task

            yield _sse_pack({"progress": 100, "status": "Completed"})
            result = {
                "file_name": file.filename,
                "file_type": detected_type,
                "extracted_text": extracted,
                "truncated": truncated,
                "chars": len(extracted),
                "ai_analysis": file_analysis,
                "analysis_mode": normalized_mode,
            }
            yield _sse_pack({"result": result, "complete": True})
        except Exception as e:
            logger.error(f"Parse file stream error: {e}", exc_info=True)
            yield _sse_pack({"error": f"File processing error: {e}", "complete": True})

    return StreamingResponse(event_stream(), media_type="text/event-stream")

@app.post("/parse_file")
async def parse_file(
    file: UploadFile = File(...),
    analysis_mode: str = Form("assistant"),
):
    if not file:
        raise HTTPException(status_code=400, detail="No file uploaded")
    normalized_mode = _normalize_analysis_mode(analysis_mode)

    data = await file.read()
    if len(data) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="File too large (max 50MB)")

    # Generate file hash for caching
    file_hash = generate_file_hash(data + normalized_mode.encode("utf-8"))
    extracted, detected_type = extract_text_by_extension(file.filename or "", data)
    extracted = extracted or ""
    extracted, truncated = _truncate_text(extracted, 12000)

    # Check cache first
    cached_analysis = await get_cached_file_analysis(file_hash, detected_type)
    if cached_analysis:
        logger.info(f" File Cache HIT for {file.filename}")
        return {
            "file_name": file.filename,
            "file_type": detected_type,
            "extracted_text": cached_analysis.get("extracted_text", extracted),
            "truncated": cached_analysis.get("truncated", truncated),
            "chars": len(extracted),
            "ai_analysis": cached_analysis.get("ai_analysis", ""),
            "analysis_mode": normalized_mode,
            "cache_hit": True
        }

    file_analysis = ""

    # Phân tích bằng AI model (ưu tiên vision cho ảnh).
    if detected_type == "image":
        target_language = await detect_content_language(extracted, file.filename or "")
        logger.info("Image file detected, starting 4-tier image analysis...")
        file_analysis = await analyze_image_with_vision_model(
            data,
            output_language=target_language,
            analysis_mode=normalized_mode,
        )
    else:
        logger.info(f"Document file detected ({detected_type}), starting AI analysis...")
        file_analysis = await analyze_file_content_with_router(
            file.filename or "",
            extracted,
            detected_type,
            analysis_mode=normalized_mode,
        )

    # Cache the analysis result
    analysis_data = {
        "extracted_text": extracted,
        "truncated": truncated,
        "ai_analysis": file_analysis
    }
    await cache_file_analysis(file_hash, detected_type, analysis_data, ttl_seconds=7200)

    return {
        "file_name": file.filename,
        "file_type": detected_type,
        "extracted_text": extracted,
        "truncated": truncated,
        "chars": len(extracted),
        "ai_analysis": file_analysis,
        "analysis_mode": normalized_mode,
        "cache_hit": False
    }

class MemoryManager:
    def __init__(self, db_file: str):
        self.db_file = db_file
        self._initialize_db()
    
    def _initialize_db(self):
        try:
            with sqlite3.connect(self.db_file) as conn:
                conn.execute('''
                    CREATE TABLE IF NOT EXISTS chat_history (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        user_id TEXT NOT NULL,
                        role TEXT NOT NULL,
                        content TEXT NOT NULL,
                        timestamp REAL NOT NULL
                    )
                ''')
                conn.execute('''
                    CREATE TABLE IF NOT EXISTS conversation_summaries (
                        user_id TEXT PRIMARY KEY,
                        summary TEXT NOT NULL,
                        updated_at REAL NOT NULL
                    )
                ''')
                conn.execute("CREATE INDEX IF NOT EXISTS idx_chat_history_user_time ON chat_history(user_id, timestamp)")
                conn.commit()
        except Exception as e:
            logger.error(f"Database initialization error: {e}")
    
    def add_message(self, user_id: str, role: str, content: str):
        try:
            with sqlite3.connect(self.db_file) as conn:
                conn.execute(
                    "INSERT INTO chat_history (user_id, role, content, timestamp) VALUES (?, ?, ?, ?)",
                    (user_id, role, content, time.time())
                )
                conn.commit()
        except Exception as e:
            logger.error(f"Error adding message to memory: {e}")

    def get_recent_messages(self, user_id: str, limit: int = MAX_SESSION_HISTORY_MESSAGES) -> List[Dict[str, str]]:
        try:
            with sqlite3.connect(self.db_file) as conn:
                rows = conn.execute(
                    """
                    SELECT role, content
                    FROM chat_history
                    WHERE user_id = ?
                    ORDER BY timestamp DESC
                    LIMIT ?
                    """,
                    (user_id, max(1, int(limit))),
                ).fetchall()
            rows = list(reversed(rows))
            return _sanitize_history_messages(
                [{"role": row[0], "content": row[1]} for row in rows],
                limit=limit,
            )
        except Exception as e:
            logger.error(f"Error reading chat history: {e}")
            return []

    def get_summary(self, user_id: str) -> str:
        try:
            with sqlite3.connect(self.db_file) as conn:
                row = conn.execute(
                    "SELECT summary FROM conversation_summaries WHERE user_id = ?",
                    (user_id,),
                ).fetchone()
            return str(row[0] or "") if row else ""
        except Exception as e:
            logger.error(f"Error reading conversation summary: {e}")
            return ""

    def save_summary(self, user_id: str, summary: str):
        try:
            with sqlite3.connect(self.db_file) as conn:
                conn.execute(
                    """
                    INSERT INTO conversation_summaries (user_id, summary, updated_at)
                    VALUES (?, ?, ?)
                    ON CONFLICT(user_id) DO UPDATE SET
                        summary = excluded.summary,
                        updated_at = excluded.updated_at
                    """,
                    (user_id, summary or "", time.time()),
                )
                conn.commit()
        except Exception as e:
            logger.error(f"Error saving conversation summary: {e}")

    def build_context_bundle(
        self,
        user_id: str,
        incoming_history: Optional[List[Dict[str, Any]]] = None,
        incoming_summary: str = "",
    ) -> Dict[str, Any]:
        stored_recent = self.get_recent_messages(user_id, limit=MAX_SESSION_HISTORY_MESSAGES)
        merged_messages = _merge_history_messages(stored_recent, incoming_history or [])
        if len(merged_messages) > MAX_HISTORY_MESSAGES_TO_MODEL:
            older_messages = merged_messages[:-MAX_HISTORY_MESSAGES_TO_MODEL]
            recent_messages = merged_messages[-MAX_HISTORY_MESSAGES_TO_MODEL:]
        else:
            older_messages = []
            recent_messages = merged_messages

        summary = _merge_context_summaries(
            self.get_summary(user_id),
            incoming_summary,
            _build_conversation_summary_from_messages(older_messages),
        )
        return {
            "recent_messages": recent_messages,
            "history_summary": summary,
            "older_messages": older_messages,
        }

    def refresh_summary(
        self,
        user_id: str,
        incoming_summary: str = "",
    ) -> str:
        recent_messages = self.get_recent_messages(user_id, limit=MAX_SESSION_HISTORY_MESSAGES)
        older_messages = recent_messages[:-MAX_HISTORY_MESSAGES_TO_MODEL] if len(recent_messages) > MAX_HISTORY_MESSAGES_TO_MODEL else []
        summary = _merge_context_summaries(
            self.get_summary(user_id),
            incoming_summary,
            _build_conversation_summary_from_messages(older_messages),
        )
        if summary:
            self.save_summary(user_id, summary)
        return summary

memory_manager = MemoryManager(DB_FILE)

def initialize_feedback_tables():
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS click_logs (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id TEXT NOT NULL,
                    query TEXT NOT NULL,
                    doc_url TEXT NOT NULL,
                    source TEXT NOT NULL,
                    dwell_ms INTEGER NOT NULL DEFAULT 0,
                    clicked INTEGER NOT NULL DEFAULT 1,
                    timestamp REAL NOT NULL
                )
                """
            )
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS ranking_feedback (
                    doc_url TEXT PRIMARY KEY,
                    views INTEGER NOT NULL DEFAULT 0,
                    clicks INTEGER NOT NULL DEFAULT 0,
                    avg_dwell_ms REAL NOT NULL DEFAULT 0,
                    ctr_score REAL NOT NULL DEFAULT 0,
                    updated_at REAL NOT NULL
                )
                """
            )
            conn.commit()
    except Exception as e:
        logger.error(f"Feedback table initialization error: {e}")

def initialize_ops_tables():
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS query_metrics (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    query TEXT NOT NULL,
                    used_search INTEGER NOT NULL,
                    has_search_data INTEGER NOT NULL,
                    latency_ms REAL NOT NULL,
                    token_count INTEGER NOT NULL DEFAULT 0,
                    language TEXT,
                    timestamp REAL NOT NULL
                )
                """
            )
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS eval_runs (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    run_name TEXT NOT NULL,
                    sample_size INTEGER NOT NULL,
                    recall_at_5 REAL NOT NULL,
                    recall_at_10 REAL NOT NULL,
                    mrr_at_10 REAL NOT NULL,
                    ndcg_at_10 REAL NOT NULL,
                    created_at REAL NOT NULL
                )
                """
            )
            conn.commit()
    except Exception as e:
        logger.error(f"Ops table initialization error: {e}")

def initialize_learning_tables():
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS ranking_impressions (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id TEXT NOT NULL,
                    query TEXT NOT NULL,
                    query_hash TEXT NOT NULL,
                    ab_variant TEXT NOT NULL,
                    rank_position INTEGER NOT NULL,
                    doc_url TEXT NOT NULL,
                    source TEXT,
                    score REAL NOT NULL,
                    timestamp REAL NOT NULL
                )
                """
            )
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS hard_negatives (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id TEXT NOT NULL,
                    query TEXT NOT NULL,
                    query_hash TEXT NOT NULL,
                    positive_url TEXT NOT NULL,
                    negative_url TEXT NOT NULL,
                    variant TEXT NOT NULL,
                    margin REAL NOT NULL DEFAULT 0,
                    created_at REAL NOT NULL
                )
                """
            )
            conn.execute("CREATE INDEX IF NOT EXISTS idx_impr_qh_time ON ranking_impressions(query_hash, timestamp)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_click_qh_time ON click_logs(query, timestamp)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_hn_qh_time ON hard_negatives(query_hash, created_at)")
            conn.commit()
    except Exception as e:
        logger.error(f"Learning table initialization error: {e}")

def _query_hash(query: str) -> str:
    return hashlib.md5((query or "").strip().lower().encode("utf-8")).hexdigest()

def _auto_assign_ab_variant(user_id: str, query: str) -> str:
    # Fully automatic variant policy:
    # 1) Adaptive exploitation by recent CTR/dwell performance
    # 2) Deterministic exploration bucket (no manual switches)
    perf = _get_recent_variant_performance(window_hours=24)
    score_a = float(perf["A"]["score"])
    score_b = float(perf["B"]["score"])
    best_variant = "A" if score_a >= score_b else "B"

    seed = f"{(user_id or 'default_user').strip().lower()}|{(query or '').strip().lower()}|{int(time.time()//3600)}"
    bucket = int(hashlib.md5(seed.encode("utf-8")).hexdigest()[:8], 16) % 100
    exploration_rate = 12  # 12% exploration, 88% exploitation
    if bucket < exploration_rate:
        return "B" if best_variant == "A" else "A"
    return best_variant

def _infer_variant_for_click(user_id: str, query: str, doc_url: str, source: str) -> str:
    source_upper = str(source or "").upper()
    if "AB:A" in source_upper:
        return "A"
    if "AB:B" in source_upper:
        return "B"
    qh = _query_hash(query)
    try:
        with sqlite3.connect(DB_FILE) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                SELECT ab_variant
                FROM ranking_impressions
                WHERE user_id = ? AND query_hash = ? AND doc_url = ?
                ORDER BY timestamp DESC
                LIMIT 1
                """,
                (user_id or "default_user", qh, doc_url),
            )
            row = cur.fetchone()
            if row and str(row[0]).upper() in {"A", "B"}:
                return str(row[0]).upper()
    except Exception:
        pass
    seed = f"{(user_id or 'default_user').strip().lower()}|{(query or '').strip().lower()}"
    bucket = int(hashlib.md5(seed.encode("utf-8")).hexdigest()[:8], 16) % 100
    return "B" if bucket >= 50 else "A"

_AB_POLICY_CACHE: Dict[str, Any] = {"ts": 0.0, "window_hours": 24, "result": None}

def _get_recent_variant_performance(window_hours: int = 24) -> Dict[str, Dict[str, float]]:
    now = time.time()
    cache_valid = (
        _AB_POLICY_CACHE.get("result") is not None
        and _AB_POLICY_CACHE.get("window_hours") == window_hours
        and (now - float(_AB_POLICY_CACHE.get("ts", 0.0))) < 30.0
    )
    if cache_valid:
        return _AB_POLICY_CACHE["result"]

    since_ts = now - (window_hours * 3600)
    impressions = {"A": 0, "B": 0}
    clicks = {"A": 0, "B": 0}
    dwell_sum = {"A": 0.0, "B": 0.0}

    with sqlite3.connect(DB_FILE) as conn:
        cur = conn.cursor()
        cur.execute(
            """
            SELECT ab_variant, COUNT(*)
            FROM ranking_impressions
            WHERE timestamp >= ?
            GROUP BY ab_variant
            """,
            (since_ts,),
        )
        for variant, count in cur.fetchall():
            v = str(variant or "").upper()
            if v in impressions:
                impressions[v] += int(count or 0)

        cur.execute(
            """
            SELECT user_id, query, doc_url, source, dwell_ms
            FROM click_logs
            WHERE timestamp >= ? AND clicked = 1
            """,
            (since_ts,),
        )
        click_rows = cur.fetchall()
        for user_id, query, doc_url, source, dwell_ms in click_rows:
            variant = _infer_variant_for_click(str(user_id), str(query), str(doc_url), str(source))
            clicks[variant] += 1
            dwell_sum[variant] += float(dwell_ms or 0.0)

    result: Dict[str, Dict[str, float]] = {}
    for variant in ("A", "B"):
        impr = float(impressions[variant])
        clk = float(clicks[variant])
        ctr = (clk / impr) if impr > 0 else 0.0
        avg_dwell = (dwell_sum[variant] / clk) if clk > 0 else 0.0
        dwell_score = min(1.0, avg_dwell / 30000.0)
        score = (0.7 * ctr) + (0.3 * dwell_score)
        result[variant] = {
            "impressions": impr,
            "clicks": clk,
            "ctr": ctr,
            "avg_dwell_ms": avg_dwell,
            "score": score,
        }

    _AB_POLICY_CACHE["ts"] = now
    _AB_POLICY_CACHE["window_hours"] = window_hours
    _AB_POLICY_CACHE["result"] = result
    return result

def _get_rank_weights(variant: str, source: str) -> Dict[str, float]:
    v = (variant or "A").upper()
    if source == "database_mcp":
        if v == "B":
            return {"lex": 0.30, "emb": 0.45, "click": 0.25}
        return {"lex": 0.45, "emb": 0.45, "click": 0.10}
    if v == "B":
        return {"lex": 0.25, "emb": 0.45, "click": 0.30}
    return {"lex": 0.40, "emb": 0.45, "click": 0.15}

def record_ranking_impressions(user_id: str, query: str, variant: str, ranked_results: List[Dict[str, Any]], top_k: int = 10):
    if not ranked_results:
        return
    now = time.time()
    qh = _query_hash(query)
    uid = user_id or "default_user"
    rows = []
    for idx, item in enumerate(ranked_results[:top_k], start=1):
        doc_url = str(item.get("url") or "").strip()
        if not doc_url:
            continue
        rows.append(
            (
                uid,
                query[:2000],
                qh,
                variant,
                idx,
                doc_url[:2000],
                str(item.get("source", ""))[:64],
                float(item.get("score", 0.0)),
                now,
            )
        )
    if not rows:
        return
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.executemany(
                """
                INSERT INTO ranking_impressions (
                    user_id, query, query_hash, ab_variant, rank_position, doc_url, source, score, timestamp
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                rows,
            )
            conn.commit()
            _AB_POLICY_CACHE["ts"] = 0.0
    except Exception as e:
        logger.warning(f"record_ranking_impressions failed: {e}")

def get_ab_ranking_stats(window_hours: int = 24) -> Dict[str, Any]:
    perf = _get_recent_variant_performance(window_hours=window_hours)
    score_a = float(perf["A"]["score"])
    score_b = float(perf["B"]["score"])
    winner = "A" if score_a >= score_b else "B"
    return {
        "window_hours": window_hours,
        "impressions": {"A": int(perf["A"]["impressions"]), "B": int(perf["B"]["impressions"])},
        "clicks": {"A": int(perf["A"]["clicks"]), "B": int(perf["B"]["clicks"])},
        "ctr": {
            "A": float(perf["A"]["ctr"]),
            "B": float(perf["B"]["ctr"]),
        },
        "avg_dwell_ms": {"A": float(perf["A"]["avg_dwell_ms"]), "B": float(perf["B"]["avg_dwell_ms"])},
        "score": {"A": score_a, "B": score_b},
        "policy": {"winner": winner, "exploration_rate": 0.12},
    }

def mine_hard_negatives(limit_queries: int = 200, negatives_per_query: int = 3) -> Dict[str, Any]:
    now = time.time()
    inserted = 0
    processed = 0
    with sqlite3.connect(DB_FILE) as conn:
        cur = conn.cursor()
        cur.execute(
            """
            SELECT user_id, query, doc_url, MAX(timestamp) as ts
            FROM click_logs
            WHERE clicked = 1
            GROUP BY user_id, query, doc_url
            ORDER BY ts DESC
            LIMIT ?
            """,
            (max(20, min(limit_queries, 2000)),),
        )
        positives = cur.fetchall()

        for user_id, query, pos_url, _ in positives:
            processed += 1
            qh = _query_hash(query)
            cur.execute(
                """
                SELECT doc_url
                FROM click_logs
                WHERE user_id = ? AND query = ? AND clicked = 1
                """,
                (user_id, query),
            )
            clicked_urls = {r[0] for r in cur.fetchall()}

            cur.execute(
                """
                SELECT doc_url, score, ab_variant
                FROM ranking_impressions
                WHERE user_id = ? AND query_hash = ?
                ORDER BY score DESC, timestamp DESC
                LIMIT 30
                """,
                (user_id, qh),
            )
            neg_candidates = cur.fetchall()

            taken = 0
            for neg_url, neg_score, variant in neg_candidates:
                if neg_url == pos_url or neg_url in clicked_urls:
                    continue
                cur.execute(
                    """
                    INSERT INTO hard_negatives (
                        user_id, query, query_hash, positive_url, negative_url, variant, margin, created_at
                    )
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        user_id,
                        query[:2000],
                        qh,
                        str(pos_url)[:2000],
                        str(neg_url)[:2000],
                        str(variant or "A")[:4],
                        float(neg_score or 0.0),
                        now,
                    ),
                )
                inserted += 1
                taken += 1
                if taken >= max(1, min(negatives_per_query, 10)):
                    break
        conn.commit()
    return {"processed_queries": processed, "inserted_hard_negatives": inserted}

def record_query_metric(
    query: str,
    used_search: bool,
    has_search_data: bool,
    latency_ms: float,
    token_count: int,
    language: str,
):
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.execute(
                """
                INSERT INTO query_metrics (query, used_search, has_search_data, latency_ms, token_count, language, timestamp)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    query[:2000],
                    1 if used_search else 0,
                    1 if has_search_data else 0,
                    float(latency_ms),
                    int(token_count),
                    (language or "")[:10],
                    time.time(),
                ),
            )
            conn.commit()
    except Exception as e:
        logger.warning(f"record_query_metric failed: {e}")

def get_click_boost(doc_url: str) -> float:
    try:
        with sqlite3.connect(DB_FILE) as conn:
            cur = conn.cursor()
            cur.execute("SELECT ctr_score FROM ranking_feedback WHERE doc_url = ?", (doc_url,))
            row = cur.fetchone()
            if not row:
                return 0.0
            return max(-0.2, min(0.4, float(row[0])))
    except Exception:
        return 0.0

def record_click_feedback(feedback: ClickFeedback) -> Dict[str, Any]:
    now = time.time()
    dwell_ms = max(0, int(feedback.dwell_ms))
    clicked_val = 1 if feedback.clicked else 0
    inferred_variant = _infer_variant_for_click(feedback.user_id, feedback.query, feedback.doc_url, feedback.source)
    enriched_source = str(feedback.source or "hybrid")
    if f"ab:{inferred_variant}".upper() not in enriched_source.upper():
        enriched_source = f"{enriched_source}|ab:{inferred_variant}"
    try:
        with sqlite3.connect(DB_FILE) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                INSERT INTO click_logs (user_id, query, doc_url, source, dwell_ms, clicked, timestamp)
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (feedback.user_id, feedback.query, feedback.doc_url, enriched_source, dwell_ms, clicked_val, now),
            )

            cur.execute("SELECT views, clicks, avg_dwell_ms FROM ranking_feedback WHERE doc_url = ?", (feedback.doc_url,))
            row = cur.fetchone()
            if row:
                views, clicks, avg_dwell = int(row[0]), int(row[1]), float(row[2])
                new_views = views + 1
                new_clicks = clicks + clicked_val
                new_avg_dwell = ((avg_dwell * views) + dwell_ms) / new_views
            else:
                new_views = 1
                new_clicks = clicked_val
                new_avg_dwell = float(dwell_ms)

            ctr = (new_clicks / new_views) if new_views > 0 else 0.0
            # Online learning signal: CTR + dwell quality
            dwell_score = min(1.0, new_avg_dwell / 30000.0)
            ctr_score = (0.7 * ctr) + (0.3 * dwell_score)

            cur.execute(
                """
                INSERT INTO ranking_feedback (doc_url, views, clicks, avg_dwell_ms, ctr_score, updated_at)
                VALUES (?, ?, ?, ?, ?, ?)
                ON CONFLICT(doc_url) DO UPDATE SET
                    views = excluded.views,
                    clicks = excluded.clicks,
                    avg_dwell_ms = excluded.avg_dwell_ms,
                    ctr_score = excluded.ctr_score,
                    updated_at = excluded.updated_at
                """,
                (feedback.doc_url, new_views, new_clicks, new_avg_dwell, ctr_score, now),
            )
            conn.commit()
            _AB_POLICY_CACHE["ts"] = 0.0

            return {
                "doc_url": feedback.doc_url,
                "views": new_views,
                "clicks": new_clicks,
                "avg_dwell_ms": round(new_avg_dwell, 2),
                "ctr_score": round(ctr_score, 4),
                "ab_variant": inferred_variant,
            }
    except Exception as e:
        logger.error(f"record_click_feedback error: {e}")
        raise

def _percentile(values: List[float], p: float) -> float:
    if not values:
        return 0.0
    vals = sorted(values)
    k = (len(vals) - 1) * p
    f = math.floor(k)
    c = math.ceil(k)
    if f == c:
        return float(vals[int(k)])
    d0 = vals[f] * (c - k)
    d1 = vals[c] * (k - f)
    return float(d0 + d1)

def get_observability_summary(window_hours: int = 24) -> Dict[str, Any]:
    since_ts = time.time() - (window_hours * 3600)
    with sqlite3.connect(DB_FILE) as conn:
        cur = conn.cursor()
        cur.execute(
            """
            SELECT latency_ms, used_search, has_search_data, token_count
            FROM query_metrics
            WHERE timestamp >= ?
            """,
            (since_ts,),
        )
        rows = cur.fetchall()
        latencies = [float(r[0]) for r in rows]
        used_search = sum(int(r[1]) for r in rows)
        has_data = sum(int(r[2]) for r in rows)
        token_counts = [int(r[3]) for r in rows]

        cur.execute(
            """
            SELECT COUNT(*), SUM(clicked), AVG(dwell_ms)
            FROM click_logs
            WHERE timestamp >= ?
            """,
            (since_ts,),
        )
        fb = cur.fetchone() or (0, 0, 0)
        cur.execute(
            """
            SELECT COUNT(*)
            FROM hard_negatives
            WHERE created_at >= ?
            """,
            (since_ts,),
        )
        hn_count = int((cur.fetchone() or (0,))[0] or 0)

    total = len(rows)
    ab_stats = get_ab_ranking_stats(window_hours=window_hours)
    return {
        "window_hours": window_hours,
        "query_count": total,
        "search_usage_rate": (used_search / total) if total else 0.0,
        "search_has_data_rate": (has_data / total) if total else 0.0,
        "latency_ms": {
            "p50": _percentile(latencies, 0.50),
            "p95": _percentile(latencies, 0.95),
            "p99": _percentile(latencies, 0.99),
            "avg": statistics.mean(latencies) if latencies else 0.0,
        },
        "tokens": {
            "avg": statistics.mean(token_counts) if token_counts else 0.0,
            "max": max(token_counts) if token_counts else 0,
        },
        "feedback": {
            "events": int(fb[0] or 0),
            "clicks": int(fb[1] or 0),
            "avg_dwell_ms": float(fb[2] or 0.0),
        },
        "ab_ranking": ab_stats,
        "hard_negatives": {
            "count": hn_count,
        },
    }

def retrain_ranking_from_click_logs() -> Dict[str, Any]:
    with sqlite3.connect(DB_FILE) as conn:
        cur = conn.cursor()
        cur.execute(
            """
            SELECT doc_url, COUNT(*) AS views, SUM(clicked) AS clicks, AVG(dwell_ms) AS avg_dwell
            FROM click_logs
            GROUP BY doc_url
            """
        )
        rows = cur.fetchall()
        updated = 0
        now = time.time()
        for doc_url, views, clicks, avg_dwell in rows:
            views = int(views or 0)
            clicks = int(clicks or 0)
            avg_dwell = float(avg_dwell or 0.0)
            if views <= 0:
                continue
            ctr = clicks / views
            dwell_score = min(1.0, avg_dwell / 30000.0)
            ctr_score = (0.7 * ctr) + (0.3 * dwell_score)
            cur.execute(
                """
                INSERT INTO ranking_feedback (doc_url, views, clicks, avg_dwell_ms, ctr_score, updated_at)
                VALUES (?, ?, ?, ?, ?, ?)
                ON CONFLICT(doc_url) DO UPDATE SET
                    views = excluded.views,
                    clicks = excluded.clicks,
                    avg_dwell_ms = excluded.avg_dwell_ms,
                    ctr_score = excluded.ctr_score,
                    updated_at = excluded.updated_at
                """,
                (doc_url, views, clicks, avg_dwell, ctr_score, now),
            )
            updated += 1
        conn.commit()
    return {"updated_docs": updated}

def _ndcg_at_k(rels: List[int], k: int) -> float:
    rels_k = rels[:k]
    if not rels_k:
        return 0.0
    dcg = 0.0
    for i, rel in enumerate(rels_k, start=1):
        dcg += (2**rel - 1) / math.log2(i + 1)
    ideal = sorted(rels, reverse=True)[:k]
    idcg = 0.0
    for i, rel in enumerate(ideal, start=1):
        idcg += (2**rel - 1) / math.log2(i + 1)
    return (dcg / idcg) if idcg > 0 else 0.0

def run_offline_eval_from_clicks(sample_size: int = 100) -> Dict[str, Any]:
    if not db_mcp_backend:
        return {"error": "database_mcp_backend unavailable"}
    with sqlite3.connect(DB_FILE) as conn:
        cur = conn.cursor()
        cur.execute(
            """
            SELECT query, doc_url
            FROM click_logs
            WHERE clicked = 1
            ORDER BY timestamp DESC
            LIMIT ?
            """,
            (max(10, min(sample_size, 500)),),
        )
        pairs = cur.fetchall()

    if not pairs:
        return {"error": "no click logs for evaluation"}

    recall5 = []
    recall10 = []
    mrr10 = []
    ndcg10 = []
    eval_count = 0
    for query_text, rel_url in pairs:
        try:
            res = db_mcp_backend.search_articles(query_text, category="all", limit=20)
            urls = [str(r.get("url") or r.get("canonical_url") or "") for r in res.get("results", [])]
            if not urls:
                continue
            eval_count += 1
            hit5 = 1 if rel_url in urls[:5] else 0
            hit10 = 1 if rel_url in urls[:10] else 0
            recall5.append(hit5)
            recall10.append(hit10)
            rr = 0.0
            rels = []
            for rank, u in enumerate(urls[:10], start=1):
                rel = 1 if u == rel_url else 0
                rels.append(rel)
                if rel == 1 and rr == 0.0:
                    rr = 1.0 / rank
            mrr10.append(rr)
            ndcg10.append(_ndcg_at_k(rels, 10))
        except Exception:
            continue

    result = {
        "sample_size": eval_count,
        "recall_at_5": statistics.mean(recall5) if recall5 else 0.0,
        "recall_at_10": statistics.mean(recall10) if recall10 else 0.0,
        "mrr_at_10": statistics.mean(mrr10) if mrr10 else 0.0,
        "ndcg_at_10": statistics.mean(ndcg10) if ndcg10 else 0.0,
    }
    with sqlite3.connect(DB_FILE) as conn:
        conn.execute(
            """
            INSERT INTO eval_runs (run_name, sample_size, recall_at_5, recall_at_10, mrr_at_10, ndcg_at_10, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                "offline_click_eval",
                int(result["sample_size"]),
                float(result["recall_at_5"]),
                float(result["recall_at_10"]),
                float(result["mrr_at_10"]),
                float(result["ndcg_at_10"]),
                time.time(),
            ),
        )
        conn.commit()
    return result

async def auto_learning_loop():
    if os.getenv("SKEMI_DEBUG_LEARNING", "0").strip().lower() in {"1", "true", "yes", "on"}:
        logger.info(
            f" Auto-learning loop started (interval={AUTO_LEARNING_INTERVAL_SECONDS}s, enabled={AUTO_LEARNING_ENABLED})"
        )
    while True:
        try:
            if AUTO_LEARNING_ENABLED:
                mined = mine_hard_negatives(limit_queries=200, negatives_per_query=2)
                retrain = retrain_ranking_from_click_logs()
                if os.getenv("SKEMI_DEBUG_LEARNING", "0").strip().lower() in {"1", "true", "yes", "on"}:
                    logger.info(
                        " Auto-learning cycle done: "
                        f"hard_negatives={mined.get('inserted_hard_negatives', 0)}, "
                        f"updated_docs={retrain.get('updated_docs', 0)}"
                    )
            await asyncio.sleep(max(30, AUTO_LEARNING_INTERVAL_SECONDS))
        except asyncio.CancelledError:
            logger.info(" Auto-learning loop cancelled")
            break
        except Exception as e:
            logger.warning(f"Auto-learning loop error: {e}")
            await asyncio.sleep(30)

initialize_feedback_tables()
initialize_ops_tables()
initialize_learning_tables()

def extract_ai_analysis_from_question(question: str) -> tuple[str, str]:
    if not question:
        return question, ""
    markers = [
        "Phân tích ảnh (AI):",
        "Phân tích AI:",
        "Phan tich anh (AI):",
        "Phan tich AI:",
        "(AI):",
        "IMAGE ANALYSIS:",
    ]
    lower = question.lower()
    start_idx = None
    marker_used = ""
    for marker in markers:
        idx = lower.find(marker.lower())
        if idx != -1 and (start_idx is None or idx < start_idx):
            start_idx = idx
            marker_used = marker
    if start_idx is None:
        return question, ""

    start_content = start_idx + len(marker_used)
    after = question[start_content:]
    end_markers = [
        "Nội dung trích xuất:",
        "URL:",
        "Tệp:",
        "Loại:",
        "Ghi chú:",
        "Văn bản OCR:",
    ]
    lower_after = after.lower()
    end_idx_rel = None
    for end_marker in end_markers:
        idx = lower_after.find(end_marker.lower())
        if idx != -1 and (end_idx_rel is None or idx < end_idx_rel):
            end_idx_rel = idx

    if end_idx_rel is None:
        analysis = after.strip()
        cleaned = question[:start_idx].rstrip().strip()
        return cleaned, analysis

    analysis = after[:end_idx_rel].strip()
    cleaned = (question[:start_idx].rstrip() + "\n" + after[end_idx_rel:].lstrip()).strip()
    return cleaned, analysis

def build_smart_prompt(language: str, has_search_data: bool, question: str, current_time: str, image_analysis: str = "", deep_research: bool = False) -> str:
    """Universal prompt builder for Skemi - Absolute Recency & Strict Grounding."""
    normalized_language = str(language or "").strip().lower()
    is_vietnamese = normalized_language == "vi"
    is_english = normalized_language in {"", "en"}
    if is_vietnamese:
        base_identity = """Bạn là Skemi - Trợ lý AI thế hệ mới, hoạt động theo nguyên tắc: **Sự thật của hiện tại là ưu tiên duy nhất**.
    
    NGUYÊN TẮC VÀNG:
    1. **Quyền năng của Tin mới nhất**: Tự động đánh giá các mốc thời gian trong ngữ cảnh, luôn ưu tiên tổng hợp thông tin mới nhất.
    2. **Trả lời Đầy đủ & Chi tiết**: Trình bày thông tin phong phú, đầy đủ ngữ cảnh và chi tiết như một trợ lý chuyên gia. Không được trả lời quá ngắn gọn hay cụt lủn.
    3. **Tuyệt đối Không lộ Backend & Không nhắc đến việc tìm kiếm**: Tổng hợp thông tin từ ngữ cảnh một cách tự nhiên như thể bạn đã luôn biết điều đó. KHÔNG BAO GIỜ dùng các cụm từ như "Theo kết quả tìm kiếm", "Dựa trên nguồn này", "Hệ thống vừa search được". KHÔNG để lộ việc bạn có sử dụng công cụ tìm kiếm.
    4. **Chỉ dựa trên sự thật (Chống bịa đặt)**: Dùng thông tin từ ngữ cảnh làm cốt lõi. Gắn kết chúng lại mạch lạc, KHÔNG bịa đặt thêm số liệu hay chi tiết không có thật. Nếu thiếu thông tin, hãy dùng kiến thức nền tảng để giải thích tổng quan nhưng không tự chế ra phiên bản/ngày tháng sự kiện cụ thể.
    5. **Trình bày chuyên nghiệp (Giống ChatGPT)**: 
       - Trình bày mạch lạc, chia đoạn văn bản hợp lý, dễ đọc.
       - Dùng `**in đậm**` cho các từ khóa quan trọng.
       - Sử dụng bullet points (danh sách) khi cần liệt kê.
       - Dùng Heading (`###` hoặc `####`) để chia các phần chính.
    6. **Xác minh đa chủ đề**: Nếu hỏi nhiều chủ đề, chia tay từng phần rõ ràng.
    """
        if has_search_data:
            search_instruction = "Tổng hợp chi tiết và tự nhiên nhất các dữ liệu sau đây, trình bày như một chuyên gia AI vô cùng hiểu biết (không nhắc đến việc bạn có dữ liệu từ đâu):"
        else:
            search_instruction = "Trả lời chi tiết dựa trên kiến thức xác thực của bạn."
    else:
        base_identity = """You are Skemi - A next-gen AI assistant. Your core directive is: **Current Fact is the Only Priority**.
    
    GOLDEN RULES:
    1. **Absolute Recency Authority**: Always evaluate dates and versions in the context, synthesizing and giving utmost priority to the newest information.
    2. **Comprehensive & Detailed**: Present information completely, with full context and richness, like an expert assistant. Never be overly brief or terse.
    3. **Invisible Backend & No Mentions of Search**: Synthesize information naturally as if you inherently know it. NEVER use phrases like "Based on my search results", "According to the sources", or "I found this". DO NOT expose that a search tool was used.
    4. **Fact-based Only (Zero Hallucination)**: Use the provided context as the absolute core truth. Blend them coherently but DO NOT invent facts, numbers, dates, or versions. If specifics are missing, explain generally based on your pre-training without fabricating.
    5. **Professional Formatting (ChatGPT Style)**: 
       - Write coherently with well-structured paragraphs.
       - Use `**bold**` for key terms.
       - Use bullet points when listing items.
       - Use Headings (`###` or `####`) to divide major sections.
    6. **Multi-topic Verification**: If asked about multiple topics, clearly divide your response.
    """
        if has_search_data:
            search_instruction = "Synthesize following data naturally and comprehensively as an expert (do not utter where or how you got this data):"
        else:
            search_instruction = "Answer concisely based on your authentic knowledge."

    # Deep research enhancement
    deep_instruction = ""
    if deep_research:
        if language == "vi":
            deep_instruction = "\n[CHẾ ĐỘ NGHIÊN CỨU SÂU]: Phân tích đa chiều, so sánh chi tiết các thông tin tìm được."
        else:
            deep_instruction = "\n[DEEP RESEARCH MODE]: Provide multi-dimensional analysis and detailed comparisons of found info."

    # Image analysis integration
    newline = "\n"
    image_prompt = ""
    if image_analysis:
        header = "KẾT QUẢ PHÂN TÍCH HÌNH ẢNH/FILE" if language == "vi" else "IMAGE/FILE ANALYSIS RESULTS"
        image_prompt = f"{newline}{newline}{'='*50}{newline} {header} {newline}{'='*50}{newline}{newline}{image_analysis}{newline}{newline}[Lưu ý]: Sử dụng thông tin trên để giải quyết yêu cầu về ảnh/file."

    time_context = f"Thoi diem hien tai o Viet Nam: {current_time}. Uu tien thong tin moi nhat theo moc thoi gian nay." if language == "vi" else f"Current Vietnam local time: {current_time}. Treat this as the true current date/time for any latest question."
    full_prompt = f"{base_identity}{newline}{newline}{time_context}{newline}{newline}{search_instruction}{newline}{deep_instruction}{image_prompt}"
    return full_prompt

def build_smart_prompt_v2(language: str, has_search_data: bool, question: str, current_time: str, image_analysis: str = "", deep_research: bool = False) -> str:
    normalized_language = str(language or "").strip().lower()
    is_vietnamese = normalized_language == "vi"
    is_english = normalized_language in {"", "en"}

    if is_vietnamese:
        base_identity = """Bạn là Skemi. Mục tiêu của bạn là trả lời đúng, mới, rõ, tự nhiên và không lộ backend."""
        search_instruction = (
            "Hãy tổng hợp tự nhiên các dữ kiện đã có, ưu tiên thông tin mới nhất."
            if has_search_data else
            "Hãy trả lời dựa trên kiến thức xác thực của bạn."
        )
        deep_instruction = "\n[CHẾ ĐỘ NGHIÊN CỨU SÂU]: Phân tích đa chiều và so sánh kỹ các dữ kiện." if deep_research else ""
        time_context = f"Thời điểm hiện tại ở Việt Nam: {current_time}. Luôn ưu tiên dữ kiện mới nhất theo mốc này."
        image_header = "KẾT QUẢ PHÂN TÍCH HÌNH ẢNH/FILE"
    else:
        base_identity = """You are Skemi. Your goal is to answer accurately, clearly, naturally, and without exposing backend/tool behavior."""
        search_instruction = (
            "Synthesize the available evidence naturally and prioritize the newest explicitly dated facts."
            if has_search_data else
            "Answer from your reliable knowledge."
        )
        if not is_english:
            search_instruction += f" Final response language: answer entirely in '{normalized_language}'."
        deep_instruction = "\n[DEEP RESEARCH MODE]: Provide multi-angle analysis and careful comparisons." if deep_research else ""
        time_context = f"Current Vietnam local time: {current_time}. Treat this as the true current date/time for any latest question."
        image_header = "IMAGE/FILE ANALYSIS RESULTS"

    newline = "\n"
    image_prompt = ""
    if image_analysis:
        image_prompt = (
            f"{newline}{newline}{'=' * 50}{newline}{image_header}{newline}{'=' * 50}"
            f"{newline}{image_analysis}{newline}"
        )

    return f"{base_identity}{newline}{newline}{time_context}{newline}{newline}{search_instruction}{deep_instruction}{image_prompt}"

def estimate_tokens(text: str) -> int:
    """Rough estimate: 1 token ≈ 4 characters for Vietnamese/English"""
    return len(text) // 4

def _truncate_text(text: str, max_chars: int = 12000) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars] + "\n\n[ĐÃ RÚT GỌN]", True

def _has_usable_ocr(text: str) -> bool:
    if not text:
        return False
    t = text.strip()
    if len(t) < OCR_MIN_CHARS:
        return False
    if t.lower().startswith("cannot ocr"):
        return False
    return True

def _extract_pdf_tables_with_pdfplumber(data: bytes) -> List[str]:
    rows_out: List[str] = []
    if pdfplumber is None:
        return rows_out
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page_idx, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables() or []
                for table_idx, table in enumerate(tables, start=1):
                    rows_out.append(f"[pdfplumber][page={page_idx}][table={table_idx}]")
                    for row in table or []:
                        cleaned = [str(cell).strip() if cell is not None else "" for cell in row]
                        if any(cleaned):
                            rows_out.append(" | ".join(cleaned))
    except Exception as e:
        rows_out.append(f"[pdfplumber_error] {e}")
    return rows_out

def _extract_pdf_tables_with_camelot(data: bytes) -> List[str]:
    rows_out: List[str] = []
    if camelot is None:
        return rows_out
    temp_path = ""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(data)
            temp_path = tmp.name
        tables = camelot.read_pdf(temp_path, pages="all", flavor="stream")
        for idx, table in enumerate(tables, start=1):
            rows_out.append(f"[camelot][table={idx}]")
            df = table.df
            for _, row in df.iterrows():
                cells = [str(cell).strip() for cell in list(row.values)]
                if any(cells):
                    rows_out.append(" | ".join(cells))
    except Exception as e:
        rows_out.append(f"[camelot_error] {e}")
    finally:
        if temp_path:
            try:
                os.remove(temp_path)
            except Exception:
                pass
    return rows_out

def _extract_text_from_pdf(data: bytes) -> str:
    parts = []
    if PdfReader is None:
        parts.append("Cannot read PDF text layer: missing PyPDF2.")
    else:
        reader = PdfReader(io.BytesIO(data))
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue

    table_lines: List[str] = []
    table_lines.extend(_extract_pdf_tables_with_pdfplumber(data))
    table_lines.extend(_extract_pdf_tables_with_camelot(data))
    if table_lines:
        parts.append("\n[TABLE/FORM EXTRACTION]\n" + "\n".join(table_lines[:6000]))
    return "\n".join(parts).strip()

def _extract_text_from_docx(data: bytes) -> str:
    if docx is None:
        return "Cannot read DOCX: missing python-docx."
    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs).strip()

def _extract_text_from_pptx(data: bytes) -> str:
    if Presentation is None:
        return "Cannot read PPTX: missing python-pptx."
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts).strip()

def _extract_text_from_xlsx(data: bytes) -> str:
    if openpyxl is None:
        return "Cannot read XLSX: missing openpyxl."
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_vals = [str(cell) if cell is not None else "" for cell in row]
            if any(row_vals):
                parts.append("\t".join(row_vals))
    return "\n".join(parts).strip()

def _extract_text_from_csv(data: bytes) -> str:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode("latin-1", errors="replace")
    reader = csv.reader(io.StringIO(text))
    lines = []
    for row in reader:
        lines.append(",".join(row))
    return "\n".join(lines).strip()

def _extract_text_from_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8", errors="replace").strip()
    except Exception:
        return data.decode("latin-1", errors="replace").strip()

def _extract_text_from_image(data: bytes) -> str:
    # OCR extraction only, not full semantic image analysis.
    if Image is None or pytesseract is None:
        return "Cannot OCR image: missing PIL or pytesseract."
    try:
        img = Image.open(io.BytesIO(data))
    except Exception as e:
        return f"Cannot OCR image: {e}"

    def _clean(text: str) -> str:
        t = (text or "").replace("\r", "\n")
        t = re.sub(r"[ \t]+", " ", t)
        t = re.sub(r"\n{3,}", "\n\n", t)
        return t.strip()

    def _try_ocr(image, lang: str) -> str:
        try:
            return pytesseract.image_to_string(image, lang=lang, config="--oem 3 --psm 6")
        except Exception:
            return ""

    langs = ["vie+eng", "vie", "eng"]
    try:
        available = set(pytesseract.get_languages(config=""))
        filtered = []
        for l in langs:
            if "+" in l:
                parts = l.split("+")
                if all(p in available for p in parts):
                    filtered.append(l)
            else:
                if l in available:
                    filtered.append(l)
        if filtered:
            langs = filtered
    except Exception:
        pass

    variants = [img]
    try:
        gray = img.convert("L")
        variants.append(gray)
        w, h = gray.size
        if max(w, h) < 1600:
            variants.append(gray.resize((w * 2, h * 2), Image.LANCZOS))
        bw = gray.point(lambda x: 0 if x < 180 else 255, '1')
        variants.append(bw)
    except Exception:
        pass

    best_text = ""
    best_score = 0
    for v in variants:
        for lang in langs:
            t = _clean(_try_ocr(v, lang))
            score = len(re.sub(r"\s+", "", t))
            if score > best_score:
                best_score = score
                best_text = t

    return best_text

def extract_text_by_extension(filename: str, data: bytes) -> tuple[str, str]:
    ext = (filename.split(".")[-1] or "").lower()
    if ext in ["pdf"]:
        return _extract_text_from_pdf(data), "pdf"
    if ext in ["docx"]:
        return _extract_text_from_docx(data), "docx"
    if ext in ["pptx"]:
        return _extract_text_from_pptx(data), "pptx"
    if ext in ["xlsx"]:
        return _extract_text_from_xlsx(data), "xlsx"
    if ext in ["csv"]:
        return _extract_text_from_csv(data), "csv"
    if ext in ["txt", "md", "log"]:
        return _extract_text_from_txt(data), "text"
    if ext in ["png", "jpg", "jpeg", "webp", "bmp", "tiff"]:
        # OCR text extraction from image (text only, no semantic reasoning).
        return _extract_text_from_image(data), "image"
    return _extract_text_from_txt(data), "text"

def _normalize_router_text(text: str) -> str:
    raw = unicodedata.normalize("NFKD", str(text or ""))
    raw = "".join(ch for ch in raw if not unicodedata.combining(ch))
    raw = raw.replace("\u0111", "d").replace("\u0110", "d")
    return re.sub(r"\s+", " ", raw.lower()).strip()


def _query_requires_live_search(question: str) -> bool:
    text = _normalize_router_text(question)
    if not text:
        return False
    if "http://" in text or "https://" in text or "www." in text or "site:" in text:
        return True
    years = [int(y) for y in re.findall(r"\b(19\d{2}|20\d{2})\b", text)]
    if years:
        return True
    if re.search(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b", text):
        return True
    return False


def _sanitize_router_search_queries(question: str, queries: List[str]) -> List[str]:
    clean_queries: List[str] = []
    seen: Set[str] = set()

    for raw_query in list(queries or []):
        query = re.sub(r"\s+", " ", str(raw_query or "")).strip()
        if not query:
            continue
        query_key = _normalize_router_text(query)
        if not query_key or query_key in seen:
            continue
        seen.add(query_key)
        clean_queries.append(query)

    if not clean_queries:
        clean_queries = [question]
    return clean_queries[:6]

async def pure_model_router(
    question: str,
    conversation_context: str = "",
    deep_research: bool = False,
) -> Dict[str, Any]:
    """Use AI to decide if search is needed and generate queries with stronger freshness controls."""
    router_models = ["devstral-2:123b-cloud", "minimax-m2:cloud"]
    now = datetime.now()
    current_year = now.year
    current_clock = get_current_datetime()["full"]
    year_matches = [int(y) for y in re.findall(r"\b(19\d{2}|20\d{2})\b", question or "")]
    year_mentioned = max(year_matches) if year_matches else None

    async def _finalize_router_result(result: Dict[str, Any], source_label: str) -> Dict[str, Any]:
        if "needs_web_search" not in result:
            result["needs_web_search"] = False
        detected_language = str(result.get("language", "")).strip().lower()
        if detected_language not in {"vi", "en"}:
            detected_language = await detect_language_with_model(question)
        result["language"] = detected_language

        queries = result.get("search_queries", [question])
        if not isinstance(queries, list):
            queries = [question]
        queries = _sanitize_router_search_queries(question, list(queries))
        if len(queries) < 5:
            time_expansions = [question]
            if str(current_year) not in _normalize_router_text(question):
                time_expansions.append(f"{question} {current_year}")
            if deep_research:
                time_expansions.append(f"{question} {now.strftime('%Y-%m')}")
            queries = _sanitize_router_search_queries(question, queries + time_expansions)
        result["search_queries"] = queries[: (7 if deep_research else 5)]

        try:
            result["answer_confidence"] = max(0.0, min(1.0, float(result.get("answer_confidence", 0.5))))
        except Exception:
            result["answer_confidence"] = 0.5

        depth = str(result.get("research_depth", "auto")).strip().lower()
        if depth not in {"auto", "standard", "deep", "exhaustive"}:
            depth = "auto"
        result["research_depth"] = depth

        if "recency_priority" not in result:
            result["recency_priority"] = "medium"

        plan = result.get("research_plan", [])
        is_bad_plan = (
            not isinstance(plan, list)
            or len(plan) == 0
            or all(str(p or "").strip().lower() in {"", question.strip().lower()} for p in plan)
        )
        if is_bad_plan:
            result["research_plan"] = [
                f"Collect information for: {queries[0]}",
                "Verify the strongest sources and dates",
                "Synthesize only the supported facts",
            ]
        else:
            result["research_plan"] = [str(item) for item in plan[:5]]

        result["needs_search"] = bool(result.get("needs_web_search", False))
        is_time_sensitive = _query_requires_live_search(question)
        if result["needs_search"] and result["answer_confidence"] >= 0.95 and not is_time_sensitive:
            logger.info(f"Router Decision ({source_label}): high confidence and not time-sensitive, skipping search")
            result["needs_search"] = False
            result["search_reason"] = "high_confidence_internalknowledge"

        effective_year = year_mentioned
        if effective_year is not None:
            if current_year - effective_year <= 2:
                result["recency_priority"] = "high"
                result["topic_type"] = "recent"
            elif effective_year < current_year - 5:
                result["recency_priority"] = "low"
                result["topic_type"] = "historical"
            else:
                result["recency_priority"] = result.get("recency_priority", "medium")
                result["topic_type"] = "general"
        else:
            result.setdefault("topic_type", "general")

        result["is_safe"] = True
        result["year_mentioned"] = effective_year
        logger.info(
            f"AI Router ({source_label}): SUCCESS - search={result['needs_search']}, lang={result['language']}, "
            f"priority={result.get('recency_priority', 'medium')}, depth={result.get('research_depth', 'auto')}, queries={result.get('search_queries', [])}"
        )
        return result

    for model_name in router_models:
        try:
            logger.info(f"Trying router model: {model_name}")
            prompt = f'''Analyze the user question and plan a search strategy. Respond in JSON only.
Question: "{question}"
Current Vietnam local time: {current_clock}
Current year: {current_year}
Conversation context:
{conversation_context or "(none)"}

Rules:
1. The current date above is authoritative.
2. If the user asks for time-sensitive, version-sensitive, or rapidly changing information, search only for current information relative to the current date.
3. Prefer official vendor domains, product docs, release notes, and major reputable press.
4. Never generate queries centered on old versions, old years, rumors, or fan sites unless the user explicitly asks for history/timeline.
5. If the question mentions multiple subjects, create separate queries for each subject.
6. Use the conversation context only to disambiguate subject/topic, not to invent facts.
7. Keep queries concise and high-signal.
8. Do not pad queries with boilerplate suffixes like "official", "update", "documentation", or "overview" unless they are necessary to disambiguate the subject.

JSON format:
{{
  "needs_web_search": true,
  "language": "vi",
  "search_queries": ["query 1", "query 2"],
  "search_reason": "brief reason",
  "recency_priority": "high",
  "research_depth": "standard",
  "answer_confidence": 0.2,
  "research_plan": ["step 1", "step 2", "step 3"]
}}'''
            router_messages = [{"role": "user", "content": prompt}]
            router_ctx = await fit_context_window_for_model(
                model_name,
                desired_ctx=12000,
                messages=router_messages,
                max_output_tokens=260,
            )

            async with httpx.AsyncClient(timeout=30) as client:
                response = await client.post(
                    OLLAMA_URL,
                    json={
                        "model": model_name,
                        "messages": router_messages,
                        "stream": False,
                        "options": {"temperature": 0.1, "num_predict": 260, "num_ctx": router_ctx},
                        "format": None if model_name == "devstral-2:123b-cloud" else "json",
                    }
                )

            logger.info(f"Router {model_name} response status: {response.status_code}")
            if response.status_code != 200:
                logger.warning(f"Router model {model_name} returned status {response.status_code}")
                continue

            content = response.json().get("message", {}).get("content", "").strip()
            logger.info(f"Router {model_name} raw response: {content[:200]}...")

            parsed_result = None
            try:
                json_blocks = re.findall(r"\{.*\}", content, re.DOTALL)
                candidate = max(json_blocks, key=len).strip() if json_blocks else content
                parsed_result = json.loads(candidate)
                logger.info(f"Router {model_name} parsed JSON successfully")
            except Exception as json_error:
                logger.warning(f"Router {model_name} JSON decode error: {json_error}, trying regex fallback")
                try:
                    parsed_result = {}
                    m_needs = re.search(r'"needs_web_search"\s*:\s*(true|false)', content, re.IGNORECASE)
                    parsed_result["needs_web_search"] = (m_needs.group(1).lower() == "true") if m_needs else _query_requires_live_search(question)
                    m_lang = re.search(r'"language"\s*:\s*"([^"]+)"', content, re.IGNORECASE)
                    parsed_result["language"] = m_lang.group(1) if m_lang else "vi"
                    m_queries = re.search(r'"search_queries"\s*:\s*\[(.*?)\]', content, re.DOTALL | re.IGNORECASE)
                    parsed_result["search_queries"] = [q.strip() for q in re.findall(r'"([^"]+)"', m_queries.group(1))] if m_queries else [question]
                    m_reason = re.search(r'"search_reason"\s*:\s*"([^"]*)"', content, re.IGNORECASE)
                    parsed_result["search_reason"] = m_reason.group(1) if m_reason else "router_regex_fallback"
                    m_priority = re.search(r'"recency_priority"\s*:\s*"([^"]+)"', content, re.IGNORECASE)
                    parsed_result["recency_priority"] = m_priority.group(1) if m_priority else "medium"
                    m_depth = re.search(r'"research_depth"\s*:\s*"([^"]+)"', content, re.IGNORECASE)
                    parsed_result["research_depth"] = m_depth.group(1) if m_depth else "auto"
                except Exception as regex_error:
                    logger.warning(f"Regex fallback failed for {model_name}: {regex_error}")
                    parsed_result = None

            if parsed_result is None:
                continue
            return await _finalize_router_result(parsed_result, model_name)
        except Exception as e:
            logger.error(f"Router model {model_name} exception: {e}", exc_info=True)
            continue

    logger.warning("ALL router models failed, using safe fallback")
    language = await detect_language_with_model(question)
    forced_search = _query_requires_live_search(question)
    fallback_result = {
        "is_safe": True,
        "language": language,
        "needs_search": forced_search,
        "needs_web_search": forced_search,
        "search_queries": _sanitize_router_search_queries(
            question,
            [question, f"{question} {current_year}"],
        )[: (7 if deep_research else 5)],
        "research_depth": "deep" if deep_research else "standard",
        "recency_priority": "high" if forced_search else "medium",
        "topic_type": "factual" if forced_search else "general",
        "year_mentioned": year_mentioned,
        "answer_confidence": 0.25 if forced_search else 0.5,
        "search_reason": "all_models_failed_fallback",
        "research_plan": [
            "Collect the strongest supporting sources",
            "Cross-check dates and key claims",
            "Answer only with verified facts",
        ],
    }
    logger.info(f"Router fallback result: search={fallback_result['needs_search']}, lang={fallback_result['language']}")
    return fallback_result

def _resolve_auto_research_plan(router_result: Dict[str, Any]) -> Dict[str, Any]:
    confidence = float(router_result.get("answer_confidence", 0.5) or 0.5)
    needs_search = bool(router_result.get("needs_search", router_result.get("needs_web_search", False)))
    depth = str(router_result.get("research_depth", "auto")).strip().lower()
    topic_type = str(router_result.get("topic_type", "general")).strip().lower()

    if depth == "auto":
        if needs_search and (confidence < 0.35 or topic_type in {"news", "factual", "technical"}):
            depth = "deep"
        else:
            depth = "standard"

    query_count = 5
    if depth == "deep":
        query_count = 5
    elif depth == "exhaustive":
        query_count = 7

    max_context = 8000
    max_tokens = 3000
    ctx_size = 24000
    if depth == "deep":
        max_context = 15000
        max_tokens = 5000
        ctx_size = 48000
    elif depth == "exhaustive":
        max_context = 20000
        max_tokens = 8000
        ctx_size = 64000

    return {
        "needs_search": needs_search,
        "depth": depth,
        "query_count": query_count,
        "max_context": max_context,
        "max_tokens": max_tokens,
        "ctx_size": ctx_size,
    }

def _select_generation_model(user_id: str, question: str) -> tuple[str, str]:
    """
    Production canary routing:
    - deterministic bucket by user+question to keep request stickiness
    - no manual keyword rules
    """
    def _runtime(model: str) -> str:
        # The configured MODEL_MAIN keeps its cloud name, but the streaming chat path
        # has no per-request fallback — so if cloud is not authenticated yet, generate
        # with the local fallback to keep answers flowing. After `ollama signin` (cloud
        # ready), this returns the cloud model unchanged.
        if _is_cloud(model) and not MODEL_MAIN_CLOUD_READY and LOCAL_PRIMARY_FALLBACK:
            return LOCAL_PRIMARY_FALLBACK
        return model

    if not CANARY_ENABLED or not MODEL_MAIN_CANARY:
        return _runtime(MODEL_MAIN), "control"

    seed = f"{(user_id or 'default_user').strip().lower()}|{(question or '').strip().lower()}"
    bucket = int(hashlib.md5(seed.encode("utf-8")).hexdigest()[:8], 16) % 10000
    threshold = int(CANARY_PERCENT * 100)
    if bucket < threshold:
        return _runtime(MODEL_MAIN_CANARY), "canary"
    return _runtime(MODEL_MAIN), "control"

def _tokenize_for_rank(text: str) -> List[str]:
    return [t for t in re.findall(r"[a-zA-Z0-9_]+", (text or "").lower()) if len(t) > 1]

def _hashed_embedding(text: str, dims: int = 256) -> Dict[int, float]:
    vec: Dict[int, float] = {}
    counts = Counter(_tokenize_for_rank(text))
    if not counts:
        return vec
    for token, cnt in counts.items():
        idx = int(hashlib.md5(token.encode("utf-8", errors="ignore")).hexdigest(), 16) % dims
        vec[idx] = vec.get(idx, 0.0) + float(cnt)
    norm = math.sqrt(sum(v * v for v in vec.values())) or 1.0
    for k in list(vec.keys()):
        vec[k] = vec[k] / norm
    return vec

def _cosine_sparse(a: Dict[int, float], b: Dict[int, float]) -> float:
    if not a or not b:
        return 0.0
    if len(a) > len(b):
        a, b = b, a
    return float(sum(v * b.get(k, 0.0) for k, v in a.items()))

def _lexical_score(query: str, text: str) -> float:
    q_tokens = _tokenize_for_rank(query)
    if not q_tokens:
        return 0.0
    text_l = (text or "").lower()
    hit = sum(1 for t in q_tokens if t in text_l)
    return hit / max(1, len(q_tokens))

def _get_embed_model():
    global _embed_model
    if not ST_AVAILABLE:
        return None
    if _embed_model is None:
        try:
            _embed_model = SentenceTransformer(EMBED_MODEL_NAME)
            logger.info(f"Loaded embed model: {EMBED_MODEL_NAME}")
        except Exception as e:
            logger.warning(f"Embed model load failed: {e}")
            _embed_model = None
    return _embed_model

def _get_rerank_model():
    global _rerank_model
    if not ST_AVAILABLE:
        return None
    if _rerank_model is None:
        try:
            _rerank_model = CrossEncoder(RERANK_MODEL_NAME)
            logger.info(f"Loaded reranker model: {RERANK_MODEL_NAME}")
        except Exception as e:
            logger.warning(f"Reranker load failed: {e}")
            _rerank_model = None
    return _rerank_model

def _semantic_score(query: str, text: str) -> float:
    model = _get_embed_model()
    if model is None:
        return _cosine_sparse(_hashed_embedding(query), _hashed_embedding(text))
    try:
        qv = model.encode([query], convert_to_numpy=True, show_progress_bar=False)[0]
        tv = model.encode([text], convert_to_numpy=True, show_progress_bar=False)[0]
        qn = math.sqrt(float((qv * qv).sum())) or 1.0
        tn = math.sqrt(float((tv * tv).sum())) or 1.0
        return float((qv @ tv) / (qn * tn))
    except Exception:
        return _cosine_sparse(_hashed_embedding(query), _hashed_embedding(text))

def _apply_reranker(query: str, candidates: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    model = _get_rerank_model()
    if model is None or not candidates:
        return candidates
    try:
        pairs = [[query, c.get("rank_text", c.get("snippet", ""))[:1200]] for c in candidates]
        scores = model.predict(pairs)
        for c, rs in zip(candidates, scores):
            rerank_score = float(rs)
            # Blend existing score with reranker output
            c["rerank_score"] = rerank_score
            c["score"] = (0.65 * float(c.get("score", 0.0))) + (0.35 * rerank_score)
        return sorted(candidates, key=lambda x: x.get("score", 0.0), reverse=True)
    except Exception as e:
        logger.warning(f"Reranker failed, fallback to base ranking: {e}")
        return candidates

async def tool_intent_router(question: str, deep_research: bool = False) -> Dict[str, Any]:
    prompt = f'''Analyze tool usage intent for the following question. Return JSON only.
Question: "{question}"
DeepResearchMode: {str(deep_research).lower()}
Current Vietnam local time: {get_current_datetime()["full"]}

Criteria:
1. use_database_mcp: true if query involves stable general knowledge or older archival data.
2. use_web_mcp: true if query requires internet-sourced context or specific web entities.
3. needs_live_web_search: true for time-sensitive, latest, current-status, or recent-version queries, and always true when DeepResearchMode is true.

JSON format:
{{
  "use_database_mcp": boolean,
  "use_web_mcp": boolean,
  "needs_live_web_search": boolean,
  "reason": "explanation"
}}'''
    try:
        async with httpx.AsyncClient(timeout=8.0) as client:
            response = await client.post(
                OLLAMA_URL,
                json={
                    "model": MODEL_ROUTER,
                    "messages": [{"role": "user", "content": prompt}],
                    "stream": False,
                    "options": {"temperature": 0.0, "num_predict": 80},
                },
            )
        if response.status_code == 200:
            content = response.json().get("message", {}).get("content", "").strip()
            match = re.search(r"\{.*\}", content, re.DOTALL)
            if match:
                parsed = json.loads(match.group(0))
                return {
                    "use_database_mcp": bool(parsed.get("use_database_mcp", True)),
                    "use_web_mcp": bool(parsed.get("use_web_mcp", deep_research)),
                    "needs_live_web_search": bool(parsed.get("needs_live_web_search", deep_research or _query_requires_live_search(question))),
                    "reason": str(parsed.get("reason", "model_intent_router")),
                }
    except Exception as e:
        logger.warning(f"tool_intent_router failed: {e}")
    return {
        "use_database_mcp": True,
        "use_web_mcp": True if deep_research else False,
        "needs_live_web_search": True if deep_research or _query_requires_live_search(question) else False,
        "reason": "intent_fallback",
    }

# ==================== DEEP RESEARCH ENHANCEMENTS ====================
# These functions add smart post-processing without extra network calls.

def _deduplicate_by_content(ranked: List[Dict[str, Any]], threshold: float = 0.85) -> List[Dict[str, Any]]:
    """Remove near-duplicate snippets based on text similarity (SequenceMatcher)."""
    if not ranked:
        return ranked
    unique = []
    for item in ranked:
        snippet = item.get('snippet', '')
        if not snippet:
            unique.append(item)
            continue
        is_dup = False
        for kept in unique:
            kept_snippet = kept.get('snippet', '')
            if not kept_snippet:
                continue
            ratio = difflib.SequenceMatcher(None, snippet, kept_snippet).ratio()
            if ratio >= threshold:
                is_dup = True
                break
        if not is_dup:
            unique.append(item)
    return unique

def _enforce_domain_diversity(ranked: List[Dict[str, Any]], min_domains: int = 4) -> List[Dict[str, Any]]:
    """Downscore oversampled domains to ensure result diversity."""
    if len(ranked) <= min_domains:
        return ranked
    domain_counts = {}
    parsed = []
    for item in ranked:
        url = item.get('url', '')
        try:
            domain = urlparse(url).netloc.lower() if url else ''
        except:
            domain = ''
        parsed.append((item, domain))
        domain_counts[domain] = domain_counts.get(domain, 0) + 1
    if len(domain_counts) >= min_domains:
        return ranked
    adjusted = []
    for item, domain in parsed:
        boost = 1.0
        if domain and domain_counts.get(domain, 0) > 1:
            boost = 0.95  # slight penalty for oversampled domain
        new_item = item.copy()
        new_item['score'] = item.get('score', 0) * boost
        adjusted.append(new_item)
    adjusted.sort(key=lambda x: x.get('score', 0), reverse=True)
    return adjusted

def _citation_chain_boost(ranked: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Boost scores of sources cited by other sources in the result set."""
    if not ranked:
        return ranked
    url_to_item = {item.get('url', ''): item for item in ranked if item.get('url')}
    boosted = []
    for item in ranked:
        snippet = item.get('snippet', '').lower()
        boost = 1.0
        for other_url in url_to_item:
            if other_url and other_url.lower() in snippet:
                boost += 0.1
                break
        new_item = item.copy()
        new_item['score'] = item.get('score', 0) * boost
        boosted.append(new_item)
    boosted.sort(key=lambda x: x.get('score', 0), reverse=True)
    return boosted


async def hybrid_mcp_retrieval(
    query: str,
    recency_priority: str = "medium",
    deep_research: bool = False,
    user_id: str = "default_user",
) -> Dict[str, Any]:
    intent = await tool_intent_router(query, deep_research=deep_research)
    candidates: List[Dict[str, Any]] = []
    ab_variant = _auto_assign_ab_variant(user_id, query)
    mcp_query = re.sub(r"[^\w\s]", " ", _normalize_router_text(query))
    mcp_query = re.sub(r"\s+", " ", mcp_query).strip() or query

    if db_mcp_backend and intent.get("use_database_mcp", True):
        try:
            db_res = db_mcp_backend.search_articles(mcp_query, category="all", limit=30 if deep_research else 12)
            for row in db_res.get("results", []):
                title = row.get("title", "")
                desc = row.get("description", "")
                content = row.get("content", "")
                url = row.get("url") or row.get("canonical_url") or ""
                newline = "\n"
                text_blob = f"{title}{newline}{desc}{newline}{content}"
                lex = _lexical_score(query, text_blob)
                emb = _semantic_score(query, text_blob)
                click = get_click_boost(url) if url else 0.0
                weights = _get_rank_weights(ab_variant, "database_mcp")
                final = (weights["lex"] * lex) + (weights["emb"] * emb) + (weights["click"] * click)
                candidates.append({
                    "title": title,
                    "url": url,
                    "snippet": (desc or content or "")[:180],
                    "rank_text": text_blob[:400 if not deep_research else 600],
                    "source": "database_mcp",
                    "lexical_score": lex,
                    "embedding_score": emb,
                    "click_boost": click,
                    "ab_variant": ab_variant,
                    "score": final,
                    "published_at": row.get("published_at"),  # new field for recency boost
                })
        except Exception as e:
            logger.warning(f"database_mcp retrieval failed: {e}")

    if web_mcp_backend and intent.get("use_web_mcp", False):
        try:
            web_local = web_mcp_backend.web_local_search(mcp_query, category="all", limit=30 if deep_research else 10)
            for row in web_local.get("results", []):
                title = row.get("title", "")
                snippet = row.get("snippet", "")
                url = row.get("canonical_url", "")
                newline = "\n"
                text_blob = f"{title}{newline}{snippet}"
                lex = _lexical_score(query, text_blob)
                emb = _semantic_score(query, text_blob)
                click = get_click_boost(url) if url else 0.0
                weights = _get_rank_weights(ab_variant, "web_mcp_local")
                final = (weights["lex"] * lex) + (weights["emb"] * emb) + (weights["click"] * click)
                candidates.append({
                    "title": title,
                    "url": url,
                    "snippet": snippet[:180] if snippet else "",
                    "rank_text": text_blob[:400 if not deep_research else 600],
                    "source": "web_mcp_local",
                    "lexical_score": lex,
                    "embedding_score": emb,
                    "click_boost": click,
                    "ab_variant": ab_variant,
                    "score": final,
                    "published_at": row.get("published_at") or row.get("date") or row.get("updated_at"),
                })
        except Exception as e:
            logger.warning(f"web_mcp local retrieval failed: {e}")

    live_web_text = ""
    live_web_sources: List[Dict[str, Any]] = []
    if intent.get("needs_live_web_search", False):
        try:
            live_web_result = await search_engine.smart_search(
                query,
                recency_priority=recency_priority,
                deep_research=deep_research,
                query_class="general",
            )
            if isinstance(live_web_result, dict):
                live_web_text = str(
                    live_web_result.get("context")
                    or live_web_result.get("knowledge_base")
                    or ""
                )
                live_web_sources = list(live_web_result.get("urls") or [])
            else:
                live_web_text = str(live_web_result or "")
        except Exception as e:
            logger.warning(f"live web search failed: {e}")

    # DOMAIN REPUTATION BOOST — multiply candidate scores by trusted domain factor
    for c in candidates:
        url = c.get('url', '')
        if url:
            try:
                domain = urlparse(url).netloc.lower()
                boost = get_quality_boost(domain)
                c['score'] = c.get('score', 0.0) * boost
            except Exception:
                pass

    # RECENCY BOOST (only when high recency priority)
    if recency_priority == "high":
        now_dt = datetime.utcnow()
        for c in candidates:
            pub_str = c.get("published_at")
            if not pub_str:
                continue
            try:
                # Handle ISO format (Z or timezone)
                pub_str_clean = pub_str.replace("Z", "+00:00") if "Z" in pub_str else pub_str
                pub_dt = datetime.fromisoformat(pub_str_clean)
                age_days = (now_dt - pub_dt).days
                if age_days < 7:
                    c["score"] = c.get("score", 0.0) * 1.2
                elif age_days < 30:
                    c["score"] = c.get("score", 0.0) * 1.05
            except Exception:
                continue

    dedup: Dict[str, Dict[str, Any]] = {}
    for c in candidates:
        key = (c.get("url") or c.get("title") or "").strip().lower()
        if not key:
            continue
        if key not in dedup or c["score"] > dedup[key]["score"]:
            dedup[key] = c

    ranked = sorted(dedup.values(), key=lambda x: x.get("score", 0.0), reverse=True)
    ranked = _apply_reranker(query, ranked[: (25 if not deep_research else 80)])  # larger pool for reranker
    if deep_research:
        ranked = _deduplicate_by_content(ranked)
        ranked = _enforce_domain_diversity(ranked)
        ranked = _citation_chain_boost(ranked)
        # Re-sort after score adjustments
        ranked.sort(key=lambda x: x.get("score", 0.0), reverse=True)
    top_ranked = ranked[: (15 if not deep_research else 60)]  # final selection meeting targets
    record_ranking_impressions(user_id=user_id, query=query, variant=ab_variant, ranked_results=top_ranked)

    context_blocks = []
    newline = "\n"
    for i, item in enumerate(top_ranked, 1):
        if deep_research:
            # Ultra-compact for 60+ sources
            title_len, url_len, snippet_len = 30, 40, 100
        else:
            # Standard: 15-25 sources
            title_len, url_len, snippet_len = 50, 60, 120
        block = (
            f"[{i}] {item.get('title','(no title)')[:title_len]}{newline}"
            f"URL: {item.get('url','')[:url_len]}{newline}"
            f"Snippet: {item.get('snippet','')[:snippet_len]}{newline}"
        )
        context_blocks.append(block)

    if live_web_text and len(live_web_text) > 100:
        context_blocks.append(f"[live_web]{newline}{live_web_text[:8000 if deep_research else 2000]}")

    return {
        "intent": intent,
        "ab_variant": ab_variant,
        "ranked_results": top_ranked,
        "context_text": "\n---\n".join(context_blocks),
        "has_context": len(context_blocks) > 0,
        "live_web_sources": live_web_sources,
    }

async def _ddgs_web_fallback(query: str, deep: bool = False) -> tuple:
    """Direct DuckDuckGo (ddgs library) search — used when SmartSearchEngine
    returns an empty/UNVERIFIED stub (its strict latest_model gate rejects almost
    everything). ddgs hits DDG's real API (not the bot-blocked HTML scrape), so it
    returns FRESH results with no key/Docker. Returns (context_text, url_rows)."""
    def _run():
        try:
            try:
                from ddgs import DDGS
            except Exception:
                from duckduckgo_search import DDGS
        except Exception:
            return "", []
        rows, ctx = [], []
        try:
            with DDGS() as d:
                for it in d.text(query, max_results=10 if deep else 6):
                    url = str(it.get("href") or it.get("url") or "").strip()
                    if not url:
                        continue
                    title = str(it.get("title") or "").strip()
                    snip = str(it.get("body") or it.get("snippet") or "").strip()
                    rows.append({"url": url, "title": title[:160], "snippet": snip[:300]})
                    ctx.append(f"• {title}\n{snip}\n({url})")
        except Exception as exc:
            logger.warning("ddgs fallback query failed: %s", exc)
            return "", []
        return ("\n\n".join(ctx))[:8000], rows
    try:
        return await asyncio.to_thread(_run)
    except Exception as exc:
        logger.warning("ddgs fallback runner failed: %s", exc)
        return "", []


async def fast_web_search(
    query: str,
    recency_priority: str = "medium",
    deep_research: bool = False,
    user_id: str = "default_user",
    query_class: str = "",
    on_progress: Optional[Any] = None,
) -> dict:

    """Restore the stable live-web search path via SmartSearchEngine."""
    def _dedupe_url_rows(rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        deduped: List[Dict[str, Any]] = []
        seen: set[str] = set()
        for row in rows or []:
            url = str((row or {}).get("url") or "").strip()
            if not url or url in seen:
                continue
            seen.add(url)
            deduped.append(row)
        return deduped

    try:
        resolved_query_class = str(query_class or "").strip().lower()
        if not resolved_query_class:
            try:
                if search_engine and hasattr(search_engine, "_is_latest_model_query") and search_engine._is_latest_model_query(query):
                    resolved_query_class = "latest_model"
            except Exception:
                resolved_query_class = ""
        if not resolved_query_class:
            resolved_query_class = "general"

        logger.info(f" Skemi Search: {query} | priority={recency_priority} | deep={deep_research} | class={resolved_query_class}")
        search_result = await search_engine.smart_search(
            query,
            recency_priority=recency_priority,
            deep_research=deep_research,
            query_class=resolved_query_class,
            on_progress=on_progress,
        )
        if isinstance(search_result, dict):
            context_text = str(search_result.get("context") or "")
            urls = _dedupe_url_rows(search_result.get("urls") or [])
            # FALLBACK: the engine (esp. its strict latest_model gate) frequently
            # returns an empty/UNVERIFIED stub even when the web has the answer →
            # the report then hallucinates from the model's OLD memory. When the
            # result is weak, fetch real fresh sources directly via ddgs.
            weak = ((not urls)
                    or context_text.startswith("**UNVERIFIED")
                    or context_text.startswith("**NO RELEVANT")
                    or len(context_text) < 200)
            if weak:
                fb_ctx, fb_urls = await _ddgs_web_fallback(query, deep_research)
                if fb_urls:
                    logger.info(" ddgs fallback used: %s urls | %s chars (query=%s)",
                                len(fb_urls), len(fb_ctx), query)
                    return {"context": fb_ctx, "urls": fb_urls}
            logger.info(
                " Stable search context: %s chars | urls=%s | class=%s",
                len(context_text),
                len(urls),
                resolved_query_class,
            )
            return {"context": context_text, "urls": urls}

        if isinstance(search_result, str):
            return {"context": search_result, "urls": []}

        return {"context": "Khong tim thay thong tin cap nhat tu web.", "urls": []}
    except Exception as e:
        logger.error(f"Search error: {e}")
        return {"context": f"Loi tim kiem: {str(e)}", "urls": []}


async def ai_worker():
    logger.info(" AI worker started. Waiting for requests...")
    
    while True:
        req_tuple = await request_queue.get()
        
        try:
            req, result_queue, user_ip = req_tuple
            logger.info(f"Processing: {req.question[:50]}...")
            request_started_at = time.time()
            
            question_text = req.question or ""
            question_for_model, image_analysis = extract_ai_analysis_from_question(question_text)
            if not question_for_model:
                question_for_model = question_text
            context_flags = await detect_query_context_with_model(question_text)
            has_image_context = bool(image_analysis or context_flags.get("is_image_request", False))
            incoming_history = [{"role": msg.role, "content": msg.content} for msg in (req.history or [])]

            # SESSION CONTEXT (fast temporary ring buffer) if session_id provided
            if req.session_id:
                touch_session(req.session_id)
                session_hist = session_get_context(req.session_id)
                if not session_hist:
                    # Auto-create session on first use
                    create_session(req.user_id, req.session_id)
                    session_hist = []
                merged = _merge_history_messages(session_hist, incoming_history)
                # Limit to model capacity
                if len(merged) > MAX_HISTORY_MESSAGES_TO_MODEL:
                    older_messages = merged[:-MAX_HISTORY_MESSAGES_TO_MODEL]
                    recent_messages = merged[-MAX_HISTORY_MESSAGES_TO_MODEL:]
                else:
                    older_messages = []
                    recent_messages = merged
                # Lightweight summary from older messages (if any)
                history_summary = _build_conversation_summary_from_messages(older_messages) if older_messages else ""
            else:
                # Fallback to permanent memory manager (DB-backed)
                context_bundle = memory_manager.build_context_bundle(
                    req.user_id,
                    incoming_history=incoming_history,
                    incoming_summary=req.history_summary,
                )
                recent_messages = context_bundle.get("recent_messages", [])
                history_summary = context_bundle.get("history_summary", "")
                older_messages = context_bundle.get("older_messages", [])
            requested_language = _normalize_language_code(req.language or "")
            question_language_hint = requested_language or _detect_language_hint_from_text(question_for_model)
            router_conversation_context = _build_router_conversation_context(history_summary, recent_messages)
            context_fingerprint = _build_context_fingerprint(
                question_for_model,
                history_summary=history_summary,
                recent_messages=recent_messages,
                force_search=req.force_search,
                deep_research=req.deep_research,
                has_image=has_image_context,
                language_hint=question_language_hint,
            )
            bypass_cache = _should_bypass_response_cache(
                question_for_model,
                force_search=req.force_search,
                deep_research=req.deep_research,
                has_image=has_image_context,
            )

            cached_response = None
            reusable_gist = ""
            reusable_style = ""
            if not bypass_cache:
                cached_response = await get_cached_response(
                    question_for_model,
                    context_fingerprint=context_fingerprint,
                )
                cached_sample = ""
                if cached_response:
                    cached_sample = str(cached_response.get("response") or cached_response.get("gist") or "")
                if cached_sample and not _is_cached_response_language_compatible(
                    cached_sample,
                    question_language_hint or requested_language,
                ):
                    logger.info(
                        "Ignoring cached response because its detected language does not match the user's message language."
                    )
                    cached_response = None

            if cached_response and cached_response.get("mode") in {"exact", "direct"} and cached_response.get("response"):
                cached_text = str(cached_response.get("response") or "")
                await result_queue.put({"status": "Using cached answer..."})
                await result_queue.put(
                    {
                        "token": cached_text,
                        "cache_hit": True,
                        "cache_mode": cached_response.get("mode"),
                    }
                )
                await result_queue.put({"complete": True})
                memory_manager.add_message(req.user_id, "user", question_text)
                memory_manager.add_message(req.user_id, "assistant", cached_text)
                if req.session_id:
                    session_append(req.session_id, "user", question_text)
                    session_append(req.session_id, "assistant", cached_text)
                memory_manager.refresh_summary(req.user_id, history_summary)
                logger.info(f" Served cached response for: {question_text[:50]}...")
                continue

            if cached_response and cached_response.get("mode") == "gist":
                reusable_gist = str(cached_response.get("gist") or "")
                reusable_style = str(cached_response.get("style") or "")
            
            if image_analysis:
                preview = image_analysis.replace("\n", " ")[:300]
                logger.info("Using image analysis from client context")
                logger.info(f"Image analysis preview: {preview}")

            # Auto language and status
            language = requested_language or "auto"
            await _emit_stream_event(
                result_queue,
                "phase_started",
                phase="analyzing",
                label="Analyzing",
                detail="Inspecting the request and active conversation context.",
                progress=0.08,
                language=question_language_hint or requested_language or "en",
            )

            # Image request without vision model: fail early with clear message
            if context_flags.get("is_image_request", False) and not image_analysis and not MODEL_VISION:
                msg = "No vision model is configured. Please install a multimodal model or set SKEMI_MODEL_VISION."
                await result_queue.put({"token": msg})
                await result_queue.put({"complete": True})
                memory_manager.add_message(req.user_id, "user", question_text)
                if req.session_id:
                    session_append(req.session_id, "user", question_text)
                continue

            # Keep image_analysis as context for the main model instead of returning raw OCR directly.
            
            # Router decision hoàn toàn bằng AI - pass deep_research mode
            router_result = None
            await _emit_stream_event(
                result_queue,
                "phase_started",
                phase="routing",
                label="Routing",
                detail="Deciding whether to answer from knowledge, cache, or live search.",
                progress=0.14,
                language=question_language_hint or requested_language or "en",
            )
            if req.confirmed_plan:
                logger.info("Using confirmed research plan from client")
                router_result = {
                    "needs_web_search": True,
                    "language": question_language_hint or await detect_language_with_model(question_for_model),
                    "search_queries": req.confirmed_plan,
                    "research_plan": req.confirmed_plan,
                    "answer_confidence": 0.1,
                    "needs_search": True
                }
            else:
                router_result = await pure_model_router(
                    question_for_model,
                    conversation_context=router_conversation_context,
                    deep_research=req.deep_research,
                )
                
            # If deep research is requested but not confirmed yet, return the plan for approval
            if req.deep_research and not req.confirmed_plan and router_result.get("needs_web_search"):
                logger.info("Deep research plan generated, sending for client approval")
                await _emit_stream_event(
                    result_queue,
                    "phase_updated",
                    phase="planning_search",
                    label="Research plan ready",
                    detail="Awaiting confirmation for the deep research plan.",
                    progress=0.24,
                    language=question_language_hint or requested_language or "en",
                )
                await result_queue.put({
                    "research_plan": router_result.get("research_plan", router_result.get("search_queries", [])),
                    "status": "Research plan ready"
                })
                await result_queue.put({"complete": True, "needs_confirmation": True})
                continue
            detected_language = requested_language or question_language_hint or _normalize_language_code(router_result.get("language", ""))
            if not detected_language:
                detected_language = await detect_language_with_model(question_for_model)
            language = "en" if FORCE_ENGLISH else (detected_language or "vi")
            router_result["language"] = language
            auto_plan = _resolve_auto_research_plan(router_result)
            effective_deep_research = auto_plan["depth"] in {"deep", "exhaustive"}
            if req.force_search or req.deep_research:
                logger.info("Manual search/deep flags detected from client, honoring request.")
            if req.deep_research:
                effective_deep_research = True
            await _emit_stream_event(
                result_queue,
                "phase_updated",
                phase="planning_search" if router_result.get("needs_search") or router_result.get("needs_web_search") else "routing",
                label="Planning search" if router_result.get("needs_search") or router_result.get("needs_web_search") else "Knowledge answer",
                detail=f"Depth={auto_plan['depth']}, query budget={auto_plan['query_count']}.",
                progress=0.2,
                language=language,
            )
            
            response_pipeline_started_at = time.time()
            web_context = ""
            token_count = 0
            all_urls: List[Dict[str, str]] = []
            search_time = 0.0
            
            if context_flags.get("is_image_request", False):
                logger.info("Context detector marked request as image-driven.")
            
            # Quyết định search
            needs_search = bool(auto_plan["needs_search"])
            if req.force_search or req.deep_research:
                needs_search = True
            if str(router_result.get("search_reason", "")).startswith("model_router_unavailable"):
                needs_search = needs_search or _query_requires_live_search(question_for_model)

            # Nếu prompt đến từ file/ảnh thì ưu tiên không search để tránh timeout.
            if context_flags.get("has_embedded_file_context", False):
                needs_search = False
            if context_flags.get("is_image_request", False):
                needs_search = False
            
            if needs_search:
                await _emit_stream_event(
                    result_queue,
                    "phase_started",
                    phase="searching",
                    label="Searching",
                    detail="Running parallel live-web retrieval across the planned queries.",
                    progress=0.28,
                    language=language,
                )
                
                # Tìm kiếm SONG SONG tất cả queries từ router.
                queries = router_result.get("search_queries", [question_for_model])
                logger.info(f" [v1-SEARCH] Executing search for: {queries}")
                
                search_start = time.time()
                newline = "\n"
                all_results = []
                
                # Lấy recency_priority từ router.
                recency_priority = router_result.get("recency_priority", "medium")
                
                # Deep research: search more queries (up to 4), normal: up to 2
                max_queries = int(auto_plan["query_count"])
                
                # Wrap search with progress reporting
                async def _search_one_query_with_progress(i, query, total):
                    try:
                        preview_sent_urls: Set[str] = set()

                        async def _emit_preview_sources() -> None:
                            try:
                                preview_rows = []
                                preview_tasks = []
                                if hasattr(search_engine, "_search_ddg_extensive"):
                                    preview_tasks.append(asyncio.create_task(search_engine._search_ddg_extensive(query, max_results=4)))
                                if hasattr(search_engine, "_search_qwant_extensive"):
                                    preview_tasks.append(asyncio.create_task(search_engine._search_qwant_extensive(query, max_results=4)))
                                if hasattr(search_engine, "_search_brave_extensive"):
                                    preview_tasks.append(asyncio.create_task(search_engine._search_brave_extensive(query, max_results=4)))

                                if preview_tasks:
                                    try:
                                        for task in asyncio.as_completed(preview_tasks, timeout=0.9):
                                            try:
                                                batch = list(await task or [])
                                            except Exception:
                                                continue
                                            if batch:
                                                preview_rows.extend(batch)
                                                break
                                    except asyncio.TimeoutError:
                                        pass
                                    finally:
                                        for task in preview_tasks:
                                            if not task.done():
                                                task.cancel()

                                if not preview_rows:
                                    return

                                deduped_preview = []
                                seen_preview_urls = set()
                                for row in list(preview_rows or []):
                                    if not isinstance(row, dict):
                                        continue
                                    row_url = str(row.get("url") or "").strip()
                                    if row_url and row_url in seen_preview_urls:
                                        continue
                                    if row_url:
                                        seen_preview_urls.add(row_url)
                                    deduped_preview.append(row)
                                    if len(deduped_preview) >= 4:
                                        break
                                preview_rows = deduped_preview
                                if not preview_rows:
                                    return
                                await _emit_stream_event(
                                    result_queue,
                                    "provider_result",
                                    phase="fetching",
                                    label="Fetching",
                                    detail=f"Initial preview surfaced {len(preview_rows)} candidate sources.",
                                    provider="smart_search_preview",
                                    query=query,
                                    result_count=len(preview_rows),
                                    language=language,
                                )
                                for source_item in preview_rows:
                                    source_payload = _build_source_event_payload(source_item)
                                    source_url = str(source_payload.get("url") or "").strip()
                                    if source_url and source_url in preview_sent_urls:
                                        continue
                                    if source_url:
                                        preview_sent_urls.add(source_url)
                                    await _emit_stream_event(
                                        result_queue,
                                        "source_candidate",
                                        phase="fetching",
                                        label="Source found",
                                        detail="Early candidate surfaced while deeper retrieval is still running.",
                                        query=query,
                                        source=source_payload,
                                        language=language,
                                    )
                            except Exception as preview_exc:
                                logger.debug(f"Preview search skipped for {query}: {preview_exc}")

                        await _emit_stream_event(
                            result_queue,
                            "query_started",
                            phase="searching",
                            label="Searching",
                            detail=f"Query {i + 1}/{total}",
                            progress=min(0.32 + (i * 0.04), 0.56),
                            query=query,
                            language=language,
                        )
                        await _emit_stream_event(
                            result_queue,
                            "provider_started",
                            phase="fetching",
                            label="Fetching",
                            detail="Collecting and filtering live web results.",
                            provider="smart_search",
                            query=query,
                            language=language,
                        )
                        await result_queue.put({
                            "status": f"🔍 Đang tìm kiếm ({i+1}/{total}): {query[:60]}{'...' if len(query)>60 else ''}"
                        })
                        preview_task = asyncio.create_task(_emit_preview_sources())
                        result = await fast_web_search(
                            query,
                            recency_priority=recency_priority,
                            deep_research=effective_deep_research,
                            user_id=req.user_id,
                        )
                        try:
                            await asyncio.wait_for(preview_task, timeout=0.25)
                        except Exception:
                            preview_task.cancel()
                        urls = list((result or {}).get("urls") or [])
                        await _emit_stream_event(
                            result_queue,
                            "provider_result",
                            phase="fetching",
                            label="Fetching",
                            detail=f"Received {len(urls)} candidate sources.",
                            provider="smart_search",
                            query=query,
                            result_count=len(urls),
                            language=language,
                        )
                        for source_item in urls[:6]:
                            source_payload = _build_source_event_payload(source_item)
                            source_url = str(source_payload.get("url") or "").strip()
                            if source_url and source_url in preview_sent_urls:
                                continue
                            await _emit_stream_event(
                                result_queue,
                                "source_candidate",
                                phase="fetching",
                                label="Source found",
                                detail="Candidate source discovered during live retrieval.",
                                query=query,
                                source=source_payload,
                                language=language,
                            )
                        return result
                    except Exception as e:
                        logger.error(f"Search error for {query}: {e}")
                        return None
                
                limited_queries = queries[:max_queries]
                for planned_query in limited_queries:
                    await _emit_stream_event(
                        result_queue,
                        "query_planned",
                        phase="planning_search",
                        label="Planned query",
                        detail="Prepared query for live search.",
                        query=planned_query,
                        language=language,
                    )
                search_tasks = [_search_one_query_with_progress(i, q, len(limited_queries)) for i, q in enumerate(limited_queries)]
                search_results = await asyncio.gather(*search_tasks, return_exceptions=True)
                
                for i, result in enumerate(search_results):
                    if isinstance(result, Exception):
                        logger.error(f"Search task exception: {result}")
                        continue
                    if result and isinstance(result, dict):
                        context = result.get("context", "")
                        urls = result.get("urls", [])
                        query_label = queries[i] if i < len(queries) else "unknown"
                        if context and len(context) > 100:
                            all_results.append(f"--- Search results for '{query_label}' ---{newline}{context}")
                            all_urls.extend(urls)
                
                if all_results:
                    web_context = "\n\n".join(all_results)
                    logger.info(f" Web context: {len(web_context)} chars | URLs: {len(all_urls)}")
                else:
                    web_context = "No useful fresh web results were found."
                    logger.warning("No web results found")
                unique_stream_sources = []
                seen_stream_urls = set()
                for url_data in all_urls:
                    url = str(url_data.get("url") or "").strip()
                    if not url or url in seen_stream_urls:
                        continue
                    seen_stream_urls.add(url)
                    unique_stream_sources.append(url_data)
                for accepted_source in unique_stream_sources[:18]:
                    await _emit_stream_event(
                        result_queue,
                        "source_accepted",
                        phase="ranking",
                        label="Source accepted",
                        detail="Kept after dedupe and ranking.",
                        source=_build_source_event_payload(accepted_source),
                        language=language,
                    )
                await _emit_stream_event(
                    result_queue,
                    "ranking_updated",
                    phase="ranking",
                    label="Ranking",
                    detail=f"Accepted {len(unique_stream_sources)} distinct sources.",
                    result_count=len(unique_stream_sources),
                    progress=0.68,
                    language=language,
                )
                search_time = time.time() - search_start
                logger.info(
                    f" Search time: {search_time:.2f}s (depth={auto_plan['depth']}, queries={max_queries})"
                )
                
                await _emit_stream_event(
                    result_queue,
                    "summary_started",
                    phase="summarizing",
                    label="Summarizing",
                    detail="Compressing the grounded evidence into a clean answer.",
                    progress=0.78,
                    language=language,
                )
            
            # Build prompt with optional image analysis.
            has_search_data = len(web_context) > 100 if web_context else False
            current_time = get_current_datetime()["full"]
            
            system_prompt = build_smart_prompt_v2(
                language=router_result["language"],
                has_search_data=has_search_data,
                question=question_for_model,
                current_time=current_time,
                image_analysis=image_analysis,
                deep_research=effective_deep_research
            )

            if history_summary:
                system_prompt += (
                    f"\n\n{'=' * 50}\n"
                    " CONVERSATION SUMMARY \n"
                    f"{'=' * 50}\n"
                    f"{history_summary}\n"
                    "Use this summary only as conversation context. Do not treat it as external evidence."
                )

            if reusable_gist and not has_search_data:
                system_prompt += (
                    f"\n\n{'=' * 50}\n"
                    " REUSABLE PRIOR ANSWER GIST \n"
                    f"{'=' * 50}\n"
                    f"Style hint: {reusable_style or 'plain'}\n"
                    f"{reusable_gist}\n"
                    "You may reuse this prior gist to save cost, but only if it still fits the current request and conversation."
                )
   
            # Nếu có dữ liệu search
            if has_search_data:
                # Tăng context cho deep research mode.
                max_context = int(auto_plan["max_context"])
                if len(web_context) > max_context:
                    web_context = web_context[:max_context] + "\n\n[ĐÃ RÚT GỌN]"
    
                # Thêm hướng dẫn bắt buộc: ưu tiên thông tin web đã tìm được.
                system_prompt += f"{newline}{newline}{'='*60}{newline}"
                system_prompt += " LATEST WEB CONTEXT \n"
                system_prompt += f"{'='*60}{newline}{newline}"
                system_prompt += web_context
                system_prompt += f"{newline}{newline}{'='*60}{newline}{newline}"
                system_prompt += """Important:
- Treat the provided web context as the only allowed factual basis for fresh information.
- Do NOT add, infer, or paraphrase any factual claim beyond what is explicitly supported by the provided web context.
- Prefer the newest explicitly dated facts. If dates conflict, keep the most recent one and mention the exact date when useful.
- Do NOT talk about older or outdated information unless the user explicitly asks for historical background or a timeline.
- If the search context says evidence is insufficient, or the latest answer cannot be verified, say that clearly instead of guessing.
- Format your response naturally: use **bold** for key terms and bullet lists for details. Only use ## headers if the response covers 3 or more clearly distinct, separate topics, never use headers for a single short answer.
- Use markdown tables (| col1 | col2 |) when comparing items or presenting structured data like versions, specs, or timelines.
- When the question covers multiple distinct topics, use a separate ## section for each topic.
- CRITICAL: When the question involves multiple subjects (for example two AI models or two events), locate and report the most recent information for each subject independently. Never reuse search results from one subject to answer about another.
- CRITICAL: When the web context contains per-query or per-subject coverage blocks, use each block only for its matching section in the answer.
- For each subject, include every concrete supported detail you can find in that subject's evidence block. Prefer partial but source-backed detail over a vague "không có dữ liệu" summary.
- Do NOT leave sections blank. If one field such as release date, context window, output limit, pricing, or benchmark is not supported, omit only that field and continue with the supported details.
- Only say evidence is missing for a subject when its dedicated coverage block has zero verified items or explicitly says verification failed.
- Do NOT collapse a whole section into "unknown" merely because one specification is missing.
- Do NOT promote a specific version name, release date, token limit, or benchmark from a single weak or snippet-only source unless the context clearly shows corroboration or an official/high-trust source.
- If a version string looks speculative or appears only once in non-official snippet evidence, label it as unverified instead of stating it as fact.
- If a subject block contains only low-confidence snippet evidence and no verified pages or supported version claims, do NOT state a concrete latest version/model name as fact for that subject.
- If a subject block says verification_status=reported_but_unverified_version_claims, summarize the reported version names only as unverified mentions from current sources; do not leave the section blank.
- If a subject block says verification_status=supporting_evidence_without_a_corroborated_version_name, summarize the supported product/update details and explicitly note that the evidence block did not corroborate a specific latest version name.
- If a subject block says verification_status=no_query_specific_evidence_retrieved, say that this search pass did not retrieve query-specific fresh evidence for that subject, but do not invent fill-in details.
- Never fill a missing field with a guess from a comparison article, rumor, marketing copy, or speculative roadmap language.
- Do NOT mention sources, references, or how you found the information.
- If information is uncertain, explicitly state uncertainty."""
            
            # Prepare messages
            messages = [{"role": "system", "content": system_prompt}]
            
            # Thêm history nếu có.
            if recent_messages:
                messages.extend(recent_messages)
            
            messages.append({"role": "user", "content": question_for_model})
            
            # Dùng model chính.
            model, model_lane = _select_generation_model(req.user_id, question_for_model)
            logger.info(f" Using MAIN model: {model} (lane={model_lane})")
            
            # Generate response
            gen_start = time.time()
            had_error = False
            response_chunks: List[str] = []
            first_token_time = 0.0
            max_tokens = int(auto_plan["max_tokens"])
            desired_ctx = int(auto_plan["ctx_size"])
            ctx_size = await fit_context_window_for_model(
                model,
                desired_ctx=desired_ctx,
                messages=messages,
                max_output_tokens=max_tokens,
            )
            available_ctx = await get_model_context_window(model)
            logger.info(
                f" Context window fit for {model}: requested={desired_ctx}, fitted={ctx_size}, available={available_ctx}"
            )
            
            try:
                await _emit_stream_event(
                    result_queue,
                    "phase_started",
                    phase="streaming_reply",
                    label="Streaming reply",
                    detail="Generating the final answer.",
                    progress=0.88,
                    language=language,
                )
                async with httpx.AsyncClient(timeout=httpx.Timeout(CHAT_TIMEOUT_SECONDS)) as client:
                    response_json = {
                        "model": model,
                        "messages": messages,
                        "stream": True,
                        "options": {
                            "temperature": 0.7,
                            "num_predict": max_tokens,
                            "num_ctx": ctx_size,
                            "top_k": 40,
                            "top_p": 0.9,
                            "repeat_penalty": 1.1
                        }
                    }
                    
                    async with client.stream("POST", OLLAMA_URL, json=response_json) as r:
                        if r.status_code != 200:
                            error_body = await r.aread()
                            logger.error(f" Ollama error {r.status_code}: {error_body}")
                            error_msg = _build_backend_failure_message("unavailable", language, status_code=r.status_code)
                            await result_queue.put({"token": error_msg})
                            await result_queue.put({"complete": True})
                            had_error = True
                            continue
                        
                        token_count = 0
                        
                        async for line in r.aiter_lines():
                            if line.strip():
                                try:
                                    data = json.loads(line)
                                    token = data.get("message", {}).get("content", "")
                                    if token:
                                        if first_token_time == 0:
                                            first_token_time = time.time() - gen_start
                                        token_count += 1
                                        response_chunks.append(token)
                                        await result_queue.put({"token": token})
                                except json.JSONDecodeError:
                                    pass
                                except Exception as e:
                                    logger.debug(f"Error parsing token: {e}")
            
            except httpx.TimeoutException:
                logger.error("Ollama timeout")
                timeout_msg = _build_backend_failure_message("timeout", language)
                await result_queue.put({"token": timeout_msg})
                await result_queue.put({"complete": True})
                had_error = True
            except Exception as e:
                logger.error(f"Generation error: {e}")
                error_msg = _build_backend_failure_message("generic", language)
                await result_queue.put({"token": error_msg})
                await result_queue.put({"complete": True})
                had_error = True

            full_response = "".join(response_chunks).strip()
            if not had_error and token_count == 0:
                empty_msg = _build_backend_failure_message("empty", language)
                await result_queue.put({"token": empty_msg})
                await result_queue.put({"complete": True})
            
            # Log summary
            end_to_end_time = time.time() - request_started_at
            response_pipeline_time = time.time() - response_pipeline_started_at
            gen_time = time.time() - gen_start if gen_start else 0
            
            if full_response and not has_search_data and not bypass_cache and token_count > 0:
                try:
                    await cache_response(
                        question_for_model,
                        full_response,
                        context_fingerprint=context_fingerprint,
                        model_name=model,
                    )
                except Exception as e:
                    logger.error(f"Cache save error: {e}")
            
            logger.info("\n" + "="*50)
            logger.info(f" SKEMI RESPONSE SUMMARY")
            logger.info(f" Question: {req.question[:60]}...")
            logger.info(f" Router: {MODEL_ROUTER}")
            logger.info(f" Main: {model}")
            logger.info(f" Lane: {model_lane}")
            logger.info(f" Search: {'Yes' if needs_search else 'No'} ({'Yes' if has_search_data else 'No'} data)")
            logger.info(f" Auto Plan: depth={auto_plan.get('depth')}, queries={auto_plan.get('query_count')}")
            logger.info(f" Language: {router_result['language']}")
            logger.info(f" End-to-end: {end_to_end_time:.2f}s")
            logger.info(f" Response pipeline: {response_pipeline_time:.2f}s")
            logger.info(f" Search phase: {search_time:.2f}s")
            logger.info(f" First token: {first_token_time:.2f}s" if first_token_time > 0 else " First token: no stream token")
            logger.info(f" Generation: {gen_time:.2f}s")
            logger.info(f" Speed: {token_count/gen_time:.1f} tokens/s" if gen_time > 0 and token_count > 0 else " No tokens generated")
            logger.info("="*50 + "\n")
            
            record_query_metric(
                query=req.question,
                used_search=bool(needs_search),
                has_search_data=bool(has_search_data),
                latency_ms=end_to_end_time * 1000.0,
                token_count=int(token_count),
                language=router_result.get("language", "auto"),
            )

            complete_data = {"complete": True, "event": "complete", "phase": "complete", "label": "Done", "ts": _event_ts()}
            if has_search_data and all_urls:
                unique_urls = []
                seen_urls = set()
                for url_data in all_urls:
                    url = url_data.get("url", "")
                    if url and url not in seen_urls:
                        seen_urls.add(url)
                        unique_urls.append(url_data)
                complete_data["sources"] = unique_urls[:30]
            
            await result_queue.put(complete_data)
            memory_manager.add_message(req.user_id, "user", req.question)
            if full_response:
                memory_manager.add_message(req.user_id, "assistant", full_response)
            if req.session_id:
                session_append(req.session_id, "user", req.question)
                if full_response:
                    session_append(req.session_id, "assistant", full_response)
            memory_manager.refresh_summary(req.user_id, history_summary)
            
        except asyncio.CancelledError:
            logger.info("Worker cancelled")
            break
        except Exception as e:
            logger.error(f"Worker error: {e}", exc_info=True)
            try:
                error_msg = "An error occurred. Please try again."
                await result_queue.put({"token": error_msg})
                await result_queue.put({"complete": True})
            except:
                pass
        finally:
            request_queue.task_done()

@app.post("/ask_stream")
async def ask_stream(request: Request, req: ChatRequest):
    user_ip = request.client.host
    
    # Rate limiting
    MAX_REQUESTS = 10
    RESET_HOURS = 6
    
    if user_ip not in user_requests:
        user_requests[user_ip] = {
            'count': 0, 
            'reset_time': datetime.now() + timedelta(hours=RESET_HOURS),
            'locked_until': None
        }
    
    user_data = user_requests[user_ip]
    
    # Check temporary lock window.
    if user_data.get('locked_until') and datetime.now() < user_data['locked_until']:
        reset_str = user_data['locked_until'].strftime("%H:%M %d/%m")
        return _sse_single_response(
            {"token": f"Rate limit reached. Please retry at {reset_str}", "complete": True}
        )
    
    # Reset counter nếu ã quá thời gian
    if datetime.now() > user_data['reset_time']:
        user_data['count'] = 0
        user_data['reset_time'] = datetime.now() + timedelta(hours=RESET_HOURS)
        user_data['locked_until'] = None
    
    # Check hard request cap.
    if user_data['count'] >= MAX_REQUESTS:
        lock_until = datetime.now() + timedelta(hours=RESET_HOURS)
        user_data['locked_until'] = lock_until
        reset_str = lock_until.strftime("%H:%M %d/%m")
        return _sse_single_response(
            {"token": f"Rate limit reached: {MAX_REQUESTS} requests. Reset at {reset_str}", "complete": True}
        )
    
    user_data['count'] += 1
    logger.info(f"Request {user_data['count']}/{MAX_REQUESTS} from {user_ip}")
    
    # Validate prompt length.
    if len(req.question.split()) > 6000:  # ~30k tokens estimate
        return _sse_single_response({"token": "Tin nhắn quá dài, vui lòng rút gọn.", "complete": True})
    
    # Thêm vào queue
    result_queue = asyncio.Queue()
    await request_queue.put((req, result_queue, user_ip))
    
    async def gen():
        try:
            while True:
                result = await result_queue.get()
                payload = dict(result)
                if "event" not in payload:
                    if "token" in payload:
                        payload["event"] = "token"
                    elif payload.get("complete"):
                        payload["event"] = "complete"
                    elif "research_plan" in payload:
                        payload["event"] = "phase_updated"
                        payload.setdefault("phase", "planning_search")
                        payload.setdefault("label", "Research plan ready")
                    elif "status" in payload:
                        payload["event"] = "phase_updated"
                        payload.setdefault("label", str(payload.get("status") or "").strip())
                payload.setdefault("ts", _event_ts())
                yield _sse_pack(payload)
                if payload.get("complete"):
                    break
        except Exception as e:
            logger.error(f"Stream error: {e}")
            yield _sse_pack({"event": "error", "token": "Connection error", "complete": True, "ts": _event_ts()})
    
    return StreamingResponse(gen(), media_type="text/event-stream")

@app.post("/clear_memory")
async def clear_memory(user_id: str):
    try:
        with sqlite3.connect(DB_FILE) as conn:
            conn.execute("DELETE FROM chat_history WHERE user_id = ?", (user_id,))
            conn.execute("DELETE FROM conversation_summaries WHERE user_id = ?", (user_id,))
        return {"status": "success", "message": "Memory cleared"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/cache_stats")

@app.get("/cache_stats")
async def cache_stats_endpoint():
    """Get cache statistics"""
    return await get_cache_stats()

@app.post("/clear_cache")
async def clear_cache_endpoint(cache_type: Optional[str] = None):
    """Clear cache by type or all"""
    if not CACHE_AVAILABLE:
        return {"success": False, "message": "Cache not available"}
    
    try:
        if cache_type == "messages":
            count = await advanced_cache.clear_by_type(CacheType.TEXT_MESSAGE)
            return {"success": True, "cleared": count, "type": "messages"}
        elif cache_type == "files":
            count = await advanced_cache.clear_by_type(CacheType.IMAGE_FILE)
            count += await advanced_cache.clear_by_type(CacheType.DOCUMENT_FILE)
            return {"success": True, "cleared": count, "type": "files"}
        elif cache_type == "all":
            count = await advanced_cache.clear_all()
            return {"success": True, "cleared": count, "type": "all"}
        else:
            return {"success": False, "message": f"Unknown cache type: {cache_type}"}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/test_router/{question}")
async def test_router(question: str):
    """Test endpoint cho router"""
    try:
        result = await pure_model_router(question)
        return {
            "question": question,
            "router_result": result
        }
    except Exception as e:
        return {"error": str(e)}

@app.post("/debug_search")
async def debug_search(question: str):
    """Debug search endpoint"""
    try:
        start = time.time()
        result = await search_engine.smart_search(question)
        elapsed = time.time() - start
        
        return {
            "question": question,
            "result_length": len(result),
            "elapsed_time": round(elapsed, 2),
            "preview": result[:500] + "..." if len(result) > 500 else result
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/search/engine_info")
async def search_engine_info():
    """Show active search engines and live provider status."""
    try:
        info = search_engine.get_engine_info() if hasattr(search_engine, "get_engine_info") else {}
        providers = dict(info.get("providers") or {})
        searxng_cfg = dict(providers.get("searxng") or {})
        searxng_status = {"enabled": False, "reachable": False}
        if searxng_cfg.get("enabled"):
            base_url = str(searxng_cfg.get("base_url") or "").strip().rstrip("/")
            started = time.time()
            try:
                async with httpx.AsyncClient(timeout=4.0, follow_redirects=True) as client:
                    response = await client.get(
                        f"{base_url}/search",
                        params={"q": "skemi health", "format": "json"},
                        headers={"Accept": "application/json"},
                    )
                result_count = 0
                if "json" in str(response.headers.get("content-type", "")).lower():
                    try:
                        result_count = len((response.json() or {}).get("results") or [])
                    except Exception:
                        result_count = 0
                searxng_status = {
                    "enabled": True,
                    "reachable": response.status_code == 200,
                    "status_code": response.status_code,
                    "latency_ms": round((time.time() - started) * 1000.0, 1),
                    "result_count": result_count,
                    "base_url": base_url,
                }
            except Exception as exc:
                searxng_status = {
                    "enabled": True,
                    "reachable": False,
                    "base_url": base_url,
                    "error": str(exc),
                }
        return {
            "success": True,
            "engine_info": info,
            "provider_status": {
                "searxng": searxng_status,
                "duckduckgo_html": {"enabled": bool((providers.get("duckduckgo_html") or {}).get("enabled", False))},
                "brave_html": {"enabled": bool((providers.get("brave_html") or {}).get("enabled", False))},
                "startpage_html": {"enabled": bool((providers.get("startpage_html") or {}).get("enabled", False))},
                "qwant_html": {"enabled": bool((providers.get("qwant_html") or {}).get("enabled", False))},
                "mojeek_html": {"enabled": bool((providers.get("mojeek_html") or {}).get("enabled", False))},
            },
            "generated_at": datetime.utcnow().isoformat(timespec="seconds") + "Z",
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/canary/status")
async def canary_status():
    return {
        "success": True,
        "canary_enabled": CANARY_ENABLED,
        "canary_percent": CANARY_PERCENT,
        "model_main": MODEL_MAIN,
        "model_main_canary": MODEL_MAIN_CANARY or None,
    }

@app.post("/feedback/click")
async def feedback_click(payload: ClickFeedback):
    """Record click/dwell feedback for online ranking updates."""
    try:
        updated = record_click_feedback(payload)
        auto_mined = None
        try:
            with sqlite3.connect(DB_FILE) as conn:
                cur = conn.cursor()
                cur.execute("SELECT COUNT(*) FROM click_logs")
                total_click_events = int((cur.fetchone() or (0,))[0] or 0)
            if total_click_events > 0 and (total_click_events % 25 == 0):
                auto_mined = mine_hard_negatives(limit_queries=150, negatives_per_query=2)
        except Exception as mine_err:
            logger.warning(f"Auto hard-negative mining skipped: {mine_err}")
        return {"success": True, "updated": updated, "auto_hard_negative_mining": auto_mined}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/feedback/top_docs")
async def feedback_top_docs(limit: int = 20):
    """Inspect top docs by learned ctr_score."""
    try:
        limit = max(1, min(limit, 100))
        with sqlite3.connect(DB_FILE) as conn:
            cur = conn.cursor()
            cur.execute(
                """
                SELECT doc_url, views, clicks, avg_dwell_ms, ctr_score, updated_at
                FROM ranking_feedback
                ORDER BY ctr_score DESC, clicks DESC
                LIMIT ?
                """,
                (limit,),
            )
            rows = cur.fetchall()
        return {
            "success": True,
            "items": [
                {
                    "doc_url": r[0],
                    "views": r[1],
                    "clicks": r[2],
                    "avg_dwell_ms": r[3],
                    "ctr_score": r[4],
                    "updated_at": r[5],
                }
                for r in rows
            ],
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/observability/metrics")
async def observability_metrics(window_hours: int = 24):
    """Operational metrics snapshot for search/ranking system."""
    try:
        window_hours = max(1, min(window_hours, 24 * 30))
        return {"success": True, "metrics": get_observability_summary(window_hours)}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/ranking/ab_stats")
async def ranking_ab_stats(window_hours: int = 24):
    try:
        window_hours = max(1, min(window_hours, 24 * 30))
        return {"success": True, "ab_stats": get_ab_ranking_stats(window_hours)}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/learning/policy")
async def learning_policy(window_hours: int = 24):
    try:
        window_hours = max(1, min(window_hours, 24 * 30))
        stats = get_ab_ranking_stats(window_hours)
        return {
            "success": True,
            "auto_learning_enabled": AUTO_LEARNING_ENABLED,
            "auto_learning_interval_seconds": AUTO_LEARNING_INTERVAL_SECONDS,
            "policy": stats.get("policy", {}),
            "score": stats.get("score", {}),
            "ctr": stats.get("ctr", {}),
            "impressions": stats.get("impressions", {}),
            "clicks": stats.get("clicks", {}),
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/learning/retrain")
async def learning_retrain():
    """Rebuild ranking feedback table from raw click logs."""
    try:
        result = retrain_ranking_from_click_logs()
        return {"success": True, "result": result}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/learning/hard_negatives/mine")
async def learning_hard_negatives_mine(limit_queries: int = 300, negatives_per_query: int = 3):
    try:
        result = mine_hard_negatives(limit_queries=limit_queries, negatives_per_query=negatives_per_query)
        return {"success": True, "result": result}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/evaluation/run")
async def evaluation_run(sample_size: int = 100):
    """Run offline retrieval evaluation using click logs as weak labels."""
    try:
        sample_size = max(10, min(sample_size, 500))
        result = run_offline_eval_from_clicks(sample_size=sample_size)
        if "error" in result:
            return {"success": False, **result}
        return {"success": True, "result": result}
    except Exception as e:
        return {"success": False, "error": str(e)}

@app.get("/dashboard/metrics", response_class=HTMLResponse)
async def dashboard_metrics():
    html = """
<!doctype html>
<html>
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Skemi Realtime Metrics</title>
  <style>
    :root {
      --bg: #0b1220;
      --card: #111d34;
      --fg: #e5efff;
      --muted: #93a6c7;
      --accent: #2dd4bf;
      --warn: #f59e0b;
    }
    body { margin:0; padding:20px; font-family: "Segoe UI", Arial, sans-serif; background: radial-gradient(circle at 20% 10%, #17305a, var(--bg)); color: var(--fg); }
    .grid { display:grid; grid-template-columns: repeat(auto-fit, minmax(240px, 1fr)); gap: 14px; }
    .card { background: var(--card); border: 1px solid #1f3358; border-radius: 12px; padding: 14px; }
    h1 { margin: 0 0 14px; font-size: 22px; }
    h2 { margin: 0 0 10px; font-size: 14px; color: var(--muted); text-transform: uppercase; letter-spacing: .06em; }
    .metric { font-size: 28px; font-weight: 700; }
    .sub { color: var(--muted); font-size: 12px; }
    .ok { color: var(--accent); }
    .warn { color: var(--warn); }
    pre { margin: 0; white-space: pre-wrap; color: var(--muted); font-size: 12px; }
  </style>
</head>
<body>
  <h1>Skemi Retrieval Ops Dashboard</h1>
  <div class="grid">
    <div class="card"><h2>Query Count (24h)</h2><div id="query_count" class="metric">0</div></div>
    <div class="card"><h2>Latency P95</h2><div id="p95" class="metric">0 ms</div></div>
    <div class="card"><h2>Search Usage</h2><div id="search_rate" class="metric">0%</div></div>
    <div class="card"><h2>AB CTR Delta</h2><div id="ctr_delta" class="metric">0%</div></div>
    <div class="card"><h2>Hard Negatives</h2><div id="hn_count" class="metric">0</div></div>
    <div class="card"><h2>Feedback Events</h2><div id="fb_events" class="metric">0</div></div>
  </div>
  <div class="card" style="margin-top: 14px;">
    <h2>Raw Metrics</h2>
    <pre id="raw">loading...</pre>
  </div>
  <script>
    async function fetchOnce() {
      const r = await fetch('/observability/metrics?window_hours=24');
      const j = await r.json();
      if (!j.success) return;
      const m = j.metrics || {};
      const a = m.ab_ranking || {};
      const ctr = a.ctr || {};
      const ctrA = Number(ctr.A || 0);
      const ctrB = Number(ctr.B || 0);
      const delta = ctrA > 0 ? ((ctrB - ctrA) / ctrA) * 100 : 0;
      document.getElementById('query_count').textContent = m.query_count ?? 0;
      document.getElementById('p95').textContent = `${Math.round((m.latency_ms || {}).p95 || 0)} ms`;
      document.getElementById('search_rate').textContent = `${Math.round(((m.search_usage_rate || 0) * 100))}%`;
      const d = document.getElementById('ctr_delta');
      d.textContent = `${delta.toFixed(2)}%`;
      d.className = `metric ${delta >= 0 ? 'ok' : 'warn'}`;
      document.getElementById('hn_count').textContent = (m.hard_negatives || {}).count || 0;
      document.getElementById('fb_events').textContent = (m.feedback || {}).events || 0;
      document.getElementById('raw').textContent = JSON.stringify(j, null, 2);
    }
    setInterval(fetchOnce, 2000);
    fetchOnce();
  </script>
</body>
</html>
"""
    return HTMLResponse(content=html)

@app.get("/observability/stream")
async def observability_stream(window_hours: int = 24):
    async def event_gen():
        while True:
            payload = {"success": True, "metrics": get_observability_summary(max(1, min(window_hours, 24 * 30)))}
            yield f"data: {json.dumps(payload)}\n\n"
            await asyncio.sleep(2.0)
    return StreamingResponse(event_gen(), media_type="text/event-stream")

# Session context management endpoints
@app.post("/session/create")
async def session_create(user_id: str, session_id: str):
    """Create a new session context."""
    create_session(user_id, session_id)
    return {"session_id": session_id, "status": "created"}

@app.delete("/session/{session_id}")
async def session_delete(session_id: str):
    """Delete a session context."""
    delete_session(session_id)
    return {"session_id": session_id, "status": "deleted"}

@app.post("/session/cleanup")
async def session_cleanup(older_than_hours: int = 24):
    """Cleanup expired sessions."""
    deleted = cleanup_expired_sessions()
    return {"deleted_sessions": deleted, "older_than_hours": older_than_hours}

if __name__ == "__main__":
    import uvicorn
    import asyncio
    
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        # Check Ollama và tự động cập nhật model
        status, models = loop.run_until_complete(check_ollama_status())
        if not status:
            print(" WARNING: Ollama not detected!")
            print("Please start Ollama first with: ollama serve")
            print("Then pull models:")
            print("ollama pull llama3.2:latest")
            print("ollama pull ministral-3:3b")
            print("ollama pull qwen3:4b")
    except:
        pass
    
    server_port = int(os.getenv("SKEMI_CHAT_AI_PORT", "8011"))
    uvicorn.run(app, host="127.0.0.1", port=server_port)

