import os
import json
import re
import zipfile
import sqlite3
import logging
import tempfile

# Thư viện xử lý văn bản/tài liệu
import docx
import pdfplumber
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
import yaml

# Thư viện xử lý ảnh
from PIL import Image
import pytesseract

# Thư viện xử lý Audio/Video
import speech_recognition as sr
from pydub import AudioSegment
from moviepy import VideoFileClip

# Thư viện nhận diện file (Tùy chọn, nếu lỗi có thể comment lại)
try:
    import magic
    USE_MAGIC = True
except ImportError:
    USE_MAGIC = False

# Cấu hình logging
logging.basicConfig(level=logging.INFO)

# =========================
# Helper: Extract ASCII/Unicode from binary
# =========================
def extract_from_binary(path):
    try:
        raw = open(path, "rb").read()
        ascii_text = re.findall(rb"[ -~]{5,}", raw)
        unicode_text = re.findall(rb"(?:[\x20-\x7E][\x00]){5,}", raw)
        ascii_text = b"\n".join(ascii_text).decode("utf-8", errors="ignore")
        unicode_text = b"\n".join(unicode_text).decode("utf-16", errors="ignore")
        return ascii_text + "\n" + unicode_text
    except:
        return ""

# =========================
# Extractors: Documents
# =========================
def extract_pdf(path):
    try:
        text = ""
        with pdfplumber.open(path) as pdf:
            for p in pdf.pages:
                t = p.extract_text()
                if t: text += t + "\n"
        return text
    except Exception as e:
        logging.error(f"PDF Error: {e}")
        return ""

def extract_docx(path):
    try:
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logging.error(f"DOCX Error: {e}")
        return ""

def extract_pptx(path):
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        logging.error(f"PPTX Error: {e}")
        return ""

def extract_excel(path):
    try:
        df_dict = pd.read_excel(path, sheet_name=None)
        text = ""
        for sheet_name, df in df_dict.items():
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string()
        return text
    except Exception as e:
        logging.error(f"Excel Error: {e}")
        return ""

def extract_csv(path):
    try:
        df = pd.read_csv(path)
        return df.to_string()
    except Exception as e:
        logging.error(f"CSV Error: {e}")
        return ""

# =========================
# Extractors: Structured & Code
# =========================
def extract_json(path):
    try:
        data = json.load(open(path, "r", encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except: return ""

def extract_yaml(path):
    try:
        data = yaml.safe_load(open(path, "r", encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except: return ""

def extract_html(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text("\n")
    except: return ""

def extract_txt(path):
    try:
        return open(path, "r", encoding="utf-8", errors="ignore").read()
    except: return ""

def extract_sqlite(path):
    try:
        conn = sqlite3.connect(path)
        cursor = conn.cursor()
        output = []
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
        tables = cursor.fetchall()
        for (table,) in tables:
            cursor.execute(f"SELECT * FROM {table}")
            rows = cursor.fetchall()
            output.append(f"===== TABLE {table} =====")
            for row in rows:
                output.append(str(row))
        conn.close()
        return "\n".join(output)
    except: return ""

def extract_zip(path):
    text = ""
    try:
        with zipfile.ZipFile(path) as z:
            for file in z.namelist():
                # Bỏ qua các file ẩn hoặc thư mục
                if file.startswith("__") or file.endswith("/"): continue
                try:
                    with z.open(file) as sub:
                        content = sub.read().decode("utf-8", errors="ignore")
                        text += f"\n\n=== {file} ===\n{content}"
                except: continue
    except: pass
    return text

# =========================
# Extractors: Image (OCR)
# =========================
def extract_image(path):
    """
    OCR cơ bản dùng Tesseract. 
    Lưu ý: Nếu bạn đã tích hợp EasyOCR ở Server.py thì hàm này có thể là dự phòng.
    """
    try:
        img = Image.open(path)
        # Thử tiếng Việt trước, fallback tiếng Anh
        t = pytesseract.image_to_string(img, lang="vie+eng")
        return t.strip()
    except Exception as e:
        logging.error(f"Image OCR Error: {e}")
        return ""

# =========================
# Extractors: Audio & Video (NEW)
# =========================
def extract_audio(path):
    """
    Chuyển đổi Audio sang Text dùng Google Speech Recognition (cần Internet).
    Hỗ trợ: mp3, wav, ogg, flac...
    """
    recognizer = sr.Recognizer()
    wav_path = None
    
    try:
        # 1. Chuyển đổi mọi định dạng sang WAV (chuẩn cho SpeechRecognition)
        audio = AudioSegment.from_file(path)
        
        # Giới hạn độ dài: Nếu file quá dài (> 2 phút), cắt lấy 2 phút đầu để demo (tránh timeout)
        # Trong thực tế bạn có thể chia nhỏ file để xử lý hết.
        if len(audio) > 120000: 
            logging.warning("Audio quá dài, chỉ xử lý 2 phút đầu.")
            audio = audio[:120000]

        # Tạo file tạm
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            wav_path = tmp.name
            audio.export(wav_path, format="wav")

        # 2. Nhận diện giọng nói
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)
            # Thử nhận diện tiếng Việt
            text = recognizer.recognize_google(audio_data, language="vi-VN")
            return f"[Audio Transcription]:\n{text}"

    except sr.UnknownValueError:
        return "[Audio]: Không nghe rõ nội dung."
    except sr.RequestError:
        return "[Audio]: Lỗi kết nối đến dịch vụ nhận diện giọng nói."
    except Exception as e:
        logging.error(f"Audio Extract Error: {e}")
        return f"[Audio Error]: {str(e)}"
    finally:
        # Dọn dẹp file tạm
        if wav_path and os.path.exists(wav_path):
            os.remove(wav_path)

def extract_video(path):
    """
    Tách Audio từ Video và chuyển thành Text.
    Hỗ trợ: mp4, avi, mkv...
    """
    audio_tmp = None
    try:
        # 1. Tách Audio từ Video
        video = VideoFileClip(path)
        
        # Tạo file audio tạm
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
            audio_tmp = tmp.name
        
        # Write audio (tắt logger để đỡ rối)
        video.audio.write_audiofile(audio_tmp, logger=None)
        
        # 2. Gọi hàm xử lý Audio
        text = extract_audio(audio_tmp)
        return f"[Video Transcription]:\n{text}"

    except Exception as e:
        logging.error(f"Video Extract Error: {e}")
        return f"[Video Error]: {str(e)}"
    finally:
        # Dọn dẹp
        try:
            if 'video' in locals(): video.close()
            if audio_tmp and os.path.exists(audio_tmp): os.remove(audio_tmp)
        except: pass

# =========================
# MAIN DISPATCHER
# =========================
def extract_universal_content(path):
    """
    Hàm chính để gọi từ Server.py. 
    Tự động nhận diện loại file và gọi hàm xử lý tương ứng.
    """
    ext = path.lower().split('.')[-1] if '.' in path else ''
    
    # 1. Ưu tiên check extension
    
    # Audio
    if ext in ['mp3', 'wav', 'ogg', 'flac', 'm4a', 'wma']:
        return extract_audio(path)
    
    # Video
    if ext in ['mp4', 'avi', 'mkv', 'mov', 'wmv', 'flv']:
        return extract_video(path)

    # Documents
    if ext == 'pdf': return extract_pdf(path)
    if ext in ['docx', 'doc']: return extract_docx(path)
    if ext in ['pptx', 'ppt']: return extract_pptx(path)
    if ext in ['xlsx', 'xls']: return extract_excel(path)
    if ext == 'csv': return extract_csv(path)
    
    # Structured / Code
    if ext == 'json': return extract_json(path)
    if ext in ['yaml', 'yml']: return extract_yaml(path)
    if ext in ['html', 'htm', 'xml']: return extract_html(path)
    if ext in ['sqlite', 'db']: return extract_sqlite(path)
    if ext == 'zip': return extract_zip(path)
    if ext in ['txt', 'md', 'py', 'js', 'ts', 'java', 'cpp', 'c', 'cs', 'php', 'css', 'log', 'ini', 'env']:
        return extract_txt(path)

    # Images
    # (Server.py ưu tiên dùng EasyOCR, nhưng giữ đây để dự phòng)
    if ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp']:
        return extract_image(path)

    # 2. Fallback dùng Magic (nếu cài đặt) cho các file không đuôi hoặc đuôi lạ
    if USE_MAGIC:
        try:
            mime = magic.from_file(path, mime=True)
            if mime.startswith('text/'): return extract_txt(path)
            if mime.startswith('image/'): return extract_image(path)
            if mime.startswith('audio/'): return extract_audio(path)
            if mime.startswith('video/'): return extract_video(path)
            if 'pdf' in mime: return extract_pdf(path)
        except: pass

    # 3. Last Resort: Extract binary strings
    return extract_from_binary(path)

# Giữ hàm extract_text để tương thích ngược nếu cần
extract_text = extract_universal_content